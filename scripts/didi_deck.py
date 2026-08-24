"""DiDi CX visual system: brand tokens, primitives and slide components.

Shared by the empty template (`build_canva_template.py`) and the populated
Deliverable 2 deck (`build_entregable2_deck.py`) so both stay identical in look.
"""

from __future__ import annotations

import os

from PIL import Image
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
ASSETS = os.path.join(ROOT, "assets")
MARK_PNG = os.path.join(ASSETS, "didi_mark_orange.png")

# ---------------------------------------------------------------- brand tokens

ORANGE = RGBColor(0xFF, 0x66, 0x00)
INK = RGBColor(0x1A, 0x1A, 0x1A)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

GREEN = RGBColor(0x16, 0xA3, 0x4A)
AMBER = RGBColor(0xF5, 0x9E, 0x0B)
RED = RGBColor(0xE1, 0x1D, 0x2E)
GREEN_TINT = RGBColor(0xE9, 0xF7, 0xEF)
AMBER_TINT = RGBColor(0xFE, 0xF4, 0xE4)
RED_TINT = RGBColor(0xFD, 0xEC, 0xEE)

GREY_TEXT = RGBColor(0x6B, 0x72, 0x80)
GREY_SOFT = RGBColor(0x9C, 0xA3, 0xAF)
GREY_LINE = RGBColor(0xE5, 0xE7, 0xEB)
GREY_ROW = RGBColor(0xF7, 0xF8, 0xF9)
GREY_PANEL = RGBColor(0xF4, 0xF5, 0xF6)
GREY_SERIES = RGBColor(0x9A, 0xA0, 0xA6)
GREY_CELL = RGBColor(0x4B, 0x55, 0x63)
DARK_RULE = RGBColor(0x3A, 0x3A, 0x3A)
COVER_SUB = RGBColor(0xC9, 0xCC, 0xD1)
COVER_CAP = RGBColor(0xA8, 0xAC, 0xB3)

FONT = "Inter"

SW, SH = 13.333, 7.5
# Canonical 16:9 widescreen (13.333... x 7.5 in). python-pptx's blank
# template keeps type="screen4x3" after you change cx/cy; PowerPoint then
# shows the "found a problem with the content" repair dialog.
WIDESCREEN_CX = Emu(12_192_000)
WIDESCREEN_CY = Emu(6_858_000)
MARGIN = 0.5
CONTENT_W = SW - 2 * MARGIN
BAND_H = 0.95
BAND_RULE_H = 0.045
BODY_TOP = 1.30

STATUS = {
    "green": (GREEN, GREEN_TINT),
    "amber": (AMBER, AMBER_TINT),
    "red": (RED, RED_TINT),
    "neutral": (GREY_TEXT, GREY_PANEL),
}

# ------------------------------------------------------------------ primitives


def _flat(shape):
    try:
        shape.shadow.inherit = False
    except (AttributeError, NotImplementedError):
        pass
    return shape


def _size(n, floor=0.01):
    """OOXML rejects zero or negative extents; PowerPoint then asks to repair."""
    return max(float(n), floor)


def rect(slide, x, y, w, h, fill=None, line=None, line_w=0.75, shape=MSO_SHAPE.RECTANGLE):
    s = slide.shapes.add_shape(
        shape, Inches(x), Inches(y), Inches(_size(w)), Inches(_size(h))
    )
    if fill is None:
        s.fill.background()
    else:
        s.fill.solid()
        s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line
        s.line.width = Pt(line_w)
    return _flat(s)


def round_rect(slide, x, y, w, h, radius, **kw):
    w, h = _size(w), _size(h)
    s = rect(slide, x, y, w, h, shape=MSO_SHAPE.ROUNDED_RECTANGLE, **kw)
    s.adjustments[0] = min(0.5, radius / min(w, h))
    return s


def pill(slide, x, y, w, h, **kw):
    return round_rect(slide, x, y, w, h, radius=h / 2, **kw)


def oval(slide, x, y, w, h, **kw):
    return rect(slide, x, y, w, h, shape=MSO_SHAPE.OVAL, **kw)


def dot(slide, cx, cy, d, fill):
    return oval(slide, cx - d / 2, cy - d / 2, d, d, fill=fill)


def configure_widescreen(prs):
    """Set 16:9 size and drop the leftover 4:3 type flag from the template."""
    prs.slide_width = WIDESCREEN_CX
    prs.slide_height = WIDESCREEN_CY
    sld_sz = prs._element.sldSz
    if sld_sz is not None and "type" in sld_sz.attrib:
        del sld_sz.attrib["type"]
    return prs


def line(slide, x1, y1, x2, y2, color=GREY_LINE, width=0.75, dash=None):
    """Axis-aligned rules are thin rectangles.

    python-pptx connectors for horizontal/vertical lines get cx="0" or cy="0",
    which some PowerPoint builds flag as corrupt content.
    """
    thick = max(width / 72.0, 0.01)
    if dash is None and abs(y1 - y2) < 1e-4:
        x = min(x1, x2)
        return rect(slide, x, y1 - thick / 2, max(abs(x2 - x1), 0.002), thick, fill=color)
    if dash is None and abs(x1 - x2) < 1e-4:
        y = min(y1, y2)
        return rect(slide, x1 - thick / 2, y, thick, max(abs(y2 - y1), 0.002), fill=color)
    c = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2)
    )
    c.line.color.rgb = color
    c.line.width = Pt(width)
    if dash is not None:
        c.line.dash_style = dash
    xfrm = c._element.spPr.xfrm
    if int(xfrm.ext.cx) == 0:
        xfrm.ext.cx = 1
    if int(xfrm.ext.cy) == 0:
        xfrm.ext.cy = 1
    return c


def text(slide, x, y, w, h, runs, size=12, bold=False, color=INK,
         align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, line_spacing=1.15,
         space_after=0):
    """`runs` is a string, a list of runs, or a list of paragraphs (list of runs)."""
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(_size(w)), Inches(_size(h)))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = anchor

    if isinstance(runs, str):
        paragraphs = [[(runs, {})]]
    elif runs and isinstance(runs[0], tuple):
        paragraphs = [runs]
    else:
        paragraphs = runs

    for i, para in enumerate(paragraphs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        p.space_after = Pt(space_after)
        if isinstance(para, str):
            para = [(para, {})]
        for content, over in para:
            r = p.add_run()
            r.text = content
            r.font.name = FONT
            r.font.size = Pt(over.get("size", size))
            r.font.bold = over.get("bold", bold)
            r.font.color.rgb = over.get("color", color)
    return box


def text_w(s, size, bold=False):
    return len(s) * size * (0.0083 if bold else 0.0079)


def picture(slide, path, x, y, w, h):
    """Fit an image inside the (x, y, w, h) box, preserving aspect ratio."""
    with Image.open(path) as im:
        iw, ih = im.size
    scale = min(w / iw, h / ih)
    dw, dh = _size(iw * scale), _size(ih * scale)
    return slide.shapes.add_picture(
        path, Inches(x + (w - dw) / 2), Inches(y + (h - dh) / 2), Inches(dw), Inches(dh)
    )


# ------------------------------------------------------------------ components


def status_badge(slide, x, cy, label, status, size=8.0):
    color, tint = STATUS[status]
    h = 0.30
    w = 0.12 + 0.11 + 0.07 + text_w(label, size, bold=True) + 0.14
    pill(slide, x, cy - h / 2, w, h, fill=tint)
    dot(slide, x + 0.175, cy, 0.11, color)
    text(slide, x + 0.30, cy - h / 2, w - 0.34, h, label,
         size=size, bold=True, color=color, anchor=MSO_ANCHOR.MIDDLE)
    return w


def status_icon(slide, cx, cy, d, status):
    color, _ = STATUS[status]
    dot(slide, cx, cy, d, color)
    t = 0.055
    if status == "green":
        arm = round_rect(slide, cx + 0.10 * d - 0.23 * d, cy - 0.03 * d - t / 2,
                         0.46 * d, t, radius=t / 2, fill=WHITE)
        arm.rotation = -45
        arm2 = round_rect(slide, cx - 0.11 * d - 0.13 * d, cy + 0.09 * d - t / 2,
                          0.26 * d, t, radius=t / 2, fill=WHITE)
        arm2.rotation = 45
    elif status == "amber":
        round_rect(slide, cx - 0.24 * d, cy - t / 2, 0.48 * d, t, radius=t / 2, fill=WHITE)
    else:
        arrow = rect(slide, cx - 0.20 * d, cy - 0.27 * d, 0.40 * d, 0.54 * d,
                     fill=WHITE, shape=MSO_SHAPE.DOWN_ARROW)
        arrow.adjustments[0] = 0.42
        arrow.adjustments[1] = 0.48


def delta_triangle(slide, cx, cy, size, color, up=True):
    t = rect(slide, cx - size / 2, cy - size / 2, size, size,
             fill=color, shape=MSO_SHAPE.ISOSCELES_TRIANGLE)
    if not up:
        t.rotation = 180
    return t


def didi_mark(slide, x, y, size):
    slide.shapes.add_picture(MARK_PNG, Inches(x), Inches(y), Inches(size), Inches(size))


def header_band(slide, title, subtitle=""):
    rect(slide, 0, 0, SW, BAND_H, fill=INK)
    rect(slide, 0, BAND_H, SW, BAND_RULE_H, fill=ORANGE)
    didi_mark(slide, MARGIN, 0.30, 0.36)
    text(slide, MARGIN + 0.44, 0.30, 1.0, 0.36, "DiDi",
         size=16, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
    line(slide, 1.60, 0.24, 1.60, 0.71, color=DARK_RULE, width=1.0)
    if subtitle:
        text(slide, 1.80, 0.22, 10.5, 0.32, title, size=18, bold=True, color=WHITE)
        text(slide, 1.80, 0.58, 10.5, 0.24, subtitle, size=9.5, color=GREY_SOFT)
    else:
        text(slide, 1.80, 0.28, 10.5, 0.40, title, size=18, bold=True, color=WHITE,
             anchor=MSO_ANCHOR.MIDDLE)


def page_tag(slide, label):
    """Small right-aligned sheet number inside the dark header band."""
    text(slide, SW - MARGIN - 3.0, 0.36, 3.0, 0.26, label,
         size=9, bold=True, color=ORANGE, align=PP_ALIGN.RIGHT)


def _is_sheet_tag(raw):
    """True for '07' or '07 / 32'."""
    if not raw:
        return False
    head = raw.split("/")[0].strip()
    return head.isdigit() and len(head) <= 2


def apply_sheet_numbers(prs):
    """Stamp every slide 01 / N, including cover and section dividers."""
    slides = list(prs.slides)
    n = len(slides)
    for i, slide in enumerate(slides, 1):
        label = f"{i:02d} / {n:02d}"
        found = False
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            if shape.left is None or shape.top is None:
                continue
            if abs(shape.left.inches - (SW - MARGIN - 3.0)) > 0.08:
                continue
            if abs(shape.top.inches - 0.36) > 0.08:
                continue
            if not _is_sheet_tag(shape.text_frame.text.strip()):
                continue
            paras = shape.text_frame.paragraphs
            if not paras:
                continue
            runs = paras[0].runs
            if runs:
                runs[0].text = label
                for extra in runs[1:]:
                    extra.text = ""
            else:
                paras[0].text = label
            found = True
        if not found:
            page_tag(slide, label)


def panel(slide, x, y, w, h, title=None, title_size=10):
    round_rect(slide, x, y, w, h, radius=0.08, fill=WHITE, line=GREY_LINE)
    if title:
        text(slide, x + 0.22, y + 0.18, w - 0.44, 0.24, title,
             size=title_size, bold=True, color=INK)
    return x + 0.22, y + (0.52 if title else 0.20), w - 0.44


def kpi_card(slide, x, y, w, h, label, value, delta, status, chip,
             value_size=30, up=True):
    color, _ = STATUS[status]
    round_rect(slide, x, y, w, h, radius=0.10, fill=WHITE, line=color, line_w=1.25)
    text(slide, x + 0.22, y + 0.18, w - 0.90, 0.22, label, size=9, color=GREY_TEXT)
    status_icon(slide, x + w - 0.40, y + 0.29, 0.32, status)
    text(slide, x + 0.20, y + 0.44, w - 0.40, 0.60, value,
         size=value_size, bold=True, color=INK)
    if delta:
        delta_triangle(slide, x + 0.27, y + h - 0.71, 0.11, color, up=up)
        text(slide, x + 0.40, y + h - 0.82, w - 0.60, 0.26, delta, size=9.5,
             color=GREY_TEXT)
    status_badge(slide, x + 0.20, y + h - 0.22, chip, status, size=8.5)


def cause_chain(slide, x, y, w, steps, h=0.70):
    """Five-step mechanism. `steps` is [(label, tag)]. tag is a metric name or ''."""
    n = len(steps)
    gap = 0.26
    box_w = (w - (n - 1) * gap) / n
    for i, (label, tag) in enumerate(steps):
        bx = x + i * (box_w + gap)
        round_rect(slide, bx, y, box_w, h, radius=0.08, fill=GREY_PANEL, line=GREY_LINE)
        rect(slide, bx, y, 0.042, h, fill=ORANGE)
        if tag:
            text(slide, bx + 0.08, y + 0.05, box_w - 0.14, 0.16, tag,
                 size=8, bold=True, color=ORANGE, align=PP_ALIGN.CENTER)
            text(slide, bx + 0.08, y + 0.22, box_w - 0.14, h - 0.28, label,
                 size=10, bold=True, color=INK, align=PP_ALIGN.CENTER,
                 anchor=MSO_ANCHOR.MIDDLE)
        else:
            text(slide, bx + 0.08, y + 0.08, box_w - 0.14, h - 0.16, label,
                 size=10, bold=True, color=INK, align=PP_ALIGN.CENTER,
                 anchor=MSO_ANCHOR.MIDDLE)
        if i < n - 1:
            text(slide, bx + box_w, y, gap, h, "→", size=14, bold=True, color=ORANGE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    return y + h


def callout(slide, x, y, w, h, headline, body, size_head=15, size_body=11):
    w, h = _size(w, 0.4), _size(h, 0.42)
    rect(slide, x, y, w, h, fill=GREY_PANEL)
    rect(slide, x, y, 0.042, h, fill=ORANGE)
    if h < 1.10:
        text(slide, x + 0.34, y + 0.05, w - 0.62, 0.22, headline,
             size=size_head, bold=True, color=INK, anchor=MSO_ANCHOR.MIDDLE)
        text(slide, x + 0.34, y + 0.26, w - 0.62, max(h - 0.32, 0.16), body,
             size=size_body, color=GREY_CELL, line_spacing=1.18)
        return
    text(slide, x + 0.34, y + 0.20, w - 0.62, 0.32, headline,
         size=size_head, bold=True, color=INK)
    text(slide, x + 0.34, y + 0.62, w - 0.62, max(h - 0.78, 0.20), body,
         size=size_body, color=GREY_CELL, line_spacing=1.32, space_after=4)


def bullets(slide, x, y, w, items, size=11, gap=0.34, color=GREY_CELL):
    """Each item is a string or a list of runs. Orange dash marker."""
    cy = y
    for item in items:
        rect(slide, x, cy + 0.075, 0.13, 0.028, fill=ORANGE)
        box = text(slide, x + 0.26, cy - 0.03, w - 0.26, gap, item,
                   size=size, color=color, line_spacing=1.28)
        n_lines = max(1, int(_estimate_lines(item, w - 0.26, size)))
        cy += gap + (n_lines - 1) * size * 0.0175
        del box
    return cy


def _estimate_lines(item, width_in, size):
    if isinstance(item, str):
        length = len(item)
    else:
        length = sum(len(t) for t, _ in item)
    return max(1, (length * size * 0.0079) / width_in + 0.35)


def data_table(slide, x, y, cols, rows, header_h=0.40, row_h=0.42,
               status_cols=(), heat_cols=(), align_right=(), size=8.5,
               header_size=8.5, bold_first=False):
    """cols: [(label, width)].

    Cells in `status_cols` are (label, status) and render as a badge; cells in
    `heat_cols` are (label, status) and render as a tinted, colour-coded value so
    a reader can scan the numbers themselves for green/amber/red.
    """
    total_w = sum(w for _, w in cols)
    rect(slide, x, y, total_w, header_h, fill=INK)
    cx = x
    for j, (label, w) in enumerate(cols):
        align = PP_ALIGN.RIGHT if j in align_right else PP_ALIGN.LEFT
        text(slide, cx + 0.14, y, w - 0.28, header_h, label, size=header_size,
             bold=True, color=WHITE, align=align, anchor=MSO_ANCHOR.MIDDLE)
        cx += w

    ry = y + header_h
    for i, row in enumerate(rows):
        if i % 2 == 1:
            rect(slide, x, ry, total_w, row_h, fill=GREY_ROW)
        cx = x
        for j, ((_, w), value) in enumerate(zip(cols, row)):
            if j in status_cols:
                label, status = value
                status_badge(slide, cx + 0.14, ry + row_h / 2, label, status, size=7.5)
            elif j in heat_cols:
                label, status = value
                color = GREY_CELL
                if status is not None:
                    color, tint = STATUS[status]
                    pad = min(0.07, row_h * 0.16)
                    round_rect(slide, cx + 0.07, ry + pad, w - 0.14, row_h - 2 * pad,
                               radius=0.05, fill=tint)
                text(slide, cx + 0.14, ry, w - 0.28, row_h, str(label), size=size,
                     color=color, bold=status is not None, align=PP_ALIGN.RIGHT,
                     anchor=MSO_ANCHOR.MIDDLE)
            else:
                align = PP_ALIGN.RIGHT if j in align_right else PP_ALIGN.LEFT
                bold = bold_first and j == 0
                text(slide, cx + 0.14, ry, w - 0.28, row_h, str(value), size=size,
                     color=INK if bold else GREY_CELL, bold=bold, align=align,
                     anchor=MSO_ANCHOR.MIDDLE)
            cx += w
        line(slide, x, ry + row_h, x + total_w, ry + row_h, color=GREY_LINE)
        ry += row_h
    return ry


def footnote(slide, note, y=None):
    text(slide, MARGIN, y if y is not None else SH - 0.42, CONTENT_W, 0.24, note,
         size=7.5, color=GREY_SOFT)


# --------------------------------------------------------------------- iconography


def icon_calendar(slide, cx, cy, s):
    body_y = cy - s / 2 + 0.07
    round_rect(slide, cx - s / 2, body_y, s, s - 0.07, radius=0.07, line=ORANGE, line_w=1.25)
    rect(slide, cx - s / 2 + 0.02, body_y + 0.02, s - 0.04, 0.075, fill=ORANGE)
    for dx in (-0.22, 0.22):
        rect(slide, cx + dx * s - 0.022, cy - s / 2 - 0.02, 0.044, 0.09, fill=ORANGE)
    for row in range(2):
        for col in range(3):
            rect(slide, cx - 0.19 * s + col * 0.19 * s - 0.022,
                 body_y + 0.19 + row * 0.14, 0.045, 0.045, fill=ORANGE)


def icon_globe(slide, cx, cy, s):
    oval(slide, cx - s / 2, cy - s / 2, s, s, line=ORANGE, line_w=1.25)
    oval(slide, cx - s * 0.21, cy - s / 2, s * 0.42, s, line=ORANGE, line_w=1.0)
    line(slide, cx - s / 2, cy, cx + s / 2, cy, color=ORANGE, width=1.0)
    line(slide, cx - s * 0.45, cy - s * 0.22, cx + s * 0.45, cy - s * 0.22,
         color=ORANGE, width=1.0)


def icon_people(slide, cx, cy, s):
    oval(slide, cx - s * 0.17, cy - s * 0.46, s * 0.34, s * 0.34, line=ORANGE, line_w=1.25)
    round_rect(slide, cx - s * 0.30, cy - s * 0.02, s * 0.60, s * 0.36,
               radius=s * 0.17, line=ORANGE, line_w=1.25)
    for sign in (-1, 1):
        oval(slide, cx + sign * s * 0.36 - s * 0.13, cy - s * 0.36, s * 0.26, s * 0.26,
             line=ORANGE, line_w=1.0)
        round_rect(slide, cx + sign * s * 0.42 - s * 0.16, cy + s * 0.04,
                   s * 0.32, s * 0.28, radius=s * 0.13, line=ORANGE, line_w=1.0)


def icon_bars(slide, cx, cy, s):
    base = cy + s * 0.40
    for i, frac in enumerate((0.42, 0.70, 0.95)):
        w, h = s * 0.20, s * frac
        round_rect(slide, cx - s * 0.36 + i * s * 0.31, base - h, w, h,
                   radius=w / 2, fill=ORANGE)


def icon_ring(slide, cx, cy, s):
    oval(slide, cx - s / 2, cy - s / 2, s, s, line=ORANGE, line_w=1.5)
    dot(slide, cx, cy, s * 0.30, ORANGE)


def icon_trend(slide, cx, cy, s):
    pts = [(cx - s * 0.45, cy + s * 0.22), (cx - s * 0.12, cy - s * 0.10),
           (cx + s * 0.10, cy + s * 0.08), (cx + s * 0.42, cy - s * 0.32)]
    for (x1, y1), (x2, y2) in zip(pts, pts[1:]):
        line(slide, x1, y1, x2, y2, color=ORANGE, width=1.5)
    dot(slide, pts[-1][0], pts[-1][1], s * 0.16, ORANGE)


def icon_bulb(slide, cx, cy, ring_d):
    oval(slide, cx - ring_d / 2, cy - ring_d / 2, ring_d, ring_d,
         fill=RGBColor(0xFD, 0xF1, 0xE8), line=ORANGE, line_w=1.0)
    b = ring_d * 0.34
    oval(slide, cx - b / 2, cy - b * 0.85, b, b, line=ORANGE, line_w=1.5)
    round_rect(slide, cx - b * 0.20, cy + b * 0.18, b * 0.40, b * 0.30,
               radius=b * 0.08, line=ORANGE, line_w=1.5)
    line(slide, cx - b * 0.16, cy + b * 0.58, cx + b * 0.16, cy + b * 0.58,
         color=ORANGE, width=1.5)


# ------------------------------------------------------------------ slide shells


def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def content_slide(prs, title, subtitle="", tag=""):
    s = blank(prs)
    header_band(s, title, subtitle)
    if tag:
        page_tag(s, tag)
    return s


def section_slide(prs, number, title, lead):
    s = blank(prs)
    rect(s, 0, 0, SW, SH, fill=INK)
    rect(s, 0, 0, 0.16, SH, fill=ORANGE)
    text(s, 1.30, 2.35, 3.0, 1.1, number, size=64, bold=True, color=ORANGE)
    text(s, 1.30, 3.55, 9.5, 0.6, title, size=30, bold=True, color=WHITE)
    rect(s, 1.30, 4.30, 1.0, 0.035, fill=ORANGE)
    text(s, 1.30, 4.58, 8.6, 0.9, lead, size=13, color=COVER_SUB, line_spacing=1.4)
    didi_mark(s, SW - 1.35, SH - 1.15, 0.42)
    text(s, SW - 0.85, SH - 1.15, 1.0, 0.42, "DiDi", size=15, bold=True, color=WHITE,
         anchor=MSO_ANCHOR.MIDDLE)
    return s


def build_mark(src=None, dst=None, target=(0xFF, 0x66, 0x00)):
    """Recolour the DiDi mark to brand orange, preserving its alpha channel."""
    src = src or os.path.join(ASSETS, "didi_favicon.png")
    dst = dst or MARK_PNG
    img = Image.open(src).convert("RGBA")
    out = Image.new("RGBA", img.size, (*target, 255))
    out.putalpha(img.getchannel("A"))
    out.save(dst)
