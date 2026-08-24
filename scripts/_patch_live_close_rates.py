"""Patch recovered deck with Phone 9/10 vs Chat 6/10 close rates. Do not rebuild."""
from __future__ import annotations

import os
import shutil
from datetime import datetime

from pptx import Presentation

ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
LIVE = os.path.join(ROOT, "entregable 2", "deck", "Entregable_2_Weekly_Performance_Report.pptx")
COPIES = [
    os.path.join(ROOT, "entregable 2", "Entregable_2_Weekly_Performance_Report.pptx"),
    os.path.join(ROOT, "entregable 2", "Entregable_2_LISTO.pptx"),
]


def set_text(shape, new):
    tf = shape.text_frame
    tf.word_wrap = True
    first = True
    for para in tf.paragraphs:
        if first:
            if para.runs:
                para.runs[0].text = new
                for extra in para.runs[1:]:
                    extra.text = ""
            else:
                para.text = new
            first = False
        else:
            for run in para.runs:
                run.text = ""


def patch_slide(slide, rules):
    """rules: list of (match, new). match is exact or prefix if it ends with * """
    n = 0
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        raw = shape.text_frame.text.strip()
        if not raw:
            continue
        for match, new in rules:
            if raw == new:
                continue
            prefix = match[:-1] if match.endswith("*") else None
            if prefix is not None:
                if raw.startswith(prefix):
                    set_text(shape, new)
                    n += 1
                    break
            elif match in raw:
                set_text(shape, raw.replace(match, new) if match != raw else new)
                if match == raw:
                    set_text(shape, new)
                else:
                    set_text(shape, raw.replace(match, new))
                n += 1
                break
    return n


def main():
    rec = os.path.join(os.path.dirname(LIVE), "recovery")
    os.makedirs(rec, exist_ok=True)
    bak = os.path.join(rec, f"Entregable_2_before_close_{datetime.now():%Y%m%d_%H%M%S}.pptx")
    shutil.copy2(LIVE, bak)
    prs = Presentation(LIVE)
    hits = {}

    hits["exec"] = patch_slide(prs.slides[2], [
        ("Live Chat scores 96.01 QA / 77.6% CSAT, while Phone scores 83.04 QA / 86.3% CSAT.",
         "Live Chat scores 96.01 QA / 77.6% CSAT, while Phone scores 83.04 QA / 86.3% CSAT. "
         "Auditor close rate (not FCR): Phone 9 of 10 assessed (87.9%), Chat 6 of 10 (63.0%). "
         "499 of 658 unresolved chats followed process."),
    ])

    hits["critical"] = patch_slide(prs.slides[4], [
        ("Resolution vs. Compliance:*",
         "Resolution vs. Compliance: CSAT tracks closure, not process. Phone closes 9 of 10 assessed (87.9%); Chat 6 of 10 (63.0%). "
         "499 of 658 unresolved chats followed process. The rubric fix makes QA honest; CSAT moves when policy and tools let agents close."),
    ])

    hits["qa_ch"] = patch_slide(prs.slides[7], [
        ("Phone QA sits below target at 83.04*",
         "Phone QA sits below target at 83.04 (-1.96 pts), while Chat reaches 96.01. Chat is 85.6% of audits, so it inflates blended QA to 94.14. "
         "Chat's score reflects a lenient form (no critical for information accuracy). Recalibrating Chat to Phone rules drops it to 80.14. "
         "Among assessed contacts (not FCR), Phone still closes 9 of 10 (87.9%) and Chat 6 of 10 (63.0%). "
         "Phone CSAT 86.3% vs Chat 77.6%. The greener QA score is the weaker closer."),
    ])

    hits["phone"] = patch_slide(prs.slides[8], [
        ("Resolution Bottleneck:*",
         "Resolution Bottleneck: Phone still closes 9 of 10 assessed contacts (87.9%; abandon 4.8%). "
         "The 1 that does not close is mostly undelivered or refund; 28 of 41 unresolved followed process."),
    ])

    hits["chat"] = patch_slide(prs.slides[9], [
        ("Recalibrated Score Impact:*",
         "Recalibrated Score Impact: Chat closes 6 of 10 assessed (63.0%; abandon 15.6%). "
         "499 of 658 unresolved chats followed process. Recalibrating the form drops QA to 80.14; it does not close those cases."),
    ])

    hits["closure"] = patch_slide(prs.slides[25], [
        ("Strategic Pivot:*",
         "Strategic Pivot: Judge channels on close rate, not QA. Phone closes 9 of 10 assessed; Chat 6 of 10. "
         "Do not coach Chat volume for a gap that is 499 process-followed non-closes."),
    ])

    hits["next"] = patch_slide(prs.slides[28], [
        ("Report closure rate alongside QA*",
         "Report auditor closure next to QA, by channel and contact reason (not FCR). "
         "Phone closes 9 of 10 assessed (87.9%); Chat 6 of 10 (63.0%). Green if 70%+ closure and 85%+ CSAT."),
        ("Phone share of audits", "Auditor close rate (Phone / Chat)"),
        ("14.4%", "87.9% / 63.0%"),
        ("30%", "90% / 75%"),
        ("1 month", "1 quarter"),
    ])

    hits["bt"] = patch_slide(prs.slides[30], [
        ("CSAT 80.74% on 46,071 surveys; order status repeats at 19.34%",
         "CSAT 80.74% on 46,071 surveys; order status repeats at 19.34%. Chat closes 6 of 10 assessed: this is that gap, not a staffing gap."),
        ("Give Phone agents limited same-call refund authority for eligible undelivered orders",
         "Give Phone agents limited same-call refund authority for eligible undelivered orders. Phone already closes 9 of 10; this attacks the remaining 1 and the Chat refund pile."),
    ])

    hits["ch_plan"] = patch_slide(prs.slides[31], [
        ("Do not add Chat headcount. The 15.99% sits on contact reasons that cannot close, mainly order status and refunds.",
         "Do not add Chat headcount. Chat closes 6 of 10 assessed (63.0%); 499 of 658 unresolved followed process. Order status and refunds will not close with more agents."),
        ("Live Chat generates 67.2% of every repeat contact in the business",
         "Chat abandon 15.6% vs Phone 4.8%. Live Chat generates 67.2% of every repeat contact."),
        ("Report closure rate weekly and assign an owner to every unclosed case",
         "Report auditor closure weekly by channel (Phone 87.9%, Chat 63.0%) and assign an owner to every unclosed case. Not FCR."),
        ("The field is already on the form. Followed process, did not solve, still averages 96.87 on QA",
         "499 Chat and 28 Phone unresolved cases followed process. The field is already on the form (QA 96.87)."),
        ("This month also correct scored Pareto defects:*",
         "Phone closes 9 of 10 assessed; Chat 6 of 10. Do not staff Chat for that gap: most unresolved chats followed process. "
         "ETA and refunds change the outcome. Pareto coaching and the Chat critical change the scorecard. Do not mix those."),
    ])

    hits["reco"] = patch_slide(prs.slides[32], [
        ("The priority is not simply more coaching. Measure resolution. Correct the scored defects this month. Fix the QA scoring gaps. Give agents a way to close the highest-volume failure reasons.",
         "The priority is not simply more coaching. Measure resolution. "
         "Phone closes 9 of 10 assessed contacts; Chat 6 of 10. Correct the scored defects this month. "
         "Fix the QA scoring gaps. Give agents a way to close the highest-volume failure reasons."),
        ("Put % closed next to QA, by contact reason.*",
         "Put % closed next to QA, by channel and contact reason. Phone closes 9 of 10 assessed (87.9%); Chat 6 of 10 (63.0%). Not FCR. "
         "Closure is strongly associated with CSAT. Close-rate vs CSAT: R² 0.64."),
        ("Policy, tools, and visibility on the highest-volume failure reasons.*",
         "Policy, tools, and visibility on the highest-volume failure reasons. "
         "Phone already closes 9 of 10 assessed; Chat leaves 4 of 10 open, mostly after following process. "
         "The priority is that constraint, not Chat coaching volume."),
    ])

    try:
        prs.save(LIVE)
        saved = LIVE
    except PermissionError:
        saved = os.path.join(ROOT, "entregable 2", "Entregable_2_LISTO.pptx")
        prs.save(saved)
        print("LIVE open; wrote", saved)

    print("backup", bak)
    print("saved", saved, "hits", hits)
    for dest in COPIES:
        if os.path.abspath(dest) == os.path.abspath(saved):
            continue
        try:
            shutil.copy2(saved, dest)
            print("copy", dest)
        except OSError as exc:
            print("skip", dest, exc)


if __name__ == "__main__":
    main()
