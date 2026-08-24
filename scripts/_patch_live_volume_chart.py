"""Replace the two volume tables with chart + heat table. Do not rebuild."""
from __future__ import annotations

import os
import shutil
import sys
from datetime import datetime

from pptx import Presentation

HERE = os.path.dirname(os.path.abspath(__file__))
ROOT = os.path.dirname(HERE)
sys.path.insert(0, HERE)

import didi_deck as dd  # noqa: E402
from build_e2_charts import chart_csat_fail_volume, chart_qa_fail_volume  # noqa: E402

LIVE = os.path.join(ROOT, "entregable 2", "deck", "Entregable_2_Weekly_Performance_Report.pptx")
COPY = os.path.join(ROOT, "entregable 2", "Entregable_2_Weekly_Performance_Report.pptx")
CHARTS = os.path.join(ROOT, "entregable 2", "deck", "charts")
QA_TITLE = "QA analysis: lowest score vs fail volume"
CSAT_TITLE = "CSAT analysis: lowest score vs detractor volume"
KEEP_HEADS = (
    "Phone QA volume focus",
    "CSAT focus is order status",
    "Completed-not-received full service is 26%",
    "The two order-status reasons plus cancellation charge",
)


def st(value, goal):
    gap = goal - value
    if gap <= 0:
        return "green"
    return "amber" if gap <= 5 else "red"


def st_n(n):
    if n < 10:
        return "red"
    if n < 20:
        return "amber"
    return "green"


def st_csat_share(pct):
    if pct >= 10:
        return "red"
    if pct >= 5:
        return "amber"
    return "green"


def find_slide(prs, title):
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame and title in shape.text_frame.text:
                return slide
    raise KeyError(title)


def wipe_band(slide, y0=1.18, y1=4.48):
    doomed = []
    for shape in slide.shapes:
        try:
            y = shape.top.inches
        except (AttributeError, TypeError, ValueError):
            continue
        if not (y0 <= y < y1):
            continue
        if shape.has_text_frame:
            blob = shape.text_frame.text
            if any(k in blob for k in KEEP_HEADS):
                continue
        doomed.append(shape._element)
    for el in doomed:
        parent = el.getparent()
        if parent is not None:
            parent.remove(el)
    return len(doomed)


def paint_qa(slide):
    n = wipe_band(slide)
    dd.picture(slide, os.path.join(CHARTS, "bar_qa_fail_volume.png"),
               dd.MARGIN, dd.BODY_TOP, 7.10, 3.18)
    x = 7.75
    dd.text(slide, x, dd.BODY_TOP, 5.10, 0.24, "Lowest QA (rate, n >= 3)",
            size=10.5, bold=True, color=dd.INK)
    dd.data_table(
        slide, x, dd.BODY_TOP + 0.32,
        [("Contact reason", 2.40), ("QA", 0.55), ("n", 0.42), ("Start here?", 1.70)],
        [
            ["Completed not received (market place)",
             ("47.5", st(47.5, 85)), ("4", st_n(4)), ("Too few", "red")],
            ["Active order, already received",
             ("65.8", st(65.8, 85)), ("12", st_n(12)), ("Watch", "amber")],
            ["Completed not received (full service)",
             ("68.2", st(68.2, 85)), ("49", st_n(49)), ("Start", "green")],
            ["Verbal aggression",
             ("76.0", st(76.0, 85)), ("5", st_n(5)), ("Too few", "red")],
            ["Refund status and conditions",
             ("76.4", st(76.4, 85)), ("25", st_n(25)), ("Start", "green")],
        ],
        header_h=0.36, row_h=0.50, heat_cols=(1, 2), status_cols=(3,),
        align_right=(1, 2), size=7.8, bold_first=True,
    )
    return n


def paint_csat(slide):
    n = wipe_band(slide)
    dd.picture(slide, os.path.join(CHARTS, "bar_csat_fail_volume.png"),
               dd.MARGIN, dd.BODY_TOP, 7.10, 3.18)
    x = 7.75
    dd.text(slide, x, dd.BODY_TOP, 5.10, 0.24, "Worst CSAT (rate, min 100 surveys)",
            size=10.5, bold=True, color=dd.INK)
    dd.data_table(
        slide, x, dd.BODY_TOP + 0.32,
        [("Contact reason", 2.35), ("CSAT", 0.70), ("Surveys", 0.78), ("Share", 1.24)],
        [
            ["After-sales fraud review",
             ("6.1%", st(6.1, 85)), "475", ("2.9%", st_csat_share(2.9))],
            ["Other (unmapped Business Type)",
             ("26.5%", st(26.5, 85)), "558", ("2.6%", st_csat_share(2.6))],
            ["Membership program renewal",
             ("28.6%", st(28.6, 85)), "248", ("1.1%", st_csat_share(1.1))],
            ["Membership program benefits",
             ("43.1%", st(43.1, 85)), "109", ("0.4%", st_csat_share(0.4))],
            ["Placing an order information",
             ("50.7%", st(50.7, 85)), "142", ("0.5%", st_csat_share(0.5))],
        ],
        header_h=0.36, row_h=0.50, heat_cols=(1, 3),
        align_right=(1, 2, 3), size=7.8, bold_first=True,
    )
    return n


def main():
    rec = os.path.join(os.path.dirname(LIVE), "recovery")
    os.makedirs(rec, exist_ok=True)
    bak = os.path.join(rec, f"Entregable_2_before_chart_{datetime.now():%Y%m%d_%H%M%S}.pptx")
    shutil.copy2(LIVE, bak)

    print("charts")
    chart_qa_fail_volume()
    chart_csat_fail_volume()

    prs = Presentation(LIVE)
    n_qa = paint_qa(find_slide(prs, QA_TITLE))
    n_csat = paint_csat(find_slide(prs, CSAT_TITLE))
    try:
        prs.save(LIVE)
        saved = LIVE
    except PermissionError:
        saved = COPY
        prs.save(saved)
        print("LIVE open; wrote", saved)
    if os.path.abspath(saved) != os.path.abspath(COPY):
        try:
            shutil.copy2(saved, COPY)
        except OSError as exc:
            print("skip copy", exc)
    print("backup", bak)
    print("saved", saved, "wiped_qa", n_qa, "wiped_csat", n_csat)


if __name__ == "__main__":
    main()
