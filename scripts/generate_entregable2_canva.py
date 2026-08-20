"""
Genera deck PowerPoint listo para importar en Canva.
Canva: Create design > Import > PPTX, o arrastra el archivo.
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parent.parent
OUT = ROOT / "entregable 2" / "canva"
CHARTS = ROOT / "entregable 2" / "charts"

ORANGE = RGBColor(255, 102, 0)
DARK = RGBColor(26, 26, 26)
WHITE = RGBColor(255, 255, 255)
GREEN = RGBColor(40, 167, 69)
AMBER = RGBColor(255, 193, 7)
RED = RGBColor(220, 53, 69)
GRAY = RGBColor(120, 120, 120)

SLIDE_W = Inches(13.333)  # 16:9
SLIDE_H = Inches(7.5)


def blank_slide(prs: Presentation):
    return prs.slides.add_slide(prs.slide_layouts[6])


def add_orange_bar(slide, height=Inches(0.12)):
    shape = slide.shapes.add_shape(1, Inches(0), Inches(0), SLIDE_W, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = ORANGE
    shape.line.fill.background()


def add_title(slide, text: str, y=Inches(0.35), size=28):
    box = slide.shapes.add_textbox(Inches(0.6), y, Inches(12), Inches(0.7))
    tf = box.text_frame
    tf.text = text
    p = tf.paragraphs[0]
    p.font.size = Pt(size)
    p.font.bold = True
    p.font.color.rgb = ORANGE


def add_body(slide, lines: list[str], y=Inches(1.2), size=16, width=12.0):
    box = slide.shapes.add_textbox(Inches(0.6), y, Inches(width), Inches(5.5))
    tf = box.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(size)
        p.font.color.rgb = DARK
        p.space_after = Pt(8)


def add_kpi_cards(slide, kpis: list[tuple[str, str, str, RGBColor]]):
    w = Inches(3.8)
    h = Inches(2.2)
    y = Inches(2.0)
    xs = [Inches(0.6), Inches(4.7), Inches(8.8)]
    for (label, value, status, color), x in zip(kpis, xs):
        card = slide.shapes.add_shape(1, x, y, w, h)
        card.fill.solid()
        card.fill.fore_color.rgb = RGBColor(248, 248, 248)
        card.line.color.rgb = color
        card.line.width = Pt(3)
        tb = slide.shapes.add_textbox(x + Inches(0.2), y + Inches(0.25), w - Inches(0.4), h)
        tf = tb.text_frame
        tf.text = label
        tf.paragraphs[0].font.size = Pt(14)
        tf.paragraphs[0].font.color.rgb = GRAY
        p2 = tf.add_paragraph()
        p2.text = value
        p2.font.size = Pt(32)
        p2.font.bold = True
        p2.font.color.rgb = DARK
        p3 = tf.add_paragraph()
        p3.text = status
        p3.font.size = Pt(12)
        p3.font.bold = True
        p3.font.color.rgb = color


def add_image(slide, path: Path, left, top, width):
    if path.exists():
        slide.shapes.add_picture(str(path), left, top, width=width)


def add_table_slide(slide, headers: list[str], rows: list[list[str]], y=Inches(1.3)):
    cols, row_count = len(headers), len(rows) + 1
    tbl = slide.shapes.add_table(row_count, cols, Inches(0.5), y, Inches(12.3), Inches(0.45 * row_count)).table
    for j, h in enumerate(headers):
        cell = tbl.cell(0, j)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = ORANGE
        for p in cell.text_frame.paragraphs:
            p.font.bold = True
            p.font.color.rgb = WHITE
            p.font.size = Pt(10)
            p.alignment = PP_ALIGN.CENTER
    for i, row in enumerate(rows, 1):
        for j, val in enumerate(row):
            cell = tbl.cell(i, j)
            cell.text = str(val)
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(9)
                p.alignment = PP_ALIGN.CENTER


def build() -> Path:
    OUT.mkdir(parents=True, exist_ok=True)
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # 1 Cover
    s = blank_slide(prs)
    add_orange_bar(s, Inches(0.18))
    box = s.shapes.add_textbox(Inches(0.8), Inches(2.2), Inches(11), Inches(2))
    tf = box.text_frame
    tf.text = "DiDi Global — CX Service Operations"
    tf.paragraphs[0].font.size = Pt(18)
    tf.paragraphs[0].font.color.rgb = GRAY
    p = tf.add_paragraph()
    p.text = "Weekly Performance Report"
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = ORANGE
    p2 = tf.add_paragraph()
    p2.text = "Semana W19  |  Entregable 2  |  CONFIDENCIAL"
    p2.font.size = Pt(16)
    p2.font.color.rgb = DARK

    # 2 Executive KPIs
    s = blank_slide(prs)
    add_orange_bar(s)
    add_title(s, "Executive Summary — KPIs vs Meta")
    add_kpi_cards(
        s,
        [
            ("QA Score", "94.14", "VERDE — Meta 85 (+9.14)", GREEN),
            ("CSAT", "79.95%", "ROJO — Meta 85% (-5.05pp)", RED),
            ("Recontact Rate", "5.83%", "ÁMBAR — Meta 5.44% (+0.39pp)", AMBER),
        ],
    )
    add_body(
        s,
        [
            "Volumen: 2,460 auditorías QA  |  77,266 encuestas CSAT  |  994,591 contactos",
            "Código: Verde = en meta  |  Ámbar = dentro de 5pp  |  Rojo = >5pp bajo meta",
        ],
        y=Inches(4.6),
        size=14,
    )

    # 3 Executive narrative
    s = blank_slide(prs)
    add_orange_bar(s)
    add_title(s, "Executive Summary — Hallazgo Crítico")
    add_body(
        s,
        [
            "QA global supera meta (94.14), impulsado por Live Chat (96.01). Phone queda en ámbar (83.04).",
            "",
            "CSAT global en rojo (79.95%). Live Chat —72% del feedback — arrastra con 77.55% vs Phone 86.26%.",
            "",
            "HALLAZGO #1 PARA LIDERAZGO:",
            "Cluster Order Status & Delays: ~42% de recontactos, CSAT 64–68%, QA verde en chat.",
            "Problema de FCR/resolución, no de cumplimiento de script auditado.",
        ],
        size=17,
    )

    # 4 QA by Channel
    s = blank_slide(prs)
    add_orange_bar(s)
    add_title(s, "QA Analysis — by Channel")
    add_table_slide(
        s,
        ["Canal", "QA Score", "Audits", "Critical Fails", "Estado"],
        [
            ["Phone", "83.04", "355", "48", "ÁMBAR"],
            ["Live Chat", "96.01", "2,105", "62", "VERDE"],
            ["Global", "94.14", "2,460", "110", "VERDE"],
        ],
    )
    add_body(
        s,
        ["Phone: Manejo del Tiempo 37.5% fail rate — presión AHT compromete calidad."],
        y=Inches(3.2),
        size=14,
    )

    # 5 QA Phone attributes
    s = blank_slide(prs)
    add_orange_bar(s)
    add_title(s, "QA Phone — Top Atributos Fallidos")
    add_table_slide(
        s,
        ["Atributo", "Fail Rate", "Critical", "Impacto"],
        [
            ["Manejo del Tiempo", "37.5%", "No", "Mayor driver de score bajo"],
            ["Información Completa y Correcta", "7.9%", "Sí", "Score automático 0"],
            ["Compensaciones/Reembolsos", "3.4%", "Sí", "Casos de entrega/reembolso"],
            ["Negación de Servicio", "2.3%", "Sí", "Critical fail"],
        ],
    )
    add_image(s, CHARTS / "02_pareto_attributes.png", Inches(0.5), Inches(3.5), Inches(12))

    # 6 QA Live Chat
    s = blank_slide(prs)
    add_orange_bar(s)
    add_title(s, "QA Live Chat — Top Atributos Fallidos")
    add_table_slide(
        s,
        ["Atributo", "Fail Rate", "Critical"],
        [
            ["Saludo e Identificación", "4.6%", "No"],
            ["Actitud de Servicio", "3.3%", "No"],
            ["Disponibilidad del Servicio", "2.2%", "Sí"],
            ["Calidad del Sondeo", "1.5%", "No"],
        ],
    )
    add_image(s, CHARTS / "03_histogram_qa_scores.png", Inches(0.5), Inches(3.4), Inches(12))

    # 7 QA by CR Lv4
    s = blank_slide(prs)
    add_orange_bar(s)
    add_title(s, "QA Analysis — by CR Lv4 (Underperformers)")
    add_table_slide(
        s,
        ["CR Lv4", "QA Score", "Audits", "Hipótesis"],
        [
            ["Order completed not received - FS", "68.2", "49", "Critical info/compensación"],
            ["Order active but received", "65.8", "12", "Tracking + negación servicio"],
            ["Courier overcharged me", "68.6", "7", "Compensación incorrecta"],
            ["Refund status and conditions", "76.4", "25", "Info incompleta"],
            ["User disagrees cancellation charge", "80.0", "19", "Política no comunicada"],
        ],
    )

    # 8 Control chart
    s = blank_slide(prs)
    add_orange_bar(s)
    add_title(s, "Herramienta 4 — Gráfico de Control QA")
    add_image(s, CHARTS / "04_control_chart_qa.png", Inches(0.4), Inches(1.1), Inches(12.5))

    # 9 CSAT overview
    s = blank_slide(prs)
    add_orange_bar(s)
    add_title(s, "CSAT / VOC Analysis — Overview")
    add_table_slide(
        s,
        ["Dimensión", "CSAT %", "Feedback", "Estado"],
        [
            ["Global", "79.95%", "77,266", "ROJO"],
            ["Phone", "86.26%", "21,304", "VERDE"],
            ["Live Chat", "77.55%", "55,962", "ROJO"],
            ["Food", "80.74%", "46,071", "ÁMBAR"],
            ["Full Service", "79.67%", "27,113", "ROJO"],
        ],
    )

    # 10 CSAT worst CRs
    s = blank_slide(prs)
    add_orange_bar(s)
    add_title(s, "CSAT — CR Lv4 con Peor Desempeño")
    add_table_slide(
        s,
        ["CR Lv4", "CSAT %", "Feedback"],
        [
            ["order status & delays", "64.68%", "5,096"],
            ["user request order status", "67.85%", "9,918"],
            ["User disagrees cancellation charge", "67.44%", "5,992"],
            ["refund status and conditions", "67.03%", "3,582"],
            ["membership program renewal", "28.63%", "248"],
        ],
    )
    add_body(
        s,
        [
            "VOC: 'Sin solución', 'espere/confiemos que llegue', 'no me regresaron el dinero'",
        ],
        y=Inches(4.8),
        size=14,
    )

    # 11 Scatter
    s = blank_slide(prs)
    add_orange_bar(s)
    add_title(s, "Herramienta 5 — QA vs CSAT (Correlación)")
    add_image(s, CHARTS / "05_scatter_qa_csat.png", Inches(0.4), Inches(1.1), Inches(12.5))
    add_body(
        s,
        ["QA alto no garantiza CSAT alto en CRs de alto volumen operacional (order status, refunds)."],
        y=Inches(6.5),
        size=13,
    )

    # 12 Recontact
    s = blank_slide(prs)
    add_orange_bar(s)
    add_title(s, "Recontact Analysis")
    add_table_slide(
        s,
        ["Canal / CR", "Rate %", "Volumen", "Estado"],
        [
            ["Global", "5.83%", "57,976 rc", "ÁMBAR"],
            ["Live Chat", "15.99%", "38,961 rc", "ROJO"],
            ["user request order status", "16.92%", "130,014 rc", "ROJO"],
            ["order status & delays", "19.34%", "7,680 rc", "ROJO"],
            ["User disagrees cancellation", "12.96%", "6,599 rc", "ROJO"],
        ],
    )

    # 13 Pareto recontact
    s = blank_slide(prs)
    add_orange_bar(s)
    add_title(s, "Herramienta 2 — Pareto Recontact")
    add_image(s, CHARTS / "08_pareto_recontact.png", Inches(0.4), Inches(1.1), Inches(12.5))

    # 14 Combined
    s = blank_slide(prs)
    add_orange_bar(s)
    add_title(s, "Combined Analysis — Insight Integrado")
    add_body(
        s,
        [
            "CR Cluster: Order Status + Delay Information + Cancellation Charge",
            "",
            "• ~24,293 recontactos (41.9% del total semanal)",
            "• CSAT 64–68% en los 3 CRs principales",
            "• QA Live Chat 95–97 (verde) — desalineación percepción vs auditoría",
            "",
            "Root cause: agentes cumplen script pero no resuelven en 1ra interacción.",
            "Falta visibilidad operativa (tracking, compensación automática).",
        ],
        size=17,
    )

    # 15 Fishbone
    s = blank_slide(prs)
    add_orange_bar(s)
    add_title(s, "Herramienta 6 — Ishikawa (Causa-Raíz)")
    add_image(s, CHARTS / "06_fishbone_order_status.png", Inches(0.3), Inches(1.0), Inches(12.7))

    # 16 Flowchart
    s = blank_slide(prs)
    add_orange_bar(s)
    add_title(s, "Herramienta 7 — Flujo FCR (Punto de Falla)")
    add_image(s, CHARTS / "07_flowchart_fcr.png", Inches(0.3), Inches(1.0), Inches(12.7))

    # 17 Check sheet
    s = blank_slide(prs)
    add_orange_bar(s)
    add_title(s, "Herramienta 1 — Check Sheet")
    add_image(s, CHARTS / "01_check_sheet.png", Inches(0.3), Inches(1.0), Inches(12.7))

    # 18 Action plans
    s = blank_slide(prs)
    add_orange_bar(s)
    add_title(s, "Action Plans — by LOB / Business Type", size=24)
    add_table_slide(
        s,
        ["LOB", "Responsable", "Plazo", "Acción"],
        [
            ["Food", "QA Ops + Training", "Sem 1–2", "Coaching Phone: Manejo Tiempo + Info Completa"],
            ["Full Service", "Ops + WFM", "Sem 1", "Staffing horas pico; reducir wait chat"],
            ["Full Service", "QA + Content", "Sem 2", "Macros compensación 'order not received'"],
            ["Market Place", "LOB Lead", "Sem 2–3", "Scripts separados marketplace vs FS"],
            ["All LOBs", "CX Leadership", "Sem 1", "War room top 3 CRs recontact hasta <5.44%"],
        ],
        y=Inches(1.1),
    )

    # 19 7 tools summary
    s = blank_slide(prs)
    add_orange_bar(s)
    add_title(s, "Anexo — 7 Herramientas de Calidad", size=24)
    add_table_slide(
        s,
        ["#", "Herramienta", "Aplicación"],
        [
            ["1", "Check Sheet", "Validación datos QA y reglas Critical"],
            ["2", "Pareto", "80/20 defectos QA y recontact"],
            ["3", "Histograma", "Distribución QA Score por canal"],
            ["4", "Control Chart", "Variación QA país/canal vs meta"],
            ["5", "Scatter", "Correlación QA vs CSAT por CR"],
            ["6", "Ishikawa", "Causa-raíz Order Status/Delays"],
            ["7", "Flowchart", "Punto de falla FCR"],
        ],
    )

    # 20 Closing
    s = blank_slide(prs)
    add_orange_bar(s, Inches(0.18))
    box = s.shapes.add_textbox(Inches(1), Inches(2.8), Inches(11), Inches(2))
    tf = box.text_frame
    tf.text = "Gracias"
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].font.size = Pt(44)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = ORANGE
    p = tf.add_paragraph()
    p.text = "DiDi CX Quality Analyst — Business Case"
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(16)
    p.font.color.rgb = GRAY

    out_path = OUT / "Entregable_2_Canva_Import.pptx"
    prs.save(out_path)
    return out_path


if __name__ == "__main__":
    path = build()
    print(f"Canva-ready deck: {path}")
