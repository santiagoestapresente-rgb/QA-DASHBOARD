"""
Entregable 2 — landscape 16:9 slide deck (PDF + PPTX).

Visual language of the DiDi CX briefing (black header, orange accent, split
data/insights layouts) with the numbered section structure of the weekly
performance dashboard. Numbers come from packaged parquet via the same KPI
functions as the Streamlit app. Official formulas are not changed.
"""
from __future__ import annotations

import json
import sys
from datetime import datetime
from pathlib import Path

import matplotlib

matplotlib.use("Agg")
import matplotlib.pyplot as plt
import numpy as np
import pandas as pd
from reportlab.lib import colors
from reportlab.lib.units import inch
from reportlab.pdfbase import pdfmetrics
from reportlab.pdfbase.ttfonts import TTFont
from reportlab.pdfgen import canvas as pdfcanvas

ROOT = Path(__file__).resolve().parent.parent
sys.path.insert(0, str(ROOT))

from config import CONTROL_TOTALS, CSAT_GOAL, QA_GOAL, RECONTACT_GOAL  # noqa: E402
from modules.data_loader import load_all_data  # noqa: E402
from modules.executive_engine import combined_operational_analysis, qa_channel_breakdown  # noqa: E402
from modules.kpis import (  # noqa: E402
    _vs_goal_status,
    channel_performance,
    csat_by_business_type,
    channel_match,
    csat_by_star_rating,
    csat_unsatisfied_by_cr,
    kpi_summary,
    recontact_by_cr,
    recontact_by_scope,
    recontact_dilution_stats,
    recontact_rate,
    top_failing_attributes,
    voc_themes_negative,
)

OUT = ROOT / "entregable 2"
CHARTS = OUT / "charts" / "slides"

# ── Brand ──────────────────────────────────────────────────────────────────
C_ORANGE = "#FF6600"
C_DARK = "#1A1A1A"
C_WHITE = "#FFFFFF"
C_GREEN = "#2E9B57"
C_AMBER = "#F2A900"
C_RED = "#D64545"
C_GRAY = "#5C5C5C"
C_MUTED = "#8A8A8A"
C_LINE = "#E6E6E6"
C_PANEL = "#FFF6F0"
C_CARD = "#F7F7F7"
C_ROW = "#FAFAFA"
STATUS_HEX = {"green": C_GREEN, "amber": C_AMBER, "red": C_RED, "neutral": C_MUTED}
STATUS_PILL = {"green": "AT GOAL", "amber": "WITHIN 5pp", "red": "OFF GOAL", "neutral": "—"}
STATUS_SHORT = {"green": "GREEN", "amber": "AMBER", "red": "RED", "neutral": "—"}
LIGHT_FILL = {"green": "#E7F5EE", "amber": "#FFF6D9", "red": "#FDECEC", "neutral": C_WHITE}

SLIDE_W = 13.333 * inch
SLIDE_H = 7.5 * inch
HEADER_H = 0.46 * inch
ACCENT_H = 0.055 * inch
ML = 0.34 * inch
MR = 0.34 * inch
FOOTER_Y = 0.13 * inch
CONTENT_TOP = SLIDE_H - HEADER_H - ACCENT_H - 0.10 * inch
CONTENT_BOTTOM = 0.40 * inch
GAP = 0.16 * inch
USABLE = SLIDE_W - ML - MR
LEFT_W = USABLE * 0.62
RIGHT_W = USABLE - LEFT_W - GAP
LEFT_X = ML
RIGHT_X = ML + LEFT_W + GAP

HEADER_TEXT = "DiDi | CX Quality Analyst — Weekly Performance Report"
FOOTER_LEFT = "CONFIDENTIAL"
FOOTER_CENTER = "DiDi Global — CX Service Operations | Internal Use Only"
BRIEFING = datetime(2026, 8, 20)

plt.rcParams.update({
    "font.family": "sans-serif",
    "font.sans-serif": ["Segoe UI", "Arial", "DejaVu Sans"],
    "axes.grid": False,
    "axes.edgecolor": "#D4D4D4",
    "axes.linewidth": 0.6,
    "figure.facecolor": "white",
    "axes.facecolor": "white",
    "xtick.color": C_GRAY,
    "ytick.color": C_GRAY,
    "text.color": C_DARK,
    "axes.labelcolor": C_GRAY,
    "legend.frameon": False,
})


def _register_fonts() -> tuple[str, str]:
    regular = Path(r"C:\Windows\Fonts\segoeui.ttf")
    bold = Path(r"C:\Windows\Fonts\segoeuib.ttf")
    if regular.exists() and bold.exists():
        pdfmetrics.registerFont(TTFont("Brand", str(regular)))
        pdfmetrics.registerFont(TTFont("Brand-Bold", str(bold)))
        return "Brand", "Brand-Bold"
    return "Helvetica", "Helvetica-Bold"


FONT, FONT_BOLD = _register_fonts()


def hx(code: str):
    return colors.HexColor(code)


def status(value: float, goal: float, higher: bool = True) -> str:
    return _vs_goal_status(value, goal, higher)


def ellipsize(text, n: int = 46) -> str:
    s = " ".join(str(text).split())
    return s if len(s) <= n else s[: n - 1] + "…"


def fmt_n(n) -> str:
    return f"{int(n):,}"


def fmt_pp(v: float, digits: int = 2) -> str:
    sign = "+" if v >= 0 else "−"
    return f"{sign}{abs(v):.{digits}f}"


def wrap_lines(text: str, font: str, size: float, max_w: float) -> list[str]:
    words = str(text).split()
    if not words:
        return [""]
    lines, cur = [], ""
    for w in words:
        trial = (cur + " " + w).strip()
        if pdfmetrics.stringWidth(trial, font, size) <= max_w:
            cur = trial
        else:
            if cur:
                lines.append(cur)
            cur = w
    if cur:
        lines.append(cur)
    return lines or [""]


def y_from_top(offset: float) -> float:
    return SLIDE_H - offset


# ── Context ────────────────────────────────────────────────────────────────

def _fail_rows(errors: pd.DataFrame, audits: pd.DataFrame, top_n: int = 3) -> list[dict]:
    n = len(audits)
    top = top_failing_attributes(errors, audits, top_n=top_n)
    rows = []
    if top.empty or n == 0:
        return rows
    for _, r in top.iterrows():
        uniq = int(errors.loc[errors["Error_Category"] == r["Error_Category"], "Audit_ID"].nunique())
        rows.append({
            "attr": str(r["Error_Category"]),
            "fails": int(r["Fail_Count"]),
            "unique": uniq,
            "rate": round(uniq / n * 100, 1),
            "critical": bool(r["Is_Critical"]),
            "pct_fails": float(r["Pct_Of_Fails"]),
        })
    return rows


def _cr_bundle(audits, csat, recontact, name: str) -> dict:
    key = name.casefold()
    qa = audits[audits["CR_Lv4"].astype(str).str.casefold() == key]
    cs = csat[csat["CR_Lv4"].astype(str).str.casefold() == key]
    rc = recontact[recontact["CR_Lv4"].astype(str).str.casefold() == key]
    qa_n = int(len(qa))
    qa_v = float(qa["Score_Pct"].mean()) if qa_n else float("nan")
    fb = float(cs["Feedback CNT"].sum()) if not cs.empty else 0.0
    sat = float(cs["Satisfied_CNT"].sum()) if not cs.empty else 0.0
    csat_v = sat / fb * 100 if fb else float("nan")
    contacts = float(rc["Contacts"].sum()) if not rc.empty else 0.0
    rvol = float(rc["Recontact Volume"].sum()) if not rc.empty else 0.0
    rc_v = rvol / contacts * 100 if contacts else float("nan")
    return {
        "name": name,
        "qa": qa_v,
        "qa_n": qa_n,
        "csat": csat_v,
        "feedback": int(fb),
        "rc": rc_v,
        "contacts": int(contacts),
        "recontacts": int(rvol),
    }


def _qa_under_n(audits: pd.DataFrame, min_n: int = 10) -> pd.DataFrame:
    g = (
        audits.groupby("CR_Lv4")
        .agg(QA_Score=("Score_Pct", "mean"), N=("Audit_ID", "count"))
        .reset_index()
    )
    g = g[g["N"] >= min_n].copy()
    g["QA_Score"] = g["QA_Score"].round(2)
    g["vs"] = (g["QA_Score"] - QA_GOAL).round(2)
    g["status"] = g["QA_Score"].apply(lambda v: status(v, QA_GOAL, True))
    return g.sort_values(["QA_Score", "N"], ascending=[True, False]).reset_index(drop=True)


def load_ctx() -> dict:
    data = load_all_data()
    audits, errors, csat, rc = (
        data["fact_audits"], data["fact_errors"], data["fact_csat"], data["fact_recontact"],
    )
    summary = kpi_summary(audits, csat, rc)
    rc_rate = float(recontact_rate(rc))
    ch_break = qa_channel_breakdown(audits, errors)
    ch_perf = channel_performance(audits, csat, rc)
    scopes = recontact_by_scope(rc)
    dilution = recontact_dilution_stats(rc)
    bt = csat_by_business_type(csat)
    stars = csat_by_star_rating(csat)
    voc = voc_themes_negative(csat, top_n=6)
    rc_cr = recontact_by_cr(rc, top_n=8)
    unsat = csat_unsatisfied_by_cr(csat)
    combined = combined_operational_analysis(audits, csat, rc)
    qa_cr = _qa_under_n(audits, 10)

    phone_a = audits[audits["Channel"] == "Phone"]
    chat_a = audits[audits["Channel"] == "Live Chat"]
    phone_e = errors[errors["Channel"] == "Phone"] if "Channel" in errors.columns else errors
    chat_e = errors[errors["Channel"] == "Live Chat"] if "Channel" in errors.columns else errors

    def ch_row(name: str) -> pd.Series:
        sub = ch_perf[ch_perf["Segment"] == name]
        return sub.iloc[0] if not sub.empty else pd.Series(dtype=float)

    phone = ch_row("Phone")
    chat = ch_row("Live Chat")
    phone_qa = float(phone["QA_Score"])
    chat_qa = float(chat["QA_Score"])
    phone_csat = float(phone["CSAT_Score"])
    chat_csat = float(chat["CSAT_Score"])
    phone_rc = float(phone["Recontact_Rate"])
    chat_rc = float(chat["Recontact_Rate"])

    def scope(key: str) -> pd.Series:
        sub = scopes[scopes["Scope_Key"] == key]
        return sub.iloc[0] if not sub.empty else pd.Series(dtype=float)

    official = scope("all")
    ex_sh = scope("ex_self_help")
    audited = scope("audited")

    cluster = [
        _cr_bundle(audits, csat, rc, "order status & delays"),
        _cr_bundle(audits, csat, rc, "User disagrees with cancellation charge/debt"),
        _cr_bundle(audits, csat, rc, "refund status and conditions"),
        _cr_bundle(audits, csat, rc, "user request order status or delay information"),
    ]

    weeks = sorted(audits["Week"].dropna().astype(str).unique().tolist())
    d0 = pd.to_datetime(audits["Fecha"]).min()
    d1 = pd.to_datetime(audits["Fecha"]).max()
    period = f"{weeks[0]}–{weeks[-1]}" if weeks else "May 2026"
    date_span = f"{int(d0.day)}–{int(d1.day)} May 2026"

    qa = float(summary["qa_score"])
    csat_v = float(summary["csat"])
    fcr = 100.0 - rc_rate
    surveys = int(csat["Feedback CNT"].sum())
    contacts = int(rc["Contacts"].sum())
    recontacts = int(rc["Recontact Volume"].sum())
    evals = int(len(audits))
    chat_surveys = int(csat.loc[channel_match(csat["Channel"], "Live Chat"), "Feedback CNT"].sum())
    phone_surveys = int(csat.loc[channel_match(csat["Channel"], "Phone"), "Feedback CNT"].sum())

    ctx = {
        "qa": qa, "csat": csat_v, "rc": rc_rate, "fcr": fcr,
        "qa_st": status(qa, QA_GOAL, True),
        "csat_st": status(csat_v, CSAT_GOAL, True),
        "rc_st": status(rc_rate, RECONTACT_GOAL, False),
        "evals": evals, "surveys": surveys, "contacts": contacts, "recontacts": recontacts,
        "phone_qa": phone_qa, "chat_qa": chat_qa,
        "phone_csat": phone_csat, "chat_csat": chat_csat,
        "phone_rc": phone_rc, "chat_rc": chat_rc,
        "phone_n": int(phone["QA_N"]), "chat_n": int(chat["QA_N"]),
        "phone_surveys": phone_surveys, "chat_surveys": chat_surveys,
        "phone_qa_st": status(phone_qa, QA_GOAL, True),
        "chat_qa_st": status(chat_qa, QA_GOAL, True),
        "phone_csat_st": status(phone_csat, CSAT_GOAL, True),
        "chat_csat_st": status(chat_csat, CSAT_GOAL, True),
        "phone_rc_st": status(phone_rc, RECONTACT_GOAL, False),
        "chat_rc_st": status(chat_rc, RECONTACT_GOAL, False),
        "phone_fatal": float(ch_break["Phone"]["pct_fatal"]),
        "chat_fatal": float(ch_break["Live Chat"]["pct_fatal"]),
        "phone_crit": int(ch_break["Phone"]["n_crit_fails"]),
        "chat_crit": int(ch_break["Live Chat"]["n_crit_fails"]),
        "phone_attrs": _fail_rows(phone_e, phone_a, 3),
        "chat_attrs": _fail_rows(chat_e, chat_a, 3),
        "phone_cr": qa_score_channel(phone_a, 10),
        "chat_cr": qa_score_channel(chat_a, 10),
        "official": official, "ex_sh": ex_sh, "audited": audited, "dilution": dilution,
        "bt": bt, "stars": stars, "voc": voc, "rc_cr": rc_cr, "unsat": unsat,
        "combined": combined, "qa_cr": qa_cr, "cluster": cluster,
        "period": period, "date_span": date_span, "weeks": weeks,
        "lob": audits["LOB"].value_counts().to_dict() if "LOB" in audits.columns else {},
        "ch_break": ch_break,
    }
    ctx["copy"] = build_copy(ctx)
    return ctx


def qa_score_channel(audits: pd.DataFrame, min_n: int = 8) -> pd.DataFrame:
    if audits.empty:
        return pd.DataFrame()
    g = (
        audits.groupby("CR_Lv4")
        .agg(QA_Score=("Score_Pct", "mean"), N=("Audit_ID", "count"))
        .reset_index()
    )
    g = g[g["N"] >= min_n].copy()
    g["QA_Score"] = g["QA_Score"].round(2)
    g["status"] = g["QA_Score"].apply(lambda v: status(v, QA_GOAL, True))
    return g.sort_values("QA_Score").head(4).reset_index(drop=True)


def build_copy(ctx: dict) -> dict:
    sh_share = float(ctx["dilution"]["share"])
    sh_rate = float(ctx["dilution"]["rate"])
    ex = ctx["ex_sh"]
    aud = ctx["audited"]
    chat_pct = ctx["chat_surveys"] / ctx["surveys"] * 100 if ctx["surveys"] else 0
    finding = (
        f"CSAT is the miss to manage: {ctx['csat']:.2f}% vs {CSAT_GOAL:.0f}% "
        f"(RED, {fmt_pp(ctx['csat'] - CSAT_GOAL)} pp) on {fmt_n(ctx['surveys'])} surveys. "
        f"Live Chat is {chat_pct:.0f}% of feedback at {ctx['chat_csat']:.2f}% CSAT (RED) while Phone is "
        f"{ctx['phone_csat']:.2f}% (GREEN). Official recontact {ctx['rc']:.2f}% is only AMBER because "
        f"Self Help is {sh_share:.0f}% of contacts at {sh_rate:.2f}% — excluding it the rate is "
        f"{ex['Rate']:.2f}%, and Phone+Chat is {aud['Rate']:.2f}%. Dissatisfaction concentrates in "
        f"order status/delay, cancellation charge/debt, and refund."
    )
    return {
        "finding": finding,
        "exec": [
            f"Global QA {ctx['qa']:.2f} (GREEN, {fmt_pp(ctx['qa'] - QA_GOAL)} vs {QA_GOAL:.0f}) is not the operating picture. Phone is {ctx['phone_qa']:.2f} (AMBER) on {fmt_n(ctx['phone_n'])} evaluations; Live Chat {ctx['chat_qa']:.2f} (GREEN) on {fmt_n(ctx['chat_n'])}.",
            f"Do not celebrate derived FCR of {ctx['fcr']:.2f}% (100 − recontact). There is no FCR target, and the official mix is diluted by Self Help.",
            "QA sample is Delivery LOB only. Action plans run by Business Type (Food, Full Service, Market Place). Pickup n=35 is excluded.",
            "Priority: Chat outcome on status/money CRs + Phone coaching. Raising audit score will not close the CSAT gap.",
        ],
        "qa_ch": [
            "Phone fail pattern is operational: Time management 37.5% of audits, then Complete and correct information (critical, 7.9%) which zeros the score, then User name 4.2%.",
            "Chat fail pattern is protocol: Greeting 4.6%, Service attitude 3.3%, Service availability (critical, 2.2%). Score stays high because most fails are non-critical.",
            "Hypothesis: Phone AHT pressure trades completeness for speed. Chat agents pass the script while customers wait on order/money outcomes the agent cannot close.",
            "Do not average the two channels. Report Phone and Live Chat separately every week.",
        ],
        "qa_cr": [
            "Every QA underperformer with n≥10 is Phone-heavy (delivery not received, already received, refund). Live Chat CRs with n≥8 all sit at or above 85.",
            "Coaching queue: order completed not received — FS (68.16, n=49) and order active but already received (65.83, n=12).",
            "Refund status at 76.40 (n=25) is the money-path QA miss — same family as the CSAT/recontact cluster, but here QA actually fails.",
            "N is the evaluation count. Do not read a weighted gap (gap × n) as audits.",
        ],
        "csat": [
            "Star mix: 75.7% five-star but 16.5% one-star. The 1-star block is large enough to keep CSAT RED.",
            "Food 80.74% (AMBER, 46,071) and Full Service 79.67% (RED, 27,113) carry the miss. Market Place 80.34% AMBER. Other 26.52% on 558 surveys is a catch-all CR, not a LOB program.",
            "Unsatisfied volume concentrates in order status, cancellation charge/debt, and refund — not in damaged/wrong order, where CSAT is already GREEN.",
            "VOC on 1-3 star comments: refund/compensation not received and no solution provided dominate tagged mentions. Outcome, not greeting.",
        ],
        "rc": [
            f"Official {ctx['rc']:.2f}% (AMBER, {fmt_pp(ctx['rc'] - RECONTACT_GOAL)} pp) is the 12-channel ratio of sums. It is not the live-agent rate.",
            f"Excluding Self Help: {ex['Rate']:.2f}% on {fmt_n(ex['Contacts'])} contacts (RED). Phone+Chat: {aud['Rate']:.2f}% on {fmt_n(aud['Contacts'])} (RED).",
            "Volume leaders: user request order status (22.4% of recontacts, rate 16.92%), order status & delays (19.34%), cancellation charge (12.96%).",
            "“User don't want the order anymore” is 10.7% of recontact volume at only 2.20% rate — volume without a rate problem. Do not treat it like status/money.",
        ],
        "combined": [
            "The status/money cluster shows HIGH QA + LOW CSAT + HIGH recontact. Agents pass the audit; the customer does not get a closed outcome.",
            "This is process (script) versus outcome (resolution). Greeting/attitude coaching will not move CSAT on these CRs.",
            "Fix tracking visibility, compensation/refund path, and first-contact close. Measure CSAT and recontact on these four CRs weekly.",
            f"Derived FCR {ctx['fcr']:.2f}% is 100 − official recontact. Do not put it on the exec scorecard.",
        ],
        "action": [
            "Sequence: (1) war room on status/money CRs from 20 Aug, (2) Phone attribute coaching by 3 Sep, (3) reclassify Other tickets.",
            "30-day success: Chat CSAT on status/money CRs moving toward 85%; Phone QA ≥85; Phone+Chat recontact published next to official 5.83%.",
            "QA staying at 94 is not a 30-day objective — it is already GREEN. Do not trade Chat CSAT work for more script audits.",
            "Pickup (n=35 surveys) is out of scope. QA LOB is Delivery only.",
        ],
    }


# ── Charts ─────────────────────────────────────────────────────────────────

def _style_ax(ax, title=""):
    ax.spines["top"].set_visible(False)
    ax.spines["right"].set_visible(False)
    ax.spines["left"].set_color("#D4D4D4")
    ax.spines["bottom"].set_color("#D4D4D4")
    ax.tick_params(length=0)
    if title:
        ax.set_title(title, loc="left", fontsize=10, fontweight="bold", color=C_DARK, pad=6)


def save_chart(fig, path: Path):
    path.parent.mkdir(parents=True, exist_ok=True)
    fig.savefig(path, dpi=180, facecolor="white", bbox_inches="tight", pad_inches=0.12)
    plt.close(fig)


def build_charts(ctx: dict) -> dict[str, Path]:
    CHARTS.mkdir(parents=True, exist_ok=True)
    paths = {}

    # Channel QA vs CSAT
    fig, ax = plt.subplots(figsize=(6.4, 2.35))
    labels = ["Phone", "Live Chat"]
    qa_v = [ctx["phone_qa"], ctx["chat_qa"]]
    cs_v = [ctx["phone_csat"], ctx["chat_csat"]]
    x = np.arange(len(labels))
    w = 0.34
    b1 = ax.bar(x - w / 2, qa_v, w, color=[STATUS_HEX[ctx["phone_qa_st"]], STATUS_HEX[ctx["chat_qa_st"]]], label="QA")
    b2 = ax.bar(x + w / 2, cs_v, w, color="#2E6FBE", label="CSAT %")
    ax.axhline(QA_GOAL, color=C_DARK, ls="--", lw=0.9)
    ax.set_xticks(x)
    ax.set_xticklabels(labels, fontsize=9)
    ax.set_ylim(0, 110)
    ax.legend(fontsize=8, loc="lower right")
    _style_ax(ax, "QA vs CSAT by channel")
    for bars in (b1, b2):
        for rect in bars:
            h = rect.get_height()
            ax.text(rect.get_x() + rect.get_width() / 2, h + 1.4, f"{h:.1f}", ha="center", fontsize=7.5, color=C_DARK)
    paths["channel"] = CHARTS / "channel_qa_csat.png"
    save_chart(fig, paths["channel"])

    # Star doughnut
    fig, ax = plt.subplots(figsize=(3.35, 2.55))
    stars = ctx["stars"]
    star_colors = ["#2E9B57", "#7BC67E", "#F2A900", "#FF6600", "#D64545"]
    sizes = stars["Count"].tolist()
    ax.pie(
        sizes, colors=star_colors, startangle=90,
        wedgeprops=dict(width=0.42, edgecolor="white", linewidth=1.2),
    )
    ax.text(0, 0.08, f"{ctx['csat']:.2f}%", ha="center", va="center", fontsize=13, fontweight="bold", color=C_DARK)
    ax.text(0, -0.18, "CSAT", ha="center", va="center", fontsize=8, color=C_GRAY)
    ax.set_aspect("equal")
    paths["stars"] = CHARTS / "csat_stars.png"
    save_chart(fig, paths["stars"])

    # Business type
    fig, ax = plt.subplots(figsize=(5.9, 2.15))
    bt = ctx["bt"][ctx["bt"]["Business_Type"] != "Pickup"].copy()
    bt = bt.sort_values("CSAT_Score")
    cols = [STATUS_HEX[status(v, CSAT_GOAL, True)] for v in bt["CSAT_Score"]]
    ax.barh(bt["Business_Type"], bt["CSAT_Score"], color=cols, height=0.55)
    ax.axvline(CSAT_GOAL, color=C_DARK, ls="--", lw=0.9)
    for _, r in bt.iterrows():
        ax.text(r["CSAT_Score"] + 0.6, r["Business_Type"], f"{r['CSAT_Score']:.1f}%", va="center", fontsize=7.5)
    ax.set_xlim(0, 100)
    _style_ax(ax, "CSAT by Business Type (Pickup n=35 omitted)")
    paths["bt"] = CHARTS / "csat_bt.png"
    save_chart(fig, paths["bt"])

    # VOC
    fig, ax = plt.subplots(figsize=(5.5, 2.35))
    voc = ctx["voc"].sort_values("Mentions")
    ax.barh(voc["Theme"], voc["Mentions"], color=C_ORANGE, height=0.55)
    for _, r in voc.iterrows():
        ax.text(r["Mentions"] + 8, r["Theme"], f"{int(r['Mentions'])}", va="center", fontsize=7)
    _style_ax(ax, "Negative VOC themes (1-3 star comments)")
    paths["voc"] = CHARTS / "voc.png"
    save_chart(fig, paths["voc"])

    # Recontact scope
    fig, ax = plt.subplots(figsize=(6.2, 2.15))
    names = ["Official\n12 channels", "Ex-Self Help", "Phone + Chat"]
    vals = [float(ctx["official"]["Rate"]), float(ctx["ex_sh"]["Rate"]), float(ctx["audited"]["Rate"])]
    st = [status(v, RECONTACT_GOAL, False) for v in vals]
    ax.bar(names, vals, color=[STATUS_HEX[s] for s in st], width=0.55)
    ax.axhline(RECONTACT_GOAL, color=C_DARK, ls="--", lw=0.9)
    for i, v in enumerate(vals):
        ax.text(i, v + 0.25, f"{v:.2f}%", ha="center", fontsize=8, fontweight="bold")
    ax.set_ylim(0, max(vals) * 1.25)
    _style_ax(ax, "Recontact rate by scope (ratio of sums)")
    paths["rc_scope"] = CHARTS / "rc_scope.png"
    save_chart(fig, paths["rc_scope"])

    # Recontact Pareto
    fig, ax = plt.subplots(figsize=(6.4, 2.55))
    rc = ctx["rc_cr"].head(6).iloc[::-1]
    ax.barh([ellipsize(x, 36) for x in rc["CR_Lv4"]], rc["Recontacts"], color=C_ORANGE, height=0.55)
    for _, r in rc.iterrows():
        ax.text(r["Recontacts"] + 80, ellipsize(r["CR_Lv4"], 36), f"{int(r['Recontacts']):,}", va="center", fontsize=7)
    _style_ax(ax, "Top CR Lv4 by recontact volume")
    paths["rc_pareto"] = CHARTS / "rc_pareto.png"
    save_chart(fig, paths["rc_pareto"])

    return paths


# ── PDF primitives ─────────────────────────────────────────────────────────

def draw_header(c, title_slide=False):
    if title_slide:
        return
    c.setFillColor(hx(C_DARK))
    c.rect(0, SLIDE_H - HEADER_H, SLIDE_W, HEADER_H, fill=1, stroke=0)
    c.setFillColor(hx(C_ORANGE))
    c.rect(0, SLIDE_H - HEADER_H - ACCENT_H, SLIDE_W, ACCENT_H, fill=1, stroke=0)
    c.setFillColor(hx(C_WHITE))
    c.setFont(FONT_BOLD, 10)
    c.drawString(ML, SLIDE_H - HEADER_H / 2 - 4, HEADER_TEXT)
    c.setFont(FONT, 8)
    c.setFillColor(hx("#CFCFCF"))
    c.drawRightString(SLIDE_W - MR, SLIDE_H - HEADER_H / 2 - 3.5, "May 2026  ·  Market All")


def draw_footer(c, page: int, total: int, title_slide=False):
    fg = C_WHITE if title_slide else C_DARK
    muted = "#FFE0CC" if title_slide else C_MUTED
    c.setStrokeColor(hx(C_ORANGE if not title_slide else C_WHITE))
    c.setLineWidth(0.9)
    c.line(ML, 0.30 * inch, SLIDE_W - MR, 0.30 * inch)
    c.setFillColor(hx(fg))
    c.setFont(FONT_BOLD, 7)
    c.drawString(ML, FOOTER_Y, FOOTER_LEFT)
    c.setFillColor(hx(muted if title_slide else C_GRAY))
    c.setFont(FONT, 7)
    c.drawCentredString(SLIDE_W / 2, FOOTER_Y, FOOTER_CENTER)
    c.setFillColor(hx(fg))
    c.setFont(FONT, 7)
    c.drawRightString(SLIDE_W - MR, FOOTER_Y, f"Page {page} of {total}")


def draw_kicker(c, number: str, title: str, y=None) -> float:
    y = CONTENT_TOP - 2 if y is None else y
    c.setFillColor(hx(C_ORANGE))
    c.setFont(FONT_BOLD, 11)
    c.drawString(ML, y - 2, f"{number}   {title}")
    return y - 18


def pill(c, x, y, text, fill, font_size=6.5, pad_x=5, h=11, right=None):
    c.setFont(FONT_BOLD, font_size)
    tw = pdfmetrics.stringWidth(text, FONT_BOLD, font_size)
    w = tw + pad_x * 2
    if right is not None:
        x = right - w
    c.setFillColor(hx(fill))
    c.roundRect(x, y, w, h, 3, fill=1, stroke=0)
    c.setFillColor(hx(C_WHITE))
    c.drawString(x + pad_x, y + 2.6, text)
    return w


def draw_kpi_card(c, x, y, w, h, label, value, meta, st: str):
    c.setFillColor(hx(C_CARD))
    c.setStrokeColor(hx(C_LINE))
    c.setLineWidth(0.6)
    c.roundRect(x, y, w, h, 3, fill=1, stroke=1)
    c.setFillColor(hx(STATUS_HEX[st]))
    c.rect(x, y, 5, h, fill=1, stroke=0)
    c.setFillColor(hx(C_MUTED))
    c.setFont(FONT, 7.5)
    c.drawString(x + 14, y + h - 16, label.upper())
    pill(c, 0, y + h - 18, STATUS_PILL[st], STATUS_HEX[st], right=x + w - 8)
    c.setFillColor(hx(C_DARK))
    c.setFont(FONT_BOLD, 22)
    c.drawString(x + 14, y + h - 42, value)
    c.setFillColor(hx(C_GRAY))
    c.setFont(FONT, 7.5)
    yy = y + h - 58
    for line in meta:
        c.drawString(x + 14, yy, line)
        yy -= 11


def draw_insights(c, x, y, w, h, bullets: list[str], title="INSIGHTS"):
    c.setFillColor(hx(C_PANEL))
    c.setStrokeColor(hx(C_ORANGE))
    c.setLineWidth(1.1)
    c.rect(x, y, w, h, fill=1, stroke=1)
    c.setFillColor(hx(C_ORANGE))
    c.rect(x, y, 4, h, fill=1, stroke=0)
    c.setFont(FONT_BOLD, 9)
    c.drawString(x + 14, y + h - 16, title)
    ty = y + h - 32
    max_w = w - 28
    for b in bullets:
        lines = wrap_lines(b, FONT, 8, max_w)
        c.setFillColor(hx(C_ORANGE))
        c.circle(x + 12, ty + 3, 2.1, fill=1, stroke=0)
        c.setFillColor(hx(C_DARK))
        c.setFont(FONT, 8)
        for i, line in enumerate(lines):
            c.drawString(x + 20, ty, line)
            ty -= 10.5
        ty -= 6
        if ty < y + 10:
            break


def draw_table(c, x, y_top, widths, header, rows, fills=None, aligns=None, row_h=15.5, font_size=7.5):
    """Draw a table downward from y_top. Returns bottom y."""
    cols = len(widths)
    aligns = aligns or ["LEFT"] * cols
    total_w = sum(widths)
    # header
    y = y_top - row_h
    c.setFillColor(hx(C_DARK))
    c.rect(x, y, total_w, row_h, fill=1, stroke=0)
    c.setFillColor(hx(C_WHITE))
    c.setFont(FONT_BOLD, font_size)
    cx = x
    for i, h in enumerate(header):
        pad = 4
        yy = y + 4.4
        if aligns[i] == "CENTER":
            c.drawCentredString(cx + widths[i] / 2, yy, str(h))
        elif aligns[i] == "RIGHT":
            c.drawRightString(cx + widths[i] - pad, yy, str(h))
        else:
            c.drawString(cx + pad, yy, str(h))
        cx += widths[i]
    for r_i, row in enumerate(rows):
        y -= row_h
        bg = LIGHT_FILL.get((fills or {}).get(r_i, ""), None)
        if bg:
            c.setFillColor(hx(bg))
        else:
            c.setFillColor(hx(C_WHITE if r_i % 2 == 0 else C_ROW))
        c.rect(x, y, total_w, row_h, fill=1, stroke=0)
        c.setStrokeColor(hx(C_LINE))
        c.setLineWidth(0.3)
        c.line(x, y, x + total_w, y)
        cx = x
        for i, val in enumerate(row):
            raw = str(val)
            is_status = raw in ("GREEN", "AMBER", "RED", "AT GOAL", "WITHIN 5pp", "OFF GOAL")
            st_key = None
            if is_status:
                st_key = {"GREEN": "green", "AMBER": "amber", "RED": "red",
                          "AT GOAL": "green", "WITHIN 5pp": "amber", "OFF GOAL": "red"}[raw]
                c.setFillColor(hx(STATUS_HEX[st_key]))
                c.roundRect(cx + 3, y + 2.5, widths[i] - 6, row_h - 5, 2, fill=1, stroke=0)
                c.setFillColor(hx(C_WHITE))
                c.setFont(FONT_BOLD, font_size - 0.5)
                c.drawCentredString(cx + widths[i] / 2, y + 4.4, raw)
            else:
                if fills and r_i in fills and i == 0:
                    pass
                c.setFillColor(hx(C_DARK))
                c.setFont(FONT, font_size)
                pad = 4
                yy = y + 4.4
                txt = ellipsize(raw, int(widths[i] / 4.4))
                if aligns[i] == "CENTER":
                    c.drawCentredString(cx + widths[i] / 2, yy, txt)
                elif aligns[i] == "RIGHT":
                    c.drawRightString(cx + widths[i] - pad, yy, txt)
                else:
                    c.drawString(cx + pad, yy, txt)
            cx += widths[i]
    c.setStrokeColor(hx(C_LINE))
    c.setLineWidth(0.5)
    c.rect(x, y, total_w, y_top - y, fill=0, stroke=1)
    return y


def draw_image(c, path: Path, x, y, w, h=None):
    if not path or not Path(path).exists():
        return
    from reportlab.lib.utils import ImageReader
    img = ImageReader(str(path))
    iw, ih = img.getSize()
    if h is None:
        h = w * ih / iw
    c.drawImage(img, x, y, width=w, height=h, preserveAspectRatio=True, mask="auto", anchor="c")


def legend_strip(c, x, y, w):
    items = [
        (C_GREEN, "GREEN  at / above goal"),
        (C_AMBER, "AMBER  within 5 pp"),
        (C_RED, "RED  more than 5 pp off"),
    ]
    c.setFont(FONT, 7)
    cx = x
    for col, label in items:
        c.setFillColor(hx(col))
        c.rect(cx, y, 8, 8, fill=1, stroke=0)
        c.setFillColor(hx(C_GRAY))
        c.drawString(cx + 11, y + 1, label)
        cx += w / 3
    c.setFillColor(hx(C_MUTED))
    c.setFont(FONT, 6.5)
    c.drawString(x, y - 11, "Recontact is lower-is-better. FCR = 100 − recontact; no FCR target.")


# ── PDF slides ─────────────────────────────────────────────────────────────

def slide_title(c, ctx, page, total):
    c.setFillColor(hx(C_ORANGE))
    c.rect(0, 0, SLIDE_W, SLIDE_H, fill=1, stroke=0)
    c.setFillColor(hx(C_WHITE))
    c.setFont(FONT_BOLD, 18)
    c.drawString(ML, SLIDE_H - 0.85 * inch, "DiDi")
    c.setStrokeColor(hx(C_WHITE))
    c.setLineWidth(1.2)
    c.line(ML, SLIDE_H - 0.98 * inch, ML + 1.4 * inch, SLIDE_H - 0.98 * inch)
    c.setFont(FONT, 11)
    c.drawString(ML, SLIDE_H - 1.28 * inch, "CX Quality Analyst  ·  Service Operations")
    c.setFont(FONT_BOLD, 28)
    for i, line in enumerate(["Deliverable 2 —", "Weekly Performance Report"]):
        c.drawString(ML, SLIDE_H - 2.25 * inch - i * 34, line)
    c.setFont(FONT, 12)
    c.drawString(ML, SLIDE_H - 3.35 * inch, "May 2026   ·   Market All   ·   Official control totals")
    c.setFont(FONT, 9.5)
    lines = [
        f"Period  {ctx['period']}  ({ctx['date_span']})",
        "QA  =  mean audit score          CSAT  =  (4-star + 5-star) / Feedback CNT",
        "Recontact  =  Σ Recontact Volume / Σ Contacts     (ratio of sums, 12-channel mix)",
        f"Snapshot  {fmt_n(ctx['evals'])} QA evaluations  ·  {fmt_n(ctx['surveys'])} surveys  ·  {fmt_n(ctx['contacts'])} contacts",
        "Briefing date  20 August 2026",
    ]
    yy = SLIDE_H - 3.85 * inch
    for line in lines:
        c.drawString(ML, yy, line)
        yy -= 16
    draw_footer(c, page, total, title_slide=True)


def slide_exec(c, ctx, charts, page, total):
    draw_header(c)
    y = draw_kicker(c, "01", "Executive Summary")
    card_w = (USABLE - 2 * 0.10 * inch) / 3
    card_h = 1.08 * inch
    cy = y - card_h
    cards = [
        ("QA SCORE", f"{ctx['qa']:.2f}", [
            f"Goal {QA_GOAL:.0f}   ·   {fmt_pp(ctx['qa'] - QA_GOAL)} pp",
            f"{fmt_n(ctx['evals'])} evaluations",
        ], ctx["qa_st"]),
        ("CSAT", f"{ctx['csat']:.2f}%", [
            f"Goal {CSAT_GOAL:.0f}%   ·   {fmt_pp(ctx['csat'] - CSAT_GOAL)} pp",
            f"{fmt_n(ctx['surveys'])} surveys   (4-star+5-star)/Feedback CNT",
        ], ctx["csat_st"]),
        ("RECONTACT", f"{ctx['rc']:.2f}%", [
            f"Goal ≤{RECONTACT_GOAL:.2f}%   ·   {fmt_pp(ctx['rc'] - RECONTACT_GOAL)} pp",
            f"{fmt_n(ctx['contacts'])} contacts   Σ RC vol / Σ contacts",
        ], ctx["rc_st"]),
    ]
    for i, (lab, val, meta, st) in enumerate(cards):
        draw_kpi_card(c, ML + i * (card_w + 0.10 * inch), cy, card_w, card_h, lab, val, meta, st)

    split_top = cy - 0.12 * inch
    split_bottom = CONTENT_BOTTOM + 8
    split_h = split_top - split_bottom

    box_h = min(1.55 * inch, split_h - 0.42 * inch)
    c.setFillColor(hx(C_DARK))
    c.roundRect(ML, split_top - box_h, LEFT_W, box_h, 3, fill=1, stroke=0)
    c.setFillColor(hx(C_ORANGE))
    c.setFont(FONT_BOLD, 8)
    c.drawString(ML + 10, split_top - 14, "MOST CRITICAL FINDING")
    c.setFillColor(hx(C_WHITE))
    c.setFont(FONT, 8)
    ty = split_top - 28
    for line in wrap_lines(ctx["copy"]["finding"], FONT, 8, LEFT_W - 20):
        c.drawString(ML + 10, ty, line)
        ty -= 11

    legend_strip(c, ML, split_top - box_h - 18, LEFT_W)

    draw_insights(c, RIGHT_X, split_bottom, RIGHT_W, split_h, ctx["copy"]["exec"])
    draw_footer(c, page, total)


def slide_qa_channel(c, ctx, charts, page, total):
    draw_header(c)
    y = draw_kicker(c, "02", "QA by Channel — Phone vs Live Chat")
    # two summary tables stacked-ish: channel KPI then attributes
    header = ["Channel", "QA", "n evals", "vs 85", "CSAT", "n surveys", "Status"]
    rows = [
        ["Phone", f"{ctx['phone_qa']:.2f}", fmt_n(ctx["phone_n"]), fmt_pp(ctx["phone_qa"] - QA_GOAL),
         f"{ctx['phone_csat']:.2f}%", fmt_n(ctx["phone_surveys"]), STATUS_SHORT[ctx["phone_qa_st"]]],
        ["Live Chat", f"{ctx['chat_qa']:.2f}", fmt_n(ctx["chat_n"]), fmt_pp(ctx["chat_qa"] - QA_GOAL),
         f"{ctx['chat_csat']:.2f}%", fmt_n(ctx["chat_surveys"]), STATUS_SHORT[ctx["chat_qa_st"]]],
        ["Global", f"{ctx['qa']:.2f}", fmt_n(ctx["evals"]), fmt_pp(ctx["qa"] - QA_GOAL),
         f"{ctx['csat']:.2f}%", fmt_n(ctx["surveys"]), STATUS_SHORT[ctx["qa_st"]]],
    ]
    widths = [1.15*inch, 0.72*inch, 0.85*inch, 0.7*inch, 0.85*inch, 0.95*inch, 0.95*inch]
    y = draw_table(c, LEFT_X, y, widths, header, rows,
                   aligns=["LEFT", "CENTER", "CENTER", "CENTER", "CENTER", "CENTER", "CENTER"], row_h=16)

    y -= 14
    c.setFillColor(hx(C_DARK))
    c.setFont(FONT_BOLD, 8)
    c.drawString(LEFT_X, y, "Top failing attributes (fail rate = unique audits / channel n)")
    y -= 6
    attr_header = ["Channel", "Attribute", "Fail rate", "Share of fails", "Critical"]
    attr_rows = []
    for ch, items in (("Phone", ctx["phone_attrs"]), ("Live Chat", ctx["chat_attrs"])):
        for a in items:
            attr_rows.append([
                ch, a["attr"], f"{a['rate']:.1f}%", f"{a['pct_fails']:.1f}%",
                "Yes" if a["critical"] else "No",
            ])
    aw = [1.05*inch, 2.55*inch, 0.85*inch, 1.15*inch, 0.75*inch]
    y = draw_table(c, LEFT_X, y, aw, attr_header, attr_rows,
                   aligns=["LEFT", "LEFT", "CENTER", "CENTER", "CENTER"], row_h=14.5, font_size=7)

    y -= 14
    c.setFillColor(hx(C_DARK))
    c.setFont(FONT_BOLD, 8)
    c.drawString(LEFT_X, y, "Lowest CRs on Phone (n≥10) — Chat CRs with n≥10 are all ≥85")
    y -= 6
    cr_h = ["CR Lv4", "QA", "n", "Status"]
    cr_rows = []
    for _, r in ctx["phone_cr"].iterrows():
        cr_rows.append([ellipsize(r["CR_Lv4"], 42), f"{r['QA_Score']:.2f}", str(int(r["N"])), STATUS_SHORT[r["status"]]])
    cw = [4.15*inch, 0.7*inch, 0.55*inch, 0.95*inch]
    y = draw_table(c, LEFT_X, y, cw, cr_h, cr_rows,
                   aligns=["LEFT", "CENTER", "CENTER", "CENTER"], row_h=14.5, font_size=7)

    ih = CONTENT_TOP - 18 - CONTENT_BOTTOM - 8
    draw_insights(c, RIGHT_X, CONTENT_BOTTOM + 8, RIGHT_W, ih, ctx["copy"]["qa_ch"])
    draw_footer(c, page, total)


def slide_qa_cr(c, ctx, charts, page, total):
    draw_header(c)
    y = draw_kicker(c, "03", "QA by CR Lv4 — ranked underperformers (n ≥ 10)")
    header = ["CR Lv4 (detail)", "QA", "n evals", "vs 85", "Status"]
    rows = []
    fills = {}
    for i, r in ctx["qa_cr"].head(10).iterrows():
        st = r["status"]
        rows.append([
            ellipsize(r["CR_Lv4"], 48), f"{r['QA_Score']:.2f}", str(int(r["N"])),
            fmt_pp(r["vs"]), STATUS_SHORT[st],
        ])
        fills[len(rows) - 1] = st
    widths = [4.35*inch, 0.7*inch, 0.85*inch, 0.75*inch, 0.95*inch]
    y = draw_table(
        c, LEFT_X, y, widths, header, rows, fills=fills,
        aligns=["LEFT", "CENTER", "CENTER", "CENTER", "CENTER"], row_h=16, font_size=7.5,
    )
    c.setFillColor(hx(C_GRAY))
    c.setFont(FONT, 7)
    n_below = int((ctx["qa_cr"]["QA_Score"] < QA_GOAL).sum())
    c.drawString(
        LEFT_X, y - 14,
        f"{n_below} contact reasons with n≥10 sit below the QA goal of 85. "
        f"n = evaluation count (not weighted gap × n). Pickup/CSAT LOBs are not on this QA sample.",
    )
    # note box
    nb_h = 0.85 * inch
    ny = max(CONTENT_BOTTOM + 12, y - 28 - nb_h)
    c.setFillColor(hx(C_CARD))
    c.roundRect(LEFT_X, ny, LEFT_W, nb_h, 3, fill=1, stroke=0)
    c.setFillColor(hx(C_DARK))
    c.setFont(FONT_BOLD, 8)
    c.drawString(LEFT_X + 10, ny + nb_h - 14, "How to read this ranking")
    c.setFont(FONT, 7.5)
    note = (
        "Rows are sorted by QA score ascending among CR Lv4 names with at least 10 evaluations. "
        "Phone accounts for the red and amber rows. Live Chat volume (2,105 of 2,460) holds the global average at 94.14. "
        "A GREEN CR can still fail CSAT — see Combined (06)."
    )
    ty = ny + nb_h - 28
    for line in wrap_lines(note, FONT, 7.5, LEFT_W - 20):
        c.drawString(LEFT_X + 10, ty, line)
        ty -= 10.5

    ih = CONTENT_TOP - 18 - CONTENT_BOTTOM - 8
    draw_insights(c, RIGHT_X, CONTENT_BOTTOM + 8, RIGHT_W, ih, ctx["copy"]["qa_cr"])
    draw_footer(c, page, total)


def slide_csat(c, ctx, charts, page, total):
    draw_header(c)
    y = draw_kicker(c, "04", "CSAT / VOC — inversion vs QA")
    # BT table
    c.setFillColor(hx(C_DARK))
    c.setFont(FONT_BOLD, 8)
    c.drawString(LEFT_X, y + 2, "By Business Type  ·  CSAT = (4-star + 5-star) / Feedback CNT")
    y -= 4
    header = ["Business Type", "CSAT", "Surveys", "vs 85", "Status"]
    rows = []
    for _, r in ctx["bt"].iterrows():
        name = str(r["Business_Type"])
        n = int(r["Feedback"])
        if name == "Pickup":
            continue
        st = status(float(r["CSAT_Score"]), CSAT_GOAL, True)
        rows.append([name, f"{r['CSAT_Score']:.2f}%", fmt_n(n), fmt_pp(float(r["CSAT_vs_Goal"])), STATUS_SHORT[st]])
    widths = [1.7*inch, 0.85*inch, 1.0*inch, 0.8*inch, 0.95*inch]
    y = draw_table(c, LEFT_X, y, widths, header, rows,
                   aligns=["LEFT", "CENTER", "CENTER", "CENTER", "CENTER"], row_h=15)

    y -= 16
    c.setFillColor(hx(C_DARK))
    c.setFont(FONT_BOLD, 8)
    c.drawString(LEFT_X, y, "Largest unsatisfied volume by CR Lv4")
    y -= 6
    uh = ["CR Lv4", "CSAT", "Unsat.", "Surveys", "Status"]
    urows = []
    for _, r in ctx["unsat"].head(6).iterrows():
        st = status(float(r["CSAT_Score"]), CSAT_GOAL, True)
        urows.append([
            ellipsize(r["CR_Lv4"], 40), f"{r['CSAT_Score']:.1f}%",
            fmt_n(r["Unsatisfied"]), fmt_n(r["Feedback"]), STATUS_SHORT[st],
        ])
    uw = [2.55*inch, 0.7*inch, 0.75*inch, 0.85*inch, 0.85*inch]
    y = draw_table(c, LEFT_X, y, uw, uh, urows,
                   aligns=["LEFT", "CENTER", "CENTER", "CENTER", "CENTER"], row_h=14.5, font_size=7)

    # charts row
    img_y = CONTENT_BOTTOM + 8
    draw_image(c, charts.get("stars"), LEFT_X, img_y, 2.15 * inch, 1.55 * inch)
    draw_image(c, charts.get("bt"), LEFT_X + 2.25 * inch, img_y, 3.55 * inch, 1.55 * inch)
    # star legend
    c.setFont(FONT, 6)
    c.setFillColor(hx(C_GRAY))
    star_items = list(zip(ctx["stars"]["Rating"], ctx["stars"]["Pct"]))
    c.drawString(LEFT_X, img_y - 2, "   ".join(f"{a} {b:.1f}%" for a, b in star_items))

    ih = CONTENT_TOP - 18 - CONTENT_BOTTOM - 8
    draw_insights(c, RIGHT_X, CONTENT_BOTTOM + 8, RIGHT_W, ih, ctx["copy"]["csat"])
    draw_footer(c, page, total)


def slide_recontact(c, ctx, charts, page, total):
    draw_header(c)
    y = draw_kicker(c, "05", "Recontact — official mix vs live channels")
    header = ["Scope", "Rate", "Contacts", "vs 5.44%", "Status"]
    scopes = [
        ("All 12 channels (official)", ctx["official"]),
        ("Excluding Self Help", ctx["ex_sh"]),
        ("Phone + Live Chat only", ctx["audited"]),
    ]
    rows = []
    for name, s in scopes:
        rate = float(s["Rate"])
        st = status(rate, RECONTACT_GOAL, False)
        rows.append([name, f"{rate:.2f}%", fmt_n(s["Contacts"]), fmt_pp(float(s["vs_goal"])), STATUS_SHORT[st]])
    widths = [2.55*inch, 0.75*inch, 1.15*inch, 0.95*inch, 0.95*inch]
    y = draw_table(c, LEFT_X, y, widths, header, rows,
                   aligns=["LEFT", "CENTER", "CENTER", "CENTER", "CENTER"], row_h=16)

    c.setFillColor(hx(C_GRAY))
    c.setFont(FONT, 7)
    dil = ctx["dilution"]
    c.drawString(
        LEFT_X, y - 12,
        f"Self Help = {float(dil['share']):.0f}% of contacts at {float(dil['rate']):.2f}% recontact. "
        f"That weight is why 5.83% can sit next to a 15% live-agent rate.",
    )

    y -= 28
    c.setFillColor(hx(C_DARK))
    c.setFont(FONT_BOLD, 8)
    c.drawString(LEFT_X, y, "Top CR Lv4 by recontact volume")
    y -= 6
    rh = ["CR Lv4", "Recontacts", "Contacts", "Share", "Rate", "Status"]
    rrows = []
    for _, r in ctx["rc_cr"].head(7).iterrows():
        st = status(float(r["Recontact_Rate"]), RECONTACT_GOAL, False)
        rrows.append([
            ellipsize(r["CR_Lv4"], 36), fmt_n(r["Recontacts"]), fmt_n(r["Contacts"]),
            f"{r['Pct']:.1f}%", f"{r['Recontact_Rate']:.2f}%", STATUS_SHORT[st],
        ])
    rw = [2.35*inch, 0.95*inch, 0.9*inch, 0.65*inch, 0.7*inch, 0.8*inch]
    y = draw_table(c, LEFT_X, y, rw, rh, rrows,
                   aligns=["LEFT", "CENTER", "CENTER", "CENTER", "CENTER", "CENTER"],
                   row_h=14.2, font_size=7)

    img_h = 1.45 * inch
    if y - CONTENT_BOTTOM > img_h + 6:
        draw_image(c, charts.get("rc_scope"), LEFT_X, CONTENT_BOTTOM + 4, LEFT_W * 0.98, img_h)

    ih = CONTENT_TOP - 18 - CONTENT_BOTTOM - 8
    draw_insights(c, RIGHT_X, CONTENT_BOTTOM + 8, RIGHT_W, ih, ctx["copy"]["rc"])
    draw_footer(c, page, total)


def slide_combined(c, ctx, charts, page, total):
    draw_header(c)
    y = draw_kicker(c, "06", "Combined operational story — 2+ KPI fail CRs")
    header = ["CR Lv4", "QA (n)", "CSAT", "Recontact", "Pattern"]
    rows = []
    for b in ctx["cluster"]:
        qa = f"{b['qa']:.1f} ({b['qa_n']})" if b["qa_n"] else "—"
        cs = f"{b['csat']:.1f}%" if pd.notna(b["csat"]) else "—"
        rc = f"{b['rc']:.2f}%" if pd.notna(b["rc"]) else "—"
        rows.append([ellipsize(b["name"], 40), qa, cs, rc, "High QA + Low CSAT + High RC"])
    # plus other 2+ from engine if not already listed
    seen = {b["name"].casefold() for b in ctx["cluster"]}
    extra = ctx["combined"]
    extra = extra[extra["Pattern"].str.contains(r"\+", regex=True, na=False)] if not extra.empty else extra
    if not extra.empty:
        for _, r in extra.iterrows():
            if str(r["CR_Lv4"]).casefold() in seen:
                continue
            if len(rows) >= 7:
                break
            qa_n = int(r["QA_N"]) if pd.notna(r.get("QA_N")) else 0
            cs = f"{r['CSAT_Score']:.1f}%" if pd.notna(r.get("CSAT_Score")) else "—"
            rc = f"{r['Recontact_Rate']:.2f}%" if pd.notna(r.get("Recontact_Rate")) else "—"
            rows.append([
                ellipsize(r["CR_Lv4"], 40),
                f"{r['QA_Score']:.1f} ({qa_n})",
                cs, rc, str(r["Pattern"]),
            ])
            seen.add(str(r["CR_Lv4"]).casefold())
    widths = [2.55*inch, 0.95*inch, 0.75*inch, 0.95*inch, 2.15*inch]
    y = draw_table(c, LEFT_X, y, widths, header, rows,
                   aligns=["LEFT", "CENTER", "CENTER", "CENTER", "LEFT"], row_h=15.5, font_size=7)

    # story box
    y -= 16
    box_h = min(2.15 * inch, y - CONTENT_BOTTOM - 8)
    by = y - box_h
    c.setFillColor(hx(C_DARK))
    c.roundRect(LEFT_X, by, LEFT_W, box_h, 3, fill=1, stroke=0)
    c.setFillColor(hx(C_ORANGE))
    c.setFont(FONT_BOLD, 8)
    c.drawString(LEFT_X + 12, by + box_h - 16, "PROCESS vs OUTCOME")
    story = (
        "On the four status/money CRs, QA is GREEN (91–98) while CSAT is 64–68% and recontact is 13–19%. "
        "The audit confirms the script; the customer still comes back. That is not a coaching miss on greeting. "
        "It is a resolution-path miss (tracking, compensation, cancellation policy). "
        f"Official recontact {ctx['rc']:.2f}% understates live-channel repeat work because Self Help dilutes the denominator. "
        "Manage the cluster, not the global average."
    )
    c.setFillColor(hx(C_WHITE))
    c.setFont(FONT, 8)
    ty = by + box_h - 32
    for line in wrap_lines(story, FONT, 8, LEFT_W - 24):
        c.drawString(LEFT_X + 12, ty, line)
        ty -= 11.5

    ih = CONTENT_TOP - 18 - CONTENT_BOTTOM - 8
    draw_insights(c, RIGHT_X, CONTENT_BOTTOM + 8, RIGHT_W, ih, ctx["copy"]["combined"])
    draw_footer(c, page, total)


def _action_rows(ctx) -> list[list[str]]:
    bt = {str(r["Business_Type"]): r for _, r in ctx["bt"].iterrows()}
    def cs(name):
        r = bt.get(name)
        if r is None:
            return "—"
        st = STATUS_SHORT[status(float(r["CSAT_Score"]), CSAT_GOAL, True)]
        return f"{float(r['CSAT_Score']):.2f}% {st}  n={fmt_n(r['Feedback'])}"

    return [
        ["Food", cs("Food"),
         "Phone coaching: Time management + Complete and correct information (critical). Expand Phone audit coverage.",
         "QA Lead + Training", "20 Aug – 3 Sep 2026"],
        ["Full Service", cs("Full Service"),
         "Status/tracking macros in Live Chat; compensation path for order-not-received (QA 68.16, n=49).",
         "Ops + QA Content", "27 Aug – 17 Sep 2026"],
        ["Market Place", cs("Market Place"),
         "Separate marketplace vs full-service scripts. Train store vs platform ownership.",
         "LOB Lead + Training", "27 Aug – 10 Sep 2026"],
        ["Other", cs("Other"),
         "Reclassify the Other catch-all (CSAT 26.52%, n=558). Not a LOB program until taxonomy is clean.",
         "Analytics + Ops", "27 Aug – 3 Sep 2026"],
        ["All channels", f"RC {ctx['rc']:.2f}% AMBER; ex-SH {float(ctx['ex_sh']['Rate']):.2f}% RED",
         "Weekly war room on top 3 recontact CRs until official rate ≤5.44%. Publish Phone+Chat rate beside official.",
         "Director, Service Ops", "From 20 Aug 2026, weekly"],
        ["QA / BI", f"Global {ctx['qa']:.2f} hides Phone {ctx['phone_qa']:.2f}",
         "Publish QA by channel on the exec view — never the average alone.",
         "QA Analyst + BI", "20–26 Aug 2026"],
    ]


def slide_actions(c, ctx, charts, page, total):
    draw_header(c)
    y = draw_kicker(c, "07", "Action plans by LOB / Business Type")
    # table uses almost full width minus a slimmer insights panel
    table_w_ratio = 0.70
    tw = USABLE * table_w_ratio
    gap = 0.12 * inch
    ix = ML + tw + gap
    iw = SLIDE_W - MR - ix

    header = ["LOB", "Finding", "What", "Who", "When"]
    rows = []
    # wrap-ish via ellipsize; action table needs more width so we use full-ish left
    raw = _action_rows(ctx)
    for r in raw:
        rows.append([
            r[0],
            ellipsize(r[1], 32),
            ellipsize(r[2], 78),
            ellipsize(r[3], 22),
            ellipsize(r[4], 26),
        ])
    widths = [1.15*inch, 1.85*inch, 3.55*inch, 1.45*inch, 1.45*inch]
    # scale if needed
    scale = tw / sum(widths)
    widths = [w * scale for w in widths]
    yb = draw_table(c, ML, y, widths, header, rows,
                    aligns=["LEFT", "LEFT", "LEFT", "LEFT", "LEFT"], row_h=28, font_size=7)

    c.setFillColor(hx(C_GRAY))
    c.setFont(FONT, 7)
    c.drawString(ML, yb - 14, "Pickup n=35 surveys omitted. Dates are relative to the 20 August 2026 briefing. QA LOB = Delivery only.")

    ih = CONTENT_TOP - 18 - CONTENT_BOTTOM - 8
    draw_insights(c, ix, CONTENT_BOTTOM + 8, iw, ih, ctx["copy"]["action"], title="INSIGHTS")
    draw_footer(c, page, total)


def slide_notes(c, ctx, charts, page, total):
    draw_header(c)
    y = draw_kicker(c, "08", "Definitions, status rules, and source notes")
    blocks = [
        ("Official formulas (do not substitute)",
         "QA Score = simple mean of Score_Pct on 2,460 Delivery evaluations. "
         "CSAT % = (4-star + 5-star) / Feedback CNT on 77,266 surveys. "
         "Recontact % = Σ Recontact Volume / Σ Contacts on 994,591 contacts (12 channels). "
         "Never average row-level rates. FCR is derived as 100 − recontact; there is no FCR target."),
        ("Traffic-light rule",
         "GREEN = at or better than goal. AMBER = within 5 percentage points. RED = more than 5 pp off. "
         "Recontact is lower-is-better, so 5.83% vs 5.44% is AMBER (+0.39). Phone recontact 13.47% is RED."),
        ("What this sample is",
         "QA audits are Delivery LOB only (May 2026, W19–W22, 4–29 May). CSAT Business Type is Food / Full Service / "
         "Market Place / Other / Pickup. Pickup (n=35) is out of action-plan scope. Control totals match the Business Case snapshot."),
        ("Self Help dilution",
         f"Self Help is {float(ctx['dilution']['share']):.0f}% of contacts at {float(ctx['dilution']['rate']):.2f}% recontact. "
         f"Ex-Self Help {float(ctx['ex_sh']['Rate']):.2f}% and Phone+Chat {float(ctx['audited']['Rate']):.2f}% are diagnostic, not replacements for the official KPI."),
        ("Weighted gaps",
         "N on QA tables is the evaluation count. Do not label Σ(gap × n) as audits. CR Lv4 names can appear with case variants in source; cluster rows on slide 06 are matched case-insensitively."),
    ]
    yy = y
    col_w = USABLE
    for title, body in blocks:
        c.setFillColor(hx(C_CARD))
        h = 0.78 * inch
        c.roundRect(ML, yy - h, col_w, h, 3, fill=1, stroke=0)
        c.setFillColor(hx(C_ORANGE))
        c.setFont(FONT_BOLD, 8.5)
        c.drawString(ML + 10, yy - 14, title)
        c.setFillColor(hx(C_DARK))
        c.setFont(FONT, 8)
        ty = yy - 28
        for line in wrap_lines(body, FONT, 8, col_w - 22):
            c.drawString(ML + 10, ty, line)
            ty -= 11
        yy -= h + 8
    draw_footer(c, page, total)


def build_pdf(ctx, charts, path: Path):
    path.parent.mkdir(parents=True, exist_ok=True)
    slides = [
        lambda c, p, n: slide_title(c, ctx, p, n),
        lambda c, p, n: slide_exec(c, ctx, charts, p, n),
        lambda c, p, n: slide_qa_channel(c, ctx, charts, p, n),
        lambda c, p, n: slide_qa_cr(c, ctx, charts, p, n),
        lambda c, p, n: slide_csat(c, ctx, charts, p, n),
        lambda c, p, n: slide_recontact(c, ctx, charts, p, n),
        lambda c, p, n: slide_combined(c, ctx, charts, p, n),
        lambda c, p, n: slide_actions(c, ctx, charts, p, n),
        lambda c, p, n: slide_notes(c, ctx, charts, p, n),
    ]
    c = pdfcanvas.Canvas(str(path), pagesize=(SLIDE_W, SLIDE_H))
    c.setTitle("Deliverable 2 — Weekly Performance Report")
    c.setAuthor("DiDi CX Quality Analyst")
    total = len(slides)
    for i, fn in enumerate(slides, 1):
        fn(c, i, total)
        c.showPage()
    c.save()
    return total


# ── PPTX ───────────────────────────────────────────────────────────────────

def _rgb(hex_color: str):
    from pptx.dml.color import RGBColor
    h = hex_color.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def build_pptx(ctx, charts, path: Path, page_count: int):
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
    from pptx.util import Inches, Pt

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    def shape_fill(slide, l, t, w, h, fill, line=None):
        sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
        sh.fill.solid()
        sh.fill.fore_color.rgb = _rgb(fill)
        if line is None:
            sh.line.fill.background()
        else:
            sh.line.color.rgb = _rgb(line)
        return sh

    def tb(slide, l, t, w, h, text, size=12, bold=False, color=C_DARK, align=PP_ALIGN.LEFT, font="Segoe UI"):
        box = slide.shapes.add_textbox(l, t, w, h)
        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(size)
        p.font.bold = bold
        p.font.color.rgb = _rgb(color)
        p.font.name = font
        p.alignment = align
        return box

    def add_header(slide):
        shape_fill(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.46), C_DARK)
        shape_fill(slide, Inches(0), Inches(0.46), Inches(13.333), Inches(0.055), C_ORANGE)
        tb(slide, Inches(0.34), Inches(0.10), Inches(10), Inches(0.28), HEADER_TEXT, size=12, bold=True, color=C_WHITE)
        tb(slide, Inches(9.6), Inches(0.12), Inches(3.4), Inches(0.24), "May 2026  ·  Market All",
           size=9, color="#CFCFCF", align=PP_ALIGN.RIGHT)

    def add_footer(slide, page, title=False):
        fg = C_WHITE if title else C_DARK
        muted = "#FFE0CC" if title else C_GRAY
        y = Inches(7.18)
        if not title:
            shape_fill(slide, Inches(0.34), Inches(7.14), Inches(12.65), Inches(0.015), C_ORANGE)
        tb(slide, Inches(0.34), y, Inches(2.2), Inches(0.22), FOOTER_LEFT, size=8, bold=True, color=fg)
        tb(slide, Inches(3.3), y, Inches(6.7), Inches(0.22), FOOTER_CENTER, size=8, color=muted, align=PP_ALIGN.CENTER)
        tb(slide, Inches(10.5), y, Inches(2.5), Inches(0.22), f"Page {page} of {page_count}",
           size=8, color=fg, align=PP_ALIGN.RIGHT)

    def kicker(slide, num, title):
        tb(slide, Inches(0.34), Inches(0.58), Inches(12.5), Inches(0.32), f"{num}   {title}",
           size=14, bold=True, color=C_ORANGE)

    def add_table(slide, headers, rows, left, top, width, col_w=None, font_size=9):
        tbl_shape = slide.shapes.add_table(len(rows) + 1, len(headers), left, top, width, Inches(0.28 * (len(rows) + 1)))
        tbl = tbl_shape.table
        if col_w:
            for i, w in enumerate(col_w):
                tbl.columns[i].width = w
        for j, h in enumerate(headers):
            cell = tbl.cell(0, j)
            cell.text = h
            cell.fill.solid()
            cell.fill.fore_color.rgb = _rgb(C_DARK)
            for p in cell.text_frame.paragraphs:
                p.font.bold = True
                p.font.size = Pt(font_size)
                p.font.color.rgb = _rgb(C_WHITE)
                p.font.name = "Segoe UI"
                p.alignment = PP_ALIGN.CENTER
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        for i, row in enumerate(rows, 1):
            for j, val in enumerate(row):
                cell = tbl.cell(i, j)
                cell.text = str(val)
                raw = str(val)
                st_map = {"GREEN": C_GREEN, "AMBER": C_AMBER, "RED": C_RED,
                          "AT GOAL": C_GREEN, "WITHIN 5pp": C_AMBER, "OFF GOAL": C_RED}
                if raw in st_map:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = _rgb(st_map[raw])
                    fg = C_WHITE
                    bold = True
                else:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = _rgb(C_WHITE if i % 2 else C_ROW)
                    fg = C_DARK
                    bold = False
                for p in cell.text_frame.paragraphs:
                    p.font.size = Pt(font_size)
                    p.font.color.rgb = _rgb(fg)
                    p.font.bold = bold
                    p.font.name = "Segoe UI"
                    p.alignment = PP_ALIGN.CENTER if j else PP_ALIGN.LEFT
                cell.vertical_anchor = MSO_ANCHOR.MIDDLE
                cell.text_frame.word_wrap = True
        return tbl_shape

    def insights_box(slide, bullets, left, top, width, height):
        sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
        sh.fill.solid()
        sh.fill.fore_color.rgb = _rgb(C_PANEL)
        sh.line.color.rgb = _rgb(C_ORANGE)
        shape_fill(slide, left, top, Inches(0.06), height, C_ORANGE)
        tb(slide, left + Inches(0.16), top + Inches(0.08), width - Inches(0.24), Inches(0.28),
           "INSIGHTS", size=11, bold=True, color=C_ORANGE)
        box = slide.shapes.add_textbox(left + Inches(0.16), top + Inches(0.38), width - Inches(0.28), height - Inches(0.48))
        tf = box.text_frame
        tf.word_wrap = True
        for i, b in enumerate(bullets):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = "•  " + b
            p.font.size = Pt(10)
            p.font.color.rgb = _rgb(C_DARK)
            p.font.name = "Segoe UI"
            p.space_after = Pt(8)

    # 1 Title
    s = prs.slides.add_slide(blank)
    shape_fill(s, Inches(0), Inches(0), Inches(13.333), Inches(7.5), C_ORANGE)
    tb(s, Inches(0.5), Inches(0.7), Inches(8), Inches(0.4), "DiDi", size=22, bold=True, color=C_WHITE)
    tb(s, Inches(0.5), Inches(1.15), Inches(10), Inches(0.3), "CX Quality Analyst  ·  Service Operations", size=14, color=C_WHITE)
    tb(s, Inches(0.5), Inches(2.1), Inches(12), Inches(1.4),
       "Deliverable 2 —\nWeekly Performance Report", size=32, bold=True, color=C_WHITE)
    tb(s, Inches(0.5), Inches(3.6), Inches(12), Inches(0.35),
       "May 2026   ·   Market All   ·   Official control totals", size=14, color=C_WHITE)
    tb(s, Inches(0.5), Inches(4.15), Inches(12), Inches(1.6),
       f"Period  {ctx['period']}  ({ctx['date_span']})\n"
       f"QA  =  mean audit score          CSAT  =  (4-star + 5-star) / Feedback CNT\n"
       f"Recontact  =  Σ Recontact Volume / Σ Contacts     (ratio of sums, 12-channel mix)\n"
       f"Snapshot  {fmt_n(ctx['evals'])} QA evaluations  ·  {fmt_n(ctx['surveys'])} surveys  ·  {fmt_n(ctx['contacts'])} contacts\n"
       f"Briefing date  20 August 2026",
       size=12, color=C_WHITE)
    add_footer(s, 1, title=True)

    # 2 Exec
    s = prs.slides.add_slide(blank)
    add_header(s)
    kicker(s, "01", "Executive Summary")
    add_footer(s, 2)
    card_specs = [
        ("QA SCORE", f"{ctx['qa']:.2f}", f"Goal {QA_GOAL:.0f}  ·  {fmt_pp(ctx['qa']-QA_GOAL)} pp\n{fmt_n(ctx['evals'])} evaluations", ctx["qa_st"]),
        ("CSAT", f"{ctx['csat']:.2f}%", f"Goal {CSAT_GOAL:.0f}%  ·  {fmt_pp(ctx['csat']-CSAT_GOAL)} pp\n{fmt_n(ctx['surveys'])} surveys", ctx["csat_st"]),
        ("RECONTACT", f"{ctx['rc']:.2f}%", f"Goal ≤{RECONTACT_GOAL:.2f}%  ·  {fmt_pp(ctx['rc']-RECONTACT_GOAL)} pp\n{fmt_n(ctx['contacts'])} contacts", ctx["rc_st"]),
    ]
    for i, (lab, val, meta, st) in enumerate(card_specs):
        x = Inches(0.34 + i * 4.22)
        card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(0.95), Inches(4.08), Inches(1.28))
        card.fill.solid()
        card.fill.fore_color.rgb = _rgb(C_CARD)
        card.line.color.rgb = _rgb(C_LINE)
        shape_fill(s, x, Inches(0.95), Inches(0.07), Inches(1.28), STATUS_HEX[st])
        tb(s, x + Inches(0.16), Inches(1.00), Inches(2.2), Inches(0.22), lab, size=9, color=C_MUTED)
        tb(s, x + Inches(2.35), Inches(1.00), Inches(1.6), Inches(0.22), STATUS_PILL[st], size=8, bold=True, color=STATUS_HEX[st], align=PP_ALIGN.RIGHT)
        tb(s, x + Inches(0.16), Inches(1.20), Inches(3.7), Inches(0.42), val, size=26, bold=True, color=C_DARK)
        tb(s, x + Inches(0.16), Inches(1.64), Inches(3.7), Inches(0.5), meta, size=9, color=C_GRAY)
    # finding
    shape_fill(s, Inches(0.34), Inches(2.40), Inches(8.05), Inches(2.35), C_DARK)
    tb(s, Inches(0.48), Inches(2.48), Inches(7.8), Inches(0.25), "MOST CRITICAL FINDING", size=10, bold=True, color=C_ORANGE)
    tb(s, Inches(0.48), Inches(2.76), Inches(7.8), Inches(1.55), ctx["copy"]["finding"], size=12, color=C_WHITE)
    tb(s, Inches(0.34), Inches(4.85), Inches(8.05), Inches(0.45),
       "GREEN  at/above goal     AMBER  within 5 pp     RED  more than 5 pp off     Recontact is lower-is-better. FCR has no target.",
       size=9, color=C_GRAY)
    insights_box(s, ctx["copy"]["exec"], Inches(8.55), Inches(2.40), Inches(4.45), Inches(4.50))

    # 3 QA channel
    s = prs.slides.add_slide(blank)
    add_header(s)
    kicker(s, "02", "QA by Channel — Phone vs Live Chat")
    add_footer(s, 3)
    add_table(s,
              ["Channel", "QA", "n evals", "vs 85", "CSAT", "n surveys", "Status"],
              [
                  ["Phone", f"{ctx['phone_qa']:.2f}", fmt_n(ctx["phone_n"]), fmt_pp(ctx["phone_qa"]-QA_GOAL),
                   f"{ctx['phone_csat']:.2f}%", fmt_n(ctx["phone_surveys"]), STATUS_SHORT[ctx["phone_qa_st"]]],
                  ["Live Chat", f"{ctx['chat_qa']:.2f}", fmt_n(ctx["chat_n"]), fmt_pp(ctx["chat_qa"]-QA_GOAL),
                   f"{ctx['chat_csat']:.2f}%", fmt_n(ctx["chat_surveys"]), STATUS_SHORT[ctx["chat_qa_st"]]],
                  ["Global", f"{ctx['qa']:.2f}", fmt_n(ctx["evals"]), fmt_pp(ctx["qa"]-QA_GOAL),
                   f"{ctx['csat']:.2f}%", fmt_n(ctx["surveys"]), STATUS_SHORT[ctx["qa_st"]]],
              ],
              Inches(0.34), Inches(0.95), Inches(8.05), font_size=9)
    attr_rows = []
    for ch, items in (("Phone", ctx["phone_attrs"]), ("Live Chat", ctx["chat_attrs"])):
        for a in items:
            attr_rows.append([ch, a["attr"], f"{a['rate']:.1f}%", f"{a['pct_fails']:.1f}%", "Yes" if a["critical"] else "No"])
    add_table(s, ["Channel", "Attribute", "Fail rate", "Share of fails", "Critical"], attr_rows,
              Inches(0.34), Inches(2.15), Inches(8.05), font_size=8)
    cr_phone = []
    for _, r in ctx["phone_cr"].iterrows():
        cr_phone.append([ellipsize(r["CR_Lv4"], 42), f"{r['QA_Score']:.2f}", str(int(r["N"])), STATUS_SHORT[r["status"]]])
    add_table(s, ["Lowest Phone CR Lv4 (n≥10)", "QA", "n", "Status"], cr_phone,
              Inches(0.34), Inches(4.55), Inches(8.05), font_size=8)
    insights_box(s, ctx["copy"]["qa_ch"], Inches(8.55), Inches(0.95), Inches(4.45), Inches(5.95))

    # 4 QA CR
    s = prs.slides.add_slide(blank)
    add_header(s)
    kicker(s, "03", "QA by CR Lv4 — ranked underperformers (n ≥ 10)")
    add_footer(s, 4)
    cr_rows = []
    for _, r in ctx["qa_cr"].head(10).iterrows():
        cr_rows.append([ellipsize(r["CR_Lv4"], 48), f"{r['QA_Score']:.2f}", str(int(r["N"])),
                        fmt_pp(r["vs"]), STATUS_SHORT[r["status"]]])
    add_table(s, ["CR Lv4 (detail)", "QA", "n evals", "vs 85", "Status"], cr_rows,
              Inches(0.34), Inches(0.95), Inches(8.05), font_size=9)
    insights_box(s, ctx["copy"]["qa_cr"], Inches(8.55), Inches(0.95), Inches(4.45), Inches(5.95))

    # 5 CSAT
    s = prs.slides.add_slide(blank)
    add_header(s)
    kicker(s, "04", "CSAT / VOC — inversion vs QA")
    add_footer(s, 5)
    bt_rows = []
    for _, r in ctx["bt"].iterrows():
        if str(r["Business_Type"]) == "Pickup":
            continue
        st = status(float(r["CSAT_Score"]), CSAT_GOAL, True)
        bt_rows.append([str(r["Business_Type"]), f"{r['CSAT_Score']:.2f}%", fmt_n(r["Feedback"]),
                        fmt_pp(float(r["CSAT_vs_Goal"])), STATUS_SHORT[st]])
    add_table(s, ["Business Type", "CSAT", "Surveys", "vs 85", "Status"], bt_rows,
              Inches(0.34), Inches(0.95), Inches(8.05), font_size=9)
    urows = []
    for _, r in ctx["unsat"].head(6).iterrows():
        st = status(float(r["CSAT_Score"]), CSAT_GOAL, True)
        urows.append([ellipsize(r["CR_Lv4"], 40), f"{r['CSAT_Score']:.1f}%",
                      fmt_n(r["Unsatisfied"]), fmt_n(r["Feedback"]), STATUS_SHORT[st]])
    add_table(s, ["CR Lv4", "CSAT", "Unsat.", "Surveys", "Status"], urows,
              Inches(0.34), Inches(2.55), Inches(8.05), font_size=8)
    if charts.get("stars") and Path(charts["stars"]).exists():
        s.shapes.add_picture(str(charts["stars"]), Inches(0.34), Inches(5.05), width=Inches(2.3))
    if charts.get("bt") and Path(charts["bt"]).exists():
        s.shapes.add_picture(str(charts["bt"]), Inches(2.7), Inches(5.05), width=Inches(5.6))
    insights_box(s, ctx["copy"]["csat"], Inches(8.55), Inches(0.95), Inches(4.45), Inches(5.95))

    # 6 Recontact
    s = prs.slides.add_slide(blank)
    add_header(s)
    kicker(s, "05", "Recontact — official mix vs live channels")
    add_footer(s, 6)
    sc_rows = []
    for name, key in (("All 12 channels (official)", "official"), ("Excluding Self Help", "ex_sh"), ("Phone + Live Chat only", "audited")):
        sc = ctx[key]
        rate = float(sc["Rate"])
        st = status(rate, RECONTACT_GOAL, False)
        sc_rows.append([name, f"{rate:.2f}%", fmt_n(sc["Contacts"]), fmt_pp(float(sc["vs_goal"])), STATUS_SHORT[st]])
    add_table(s, ["Scope", "Rate", "Contacts", "vs 5.44%", "Status"], sc_rows,
              Inches(0.34), Inches(0.95), Inches(8.05), font_size=9)
    rc_rows = []
    for _, r in ctx["rc_cr"].head(7).iterrows():
        st = status(float(r["Recontact_Rate"]), RECONTACT_GOAL, False)
        rc_rows.append([ellipsize(r["CR_Lv4"], 36), fmt_n(r["Recontacts"]), fmt_n(r["Contacts"]),
                        f"{r['Pct']:.1f}%", f"{r['Recontact_Rate']:.2f}%", STATUS_SHORT[st]])
    add_table(s, ["CR Lv4", "Recontacts", "Contacts", "Share", "Rate", "Status"], rc_rows,
              Inches(0.34), Inches(2.15), Inches(8.05), font_size=8)
    if charts.get("rc_scope") and Path(charts["rc_scope"]).exists():
        s.shapes.add_picture(str(charts["rc_scope"]), Inches(0.34), Inches(5.15), width=Inches(8.0))
    insights_box(s, ctx["copy"]["rc"], Inches(8.55), Inches(0.95), Inches(4.45), Inches(5.95))

    # 7 Combined
    s = prs.slides.add_slide(blank)
    add_header(s)
    kicker(s, "06", "Combined operational story — 2+ KPI fail CRs")
    add_footer(s, 7)
    comb_rows = []
    for b in ctx["cluster"]:
        qa = f"{b['qa']:.1f} ({b['qa_n']})" if b["qa_n"] else "—"
        cs = f"{b['csat']:.1f}%" if pd.notna(b["csat"]) else "—"
        rc = f"{b['rc']:.2f}%" if pd.notna(b["rc"]) else "—"
        comb_rows.append([ellipsize(b["name"], 40), qa, cs, rc, "High QA + Low CSAT + High RC"])
    add_table(s, ["CR Lv4", "QA (n)", "CSAT", "Recontact", "Pattern"], comb_rows,
              Inches(0.34), Inches(0.95), Inches(8.05), font_size=9)
    shape_fill(s, Inches(0.34), Inches(2.85), Inches(8.05), Inches(3.9), C_DARK)
    tb(s, Inches(0.5), Inches(3.0), Inches(7.7), Inches(0.3), "PROCESS vs OUTCOME", size=12, bold=True, color=C_ORANGE)
    tb(s, Inches(0.5), Inches(3.35), Inches(7.7), Inches(3.2),
       "On the four status/money CRs, QA is GREEN (91–98) while CSAT is 64–68% and recontact is 13–19%. "
       "The audit confirms the script; the customer still comes back. That is not a coaching miss on greeting. "
       "It is a resolution-path miss (tracking, compensation, cancellation policy). "
       f"Official recontact {ctx['rc']:.2f}% understates live-channel repeat work because Self Help dilutes the denominator. "
       "Manage the cluster, not the global average.",
       size=13, color=C_WHITE)
    insights_box(s, ctx["copy"]["combined"], Inches(8.55), Inches(0.95), Inches(4.45), Inches(5.95))

    # 8 Actions
    s = prs.slides.add_slide(blank)
    add_header(s)
    kicker(s, "07", "Action plans by LOB / Business Type")
    add_footer(s, 8)
    act = _action_rows(ctx)
    add_table(s, ["LOB", "Finding", "What", "Who", "When"], act,
              Inches(0.28), Inches(0.95), Inches(9.0), font_size=8)
    insights_box(s, ctx["copy"]["action"], Inches(9.4), Inches(0.95), Inches(3.6), Inches(5.95))

    # 9 Notes
    s = prs.slides.add_slide(blank)
    add_header(s)
    kicker(s, "08", "Definitions, status rules, and source notes")
    add_footer(s, 9)
    notes = [
        ("Official formulas (do not substitute)",
         "QA = mean Score_Pct on 2,460 Delivery evaluations. CSAT = (4-star+5-star)/Feedback CNT on 77,266 surveys. "
         "Recontact = Σ Recontact Volume / Σ Contacts on 994,591 contacts. FCR = 100 − recontact; no FCR target."),
        ("Traffic-light rule",
         "GREEN = at/above goal. AMBER = within 5 pp. RED = more than 5 pp off. Recontact is lower-is-better."),
        ("Sample",
         "QA is Delivery LOB only, May 2026 W19–W22 (4–29 May). Pickup n=35 omitted from action plans. "
         "Control totals match the Business Case snapshot: QA 94.14 · CSAT 79.95% · Recontact 5.83%."),
        ("Self Help",
         f"Self Help is {float(ctx['dilution']['share']):.0f}% of contacts at {float(ctx['dilution']['rate']):.2f}%. "
         f"Ex-SH {float(ctx['ex_sh']['Rate']):.2f}% and Phone+Chat {float(ctx['audited']['Rate']):.2f}% are diagnostic."),
    ]
    for i, (t, b) in enumerate(notes):
        top = Inches(0.95 + i * 1.35)
        shape_fill(s, Inches(0.34), top, Inches(12.65), Inches(1.22), C_CARD)
        tb(s, Inches(0.5), top + Inches(0.08), Inches(12.3), Inches(0.28), t, size=13, bold=True, color=C_ORANGE)
        tb(s, Inches(0.5), top + Inches(0.38), Inches(12.3), Inches(0.75), b, size=12, color=C_DARK)

    path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(path))


def main():
    OUT.mkdir(parents=True, exist_ok=True)
    ctx = load_ctx()
    # Guardrail: official control totals
    assert abs(ctx["qa"] - CONTROL_TOTALS["qa"]) < 0.02, ctx["qa"]
    assert abs(ctx["csat"] - CONTROL_TOTALS["csat"]) < 0.02, ctx["csat"]
    assert abs(ctx["rc"] - CONTROL_TOTALS["recontact"]) < 0.02, ctx["rc"]
    charts = build_charts(ctx)
    pdf_path = OUT / "Entregable_2_Weekly_Performance_Report.pdf"
    pptx_path = OUT / "Entregable_2_Weekly_Performance_Report.pptx"
    n = build_pdf(ctx, charts, pdf_path)
    build_pptx(ctx, charts, pptx_path, n)
    meta = {
        "generated": datetime.now().isoformat(),
        "format": "16:9 slide deck",
        "pdf": str(pdf_path),
        "pptx": str(pptx_path),
        "slides": n,
        "kpis": {"qa": ctx["qa"], "csat": ctx["csat"], "recontact": round(ctx["rc"], 2)},
        "headers_footers": "black header + orange accent + confidential footer on every content page; title is full orange with footer",
    }
    (OUT / "report_metadata.json").write_text(json.dumps(meta, indent=2), encoding="utf-8")
    print(f"PDF  {pdf_path}")
    print(f"PPTX {pptx_path}")
    print(f"Slides {n}")
    print(f"KPIs QA={ctx['qa']:.2f} CSAT={ctx['csat']:.2f} RC={ctx['rc']:.2f}")
    print(f"Phone QA={ctx['phone_qa']:.2f} Chat QA={ctx['chat_qa']:.2f}")
    print(f"Phone CSAT={ctx['phone_csat']:.2f} Chat CSAT={ctx['chat_csat']:.2f}")


if __name__ == "__main__":
    main()
