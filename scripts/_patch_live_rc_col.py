"""Add recontact rate to the closure vs CSAT tables. Do not rebuild."""
from __future__ import annotations

import os
import shutil
import sys
from datetime import datetime

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

HERE = os.path.dirname(os.path.abspath(__file__))
ROOT = os.path.dirname(HERE)
sys.path.insert(0, HERE)
sys.path.insert(0, ROOT)

import didi_deck as dd  # noqa: E402
from cr_labels import cr_label  # noqa: E402
from modules.data_loader import load_all_data  # noqa: E402
from modules.resolution_csat import resolution_story  # noqa: E402
from modules import kpis as K  # noqa: E402

LIVE = os.path.join(ROOT, "entregable 2", "deck", "Entregable_2_Weekly_Performance_Report.pptx")
COPY = os.path.join(ROOT, "entregable 2", "Entregable_2_Weekly_Performance_Report.pptx")

COLS = [
    ("Contact reason (Lv4)", 1.48),
    ("Resolved", 0.56),
    ("CSAT", 0.54),
    ("Recontact", 0.80),
    ("Audits", 0.62),
    ("Surveys", 0.70),
]


def st(value, goal, lower_better=False):
    gap = (value - goal) if lower_better else (goal - value)
    if gap <= 0:
        return "green"
    return "amber" if gap <= 5 else "red"


def st_res(value):
    if value >= 70:
        return "green"
    if value > 50:
        return "amber"
    return "red"


def rows_for(nodes, rc_map):
    out = []
    for p in nodes[:4]:
        rate = rc_map.get(str(p["name"]).strip().casefold())
        rc_cell = ((f"{rate:.1f}%", st(rate, 5.44, True))
                   if rate is not None else ("—", None))
        out.append([
            cr_label(p["name"], 24),
            (f"{p['pct_res']:.0f}%", st_res(p["pct_res"])),
            (f"{p['csat']:.1f}%", st(p["csat"], 85)),
            rc_cell,
            f"{p['n']:,}",
            f"{p['fb']:,}",
        ])
    return out


def find_slide(prs):
    for i, slide in enumerate(prs.slides):
        blob = " ".join(sh.text_frame.text for sh in slide.shapes if sh.has_text_frame)
        if "Reasons with lower closure" in blob and "Reasons with higher closure" in blob:
            return i, slide
    raise KeyError("closure tables")


def main():
    rec = os.path.join(os.path.dirname(LIVE), "recovery")
    os.makedirs(rec, exist_ok=True)
    bak = os.path.join(rec, f"Entregable_2_before_rc_col_{datetime.now():%Y%m%d_%H%M%S}.pptx")
    shutil.copy2(LIVE, bak)

    data = load_all_data()
    R = resolution_story(data["fact_audits"], data["fact_csat"])
    rc_all = K.recontact_by_cr(data["fact_recontact"], top_n=None)
    rc_map = {
        str(k).strip().casefold(): float(v)
        for k, v in zip(rc_all["CR_Lv4"], rc_all["Recontact_Rate"])
    }

    prs = Presentation(LIVE)
    idx, slide = find_slide(prs)
    groups = [sh for sh in slide.shapes if sh.shape_type == MSO_SHAPE_TYPE.GROUP]
    if len(groups) != 2:
        raise RuntimeError(f"expected 2 grouped tables, found {len(groups)}")

    # Visual order: lower-closure table sits above higher-closure.
    groups = sorted(groups, key=lambda g: g.top)
    specs = [
        (groups[0], R["cr_b"]),
        (groups[1], R["cr_a"]),
    ]
    laid = []
    for group, nodes in specs:
        laid.append((group.left.inches, group.top.inches, group.width.inches, group.height.inches, nodes))
        group._element.getparent().remove(group._element)

    for x, y, w, h, nodes in laid:
        header_h = 0.24
        row_h = max(0.22, (h - header_h) / 4)
        dd.data_table(
            slide, x, y, COLS, rows_for(nodes, rc_map),
            header_h=header_h, row_h=row_h, heat_cols=(1, 2, 3),
            align_right=(1, 2, 3, 4, 5), size=7.2, header_size=7.0,
            bold_first=True,
        )

    try:
        prs.save(LIVE)
        saved = LIVE
    except PermissionError:
        saved = COPY
        prs.save(saved)
        print("LIVE open; wrote", saved)
    if os.path.abspath(saved) != os.path.abspath(COPY):
        try:
            shutil.copy2(saved, COPY)
        except OSError as exc:
            print("skip copy", exc)
    print("backup", bak)
    print("saved", saved, "slide", idx + 1)


if __name__ == "__main__":
    main()
