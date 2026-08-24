"""Patch the recovered Entregable 2 PPTX in place. Do not run build_entregable2_deck.py."""
from __future__ import annotations

import os
import shutil
import sys
from datetime import datetime

from pptx import Presentation

HERE = os.path.dirname(os.path.abspath(__file__))
ROOT = os.path.dirname(HERE)
sys.path.insert(0, HERE)

import didi_deck as dd  # noqa: E402

LIVE = os.path.join(ROOT, "entregable 2", "deck", "Entregable_2_Weekly_Performance_Report.pptx")
COPIES = [
    os.path.join(ROOT, "entregable 2", "Entregable_2_Weekly_Performance_Report.pptx"),
    os.path.join(ROOT, "entregable 2", "Entregable_2_LISTO.pptx"),
]
FOOTER = "DiDi Global, CX Service Operations   |   Confidential. Internal use only"
FOCUS_TITLE = "QA analysis: where we concentrate this month"


def pick(*names):
    for name in names:
        if hasattr(dd, name):
            return getattr(dd, name)
    raise AttributeError(names)


content_slide = pick("content_slide", "content_slide")
footnote = pick("footnote", "footnote")
data_table = pick("data_table", "data_table")
callout = pick("callout")
apply_sheet_numbers = pick("apply_sheet_numbers", "apply_sheet_numbers")
MARGIN = dd.MARGIN
CONTENT_W = pick("CONTENT_W", "CONTENT_W")
BODY_TOP = pick("BODY_TOP", "BODY_TOP")


def call_callout(slide, x, y, w, h, headline, body, head=12, body_size=10):
    import inspect
    kwargs = {}
    for name in inspect.signature(callout).parameters:
        if "head" in name and "size" in name:
            kwargs[name] = head
        elif "body" in name and "size" in name:
            kwargs[name] = body_size
    return callout(slide, x, y, w, h, headline, body, **kwargs)


def set_shape_text(shape, new):
    tf = shape.text_frame
    tf.word_wrap = True
    first = True
    for para in tf.paragraphs:
        if first:
            if para.runs:
                para.runs[0].text = new
                for extra in para.runs[1:]:
                    extra.text = ""
            else:
                para.text = new
            first = False
        else:
            for run in para.runs:
                run.text = ""


def replace_startswith(slide, mapping):
    """mapping: prefix -> new full text. Skip if new text already present."""
    n = 0
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        raw = shape.text_frame.text.strip()
        if not raw:
            continue
        for prefix, new in mapping.items():
            if new in raw:
                continue
            if raw.startswith(prefix) or raw == prefix.rstrip(":"):
                set_shape_text(shape, new)
                n += 1
                break
    return n


def move_slide(prs, old_index, new_index):
    xml_slides = prs.slides._sldIdLst
    slides = list(xml_slides)
    el = slides[old_index]
    xml_slides.remove(el)
    xml_slides.insert(new_index, el)


def has_focus_slide(prs):
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame and FOCUS_TITLE in shape.text_frame.text:
                return True
    return False


def add_focus_slide(prs):
    s = content_slide(
        prs,
        FOCUS_TITLE,
        "Correct the defects we already score. Chat information accuracy is a scoring change, separate from this list.",
        "12",
    )
    footnote(s, FOOTER)
    cols = [
        ("Defect", 2.55),
        ("Channel", 1.15),
        ("How we correct it", 5.35),
        ("Who", 1.85),
        ("When", 0.93),
    ]
    rows = [
        ["Time management (133 fails, 64.3%)",
         "Phone",
         "System-check script at the start of money and undelivered calls. Escalate if no data in 2 minutes.",
         "Supervisors", "27 Jun"],
        ["Complete and correct information (28 fails, -6.55 pp)",
         "Phone",
         "1-page job aid for undelivered and refund CRs. Coach the agents on those 28 audits.",
         "QA Lead + Supervisors", "27 Jun"],
        ["Greeting and identification (96 fails, 30.9%)",
         "Live Chat",
         "Mandatory open-chat greeting and identification macro.",
         "Chat Ops + QA Lead", "13 Jun"],
        ["Service attitude (69 fails, 22.2%)",
         "Live Chat",
         "Next-step macros on stuck CRs. Recalibrate attitude when the case cannot close.",
         "Chat Ops + QA Lead", "27 Jun"],
        ["Service availability (46 fails, -2.10 pp)",
         "Live Chat",
         "WFM and Chat Ops on queue, wrap-up, and chats left hanging. This is the Chat critical we already score.",
         "WFM + Chat Ops", "13 Jun"],
    ]
    data_table(s, MARGIN, BODY_TOP, cols, rows, header_h=0.40, row_h=0.70,
               size=8, bold_first=True)
    y = BODY_TOP + 0.40 + 5 * 0.70 + 0.16
    call_callout(
        s, MARGIN, y, CONTENT_W, 1.28,
        "This is defect correction, not a scoring change",
        "These five attributes are already on the form. Adding Chat correct-information "
        "as a critical is a separate scorecard change (Chat QA would read 80.14). "
        "It does not lift CSAT. Cases that followed process and still did not close "
        "belong to process, not to the agent.",
        head=12, body_size=10,
    )
    return s


def patch(prs):
    slides = list(prs.slides)
    counts = {}

    counts["phone"] = replace_startswith(slides[8], {
        "Primary Volume Driver:": (
            "Primary Volume Driver: Time Management is 64.3% of Phone fails (133). "
            "This month: system-check script at the start of money and undelivered calls; "
            "escalate if no data in 2 minutes."
        ),
        "Disproportionate Score Drag:": (
            "Disproportionate Score Drag: Complete & Correct Information is 28 fails and "
            "-6.55 pp (critical). This month: 1-page job aid for undelivered/refund CRs; "
            "coach the agents on those 28 audits."
        ),
        "Concentrated Root Causes:": (
            "Concentrated Root Causes: Together 77.8% of Phone defects. "
            "Month 1 Phone work is these two attributes only."
        ),
    })

    counts["chat"] = replace_startswith(slides[9], {
        "Soft Skill Bias:": (
            "Soft Skill Bias: Greeting (30.9%) and Service Attitude (22.2%) are 53.1% of "
            "Chat defects. This month: mandatory open-chat greeting/ID macro, and "
            "next-step macros on stuck CRs."
        ),
        "Highest Score Drag:": (
            "Highest Score Drag: Service Availability is 46 fails and -2.10 pp (critical). "
            "This month: WFM/Chat Ops on queue, wrap-up, and not leaving the chat hanging."
        ),
        "Rubric Blind Spot:": (
            "Rubric Blind Spot: The form still does not score information accuracy. "
            "That is a measurement change, separate from this month's defect correction."
        ),
    })

    counts["defects"] = replace_startswith(slides[10], {
        "Volume vs. Severity Trade-off:": (
            "Volume vs. Severity Trade-off: Time Management is the highest Phone volume fail; "
            "Complete & Correct Information is the highest Phone score drag. "
            "Month 1: those two plus Chat greeting, attitude, and availability."
        ),
    })

    counts["next"] = replace_startswith(slides[27], {
        "Add correct information as a critical Chat attribute": (
            "This month: correct Pareto defects we already score (Phone time management and "
            "complete information; Chat greeting, attitude, availability). That is operational, "
            "not a form change."
        ),
        "Track the official CSAT for each targeted contact reason": (
            "Separately, add Chat correct-information as a critical so the scorecard stops "
            "inflating (Chat QA would read 80.14). That does not lift CSAT. Track CSAT on "
            "targeted contact reasons."
        ),
    })

    counts["channel"] = replace_startswith(slides[30], {
        "Coach Time Management using a system-check script for money disputes": (
            "Coach Time Management: system-check script at start; escalate if no data in 2 minutes"
        ),
        "Closure reporting and the Chat critical change the scorecard.": (
            "This month also correct scored Pareto defects: Phone time management and complete "
            "information; Chat greeting, attitude, availability. Closure reporting and the Chat "
            "critical change the scorecard. ETA and refunds change the outcome. Do not mix those."
        ),
    })

    counts["reco"] = replace_startswith(slides[31], {
        "The priority is not simply more coaching. Measure resolution. Fix the QA scoring gaps.": (
            "The priority is not simply more coaching. Measure resolution. "
            "Correct this month's scored defects. Fix the QA scoring gaps. "
            "Give agents a way to close the highest-volume failure reasons."
        ),
        "Fix the QA scoring gaps": "Fix scored defects, then the form",
        "Current Chat QA: 96.01.": (
            "This month: Phone time management and complete information; Chat greeting, "
            "attitude, availability. Then close the Chat scoring gap (current 96.01 vs 80.14 "
            "if process-fails scored 0). The second number is measurement, not CSAT."
        ),
    })

    inserted = False
    if not has_focus_slide(prs):
        add_focus_slide(prs)
        move_slide(prs, len(list(prs.slides)) - 1, 11)
        inserted = True

    apply_sheet_numbers(prs)
    return counts, inserted


def backup(path):
    rec = os.path.join(os.path.dirname(path), "recovery")
    os.makedirs(rec, exist_ok=True)
    stamp = datetime.now().strftime("%Y%m%d_%H%M%S")
    dest = os.path.join(rec, f"Entregable_2_before_month1_{stamp}.pptx")
    shutil.copy2(path, dest)
    return dest


def main():
    bak = backup(LIVE)
    print("backup", bak)
    prs = Presentation(LIVE)
    counts, inserted = patch(prs)
    try:
        prs.save(LIVE)
        saved = LIVE
    except PermissionError:
        saved = os.path.join(ROOT, "entregable 2", "Entregable_2_LISTO.pptx")
        prs.save(saved)
        print("LIVE file is open; wrote", saved)
    n = len(list(Presentation(saved).slides))
    print("saved", saved, "slides", n, "inserted", inserted, "counts", counts)
    for dest in COPIES:
        if os.path.abspath(dest) == os.path.abspath(saved):
            continue
        try:
            shutil.copy2(saved, dest)
            print("copy", dest)
        except OSError as exc:
            print("skip", dest, exc)


if __name__ == "__main__":
    main()
