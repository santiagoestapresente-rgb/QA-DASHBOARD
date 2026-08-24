"""Paint traffic-light heat on the two volume slides. Do not rebuild."""
from __future__ import annotations

import os
import shutil
import sys
from datetime import datetime

from pptx import Presentation

HERE = os.path.dirname(os.path.abspath(__file__))
ROOT = os.path.dirname(HERE)
sys.path.insert(0, HERE)

import didi_deck as dd  # noqa: E402

LIVE = os.path.join(ROOT, "entregable 2", "deck", "Entregable_2_Weekly_Performance_Report.pptx")
COPY = os.path.join(ROOT, "entregable 2", "Entregable_2_Weekly_Performance_Report.pptx")
QA_TITLE = "QA analysis: lowest score vs fail volume"
CSAT_TITLE = "CSAT analysis: lowest score vs detractor volume"
BAND_Y0 = 1.55
BAND_Y1 = 4.42


def st(value, goal):
    gap = goal - value
    if gap <= 0:
        return "green"
    return "amber" if gap <= 5 else "red"


def st_n(n):
    if n < 10:
        return "red"
    if n < 20:
        return "amber"
    return "green"


def st_share(pct):
    if pct >= 7:
        return "red"
    if pct >= 4:
        return "amber"
    return "green"


def st_csat_share(pct):
    if pct >= 10:
        return "red"
    if pct >= 5:
        return "amber"
    return "green"


def st_below(n):
    if n >= 10:
        return "red"
    return "amber"


def wipe_band(slide, y0, y1):
    doomed = []
    for shape in slide.shapes:
        try:
            y = shape.top.inches
        except (AttributeError, TypeError, ValueError):
            continue
        if y0 <= y < y1:
            doomed.append(shape._element)
    for el in doomed:
        parent = el.getparent()
        if parent is not None:
            parent.remove(el)
    return len(doomed)


def find_slide(prs, title):
    for i, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if shape.has_text_frame and title in shape.text_frame.text:
                return i, slide
    raise KeyError(title)


def paint_qa(slide):
    n = wipe_band(slide, BAND_Y0, BAND_Y1)
    dd.data_table(
        slide, dd.MARGIN, dd.BODY_TOP + 0.32,
        [("Contact reason", 3.35), ("QA", 0.70), ("n", 0.55), ("Start here?", 1.45)],
        [
            ["Completed not received (market place)",
             ("47.5", st(47.5, 85)), ("4", st_n(4)), ("Too few", "red")],
            ["Active order, already received",
             ("65.8", st(65.8, 85)), ("12", st_n(12)), ("Watch", "amber")],
            ["Completed not received (full service)",
             ("68.2", st(68.2, 85)), ("49", st_n(49)), ("Start", "green")],
            ["Verbal aggression",
             ("76.0", st(76.0, 85)), ("5", st_n(5)), ("Too few", "red")],
            ["Refund status and conditions",
             ("76.4", st(76.4, 85)), ("25", st_n(25)), ("Start · Phone", "green")],
        ],
        header_h=0.36, row_h=0.46, heat_cols=(1, 2), status_cols=(3,),
        align_right=(1, 2), size=8, bold_first=True,
    )
    dd.data_table(
        slide, 7.15, dd.BODY_TOP + 0.32,
        [("Contact reason", 3.20), ("Fails", 0.70), ("Share", 0.70), ("Below 85", 0.90)],
        [
            ["Completed not received (full service)",
             ("37", st_share(7.1)), ("7.1%", st_share(7.1)), ("14", st_below(14))],
            ["Cancel the order",
             ("25", st_share(4.8)), ("4.8%", st_share(4.8)), ("5", st_below(5))],
            ["After-sales fraud review",
             ("23", st_share(4.4)), ("4.4%", st_share(4.4)), ("4", st_below(4))],
            ["Cash order blocked (antifraud)",
             ("21", st_share(4.1)), ("4.1%", st_share(4.1)), ("4", st_below(4))],
            ["Incomplete order",
             ("21", st_share(4.1)), ("4.1%", st_share(4.1)), ("7", st_below(7))],
        ],
        header_h=0.36, row_h=0.46, heat_cols=(1, 2, 3),
        align_right=(1, 2, 3), size=8, bold_first=True,
    )
    return n


def paint_csat(slide):
    n = wipe_band(slide, BAND_Y0, BAND_Y1)
    dd.data_table(
        slide, dd.MARGIN, dd.BODY_TOP + 0.32,
        [("Contact reason", 3.20), ("CSAT", 0.72), ("Surveys", 0.85), ("Share of pile", 1.20)],
        [
            ["After-sales fraud review",
             ("6.1%", st(6.1, 85)), "475", ("2.9%", st_csat_share(2.9))],
            ["Other (unmapped Business Type)",
             ("26.5%", st(26.5, 85)), "558", ("2.6%", st_csat_share(2.6))],
            ["Membership program renewal",
             ("28.6%", st(28.6, 85)), "248", ("1.1%", st_csat_share(1.1))],
            ["Membership program benefits",
             ("43.1%", st(43.1, 85)), "109", ("0.4%", st_csat_share(0.4))],
            ["Placing an order information",
             ("50.7%", st(50.7, 85)), "142", ("0.5%", st_csat_share(0.5))],
        ],
        header_h=0.36, row_h=0.46, heat_cols=(1, 3),
        align_right=(1, 2, 3), size=8, bold_first=True,
    )
    dd.data_table(
        slide, 7.15, dd.BODY_TOP + 0.32,
        [("Contact reason", 3.05), ("Unsat.", 0.78), ("Share", 0.70), ("CSAT", 0.72)],
        [
            ["Order status / delay information",
             ("3,189", st_csat_share(20.6)), ("20.6%", st_csat_share(20.6)), ("67.8%", st(67.8, 85))],
            ["Disagrees with cancellation charge",
             ("1,951", st_csat_share(12.6)), ("12.6%", st_csat_share(12.6)), ("67.4%", st(67.4, 85))],
            ["Order status & delays",
             ("1,800", st_csat_share(11.6)), ("11.6%", st_csat_share(11.6)), ("64.7%", st(64.7, 85))],
            ["No longer wants the order",
             ("1,229", st_csat_share(7.9)), ("7.9%", st_csat_share(7.9)), ("88.3%", st(88.3, 85))],
            ["Refund status and conditions",
             ("1,181", st_csat_share(7.6)), ("7.6%", st_csat_share(7.6)), ("67.0%", st(67.0, 85))],
        ],
        header_h=0.36, row_h=0.46, heat_cols=(1, 2, 3),
        align_right=(1, 2, 3), size=8, bold_first=True,
    )
    return n


def main():
    rec = os.path.join(os.path.dirname(LIVE), "recovery")
    os.makedirs(rec, exist_ok=True)
    bak = os.path.join(rec, f"Entregable_2_before_heat_{datetime.now():%Y%m%d_%H%M%S}.pptx")
    shutil.copy2(LIVE, bak)

    prs = Presentation(LIVE)
    _, qa = find_slide(prs, QA_TITLE)
    _, csat = find_slide(prs, CSAT_TITLE)
    n_qa = paint_qa(qa)
    n_csat = paint_csat(csat)
    try:
        prs.save(LIVE)
        saved = LIVE
    except PermissionError:
        saved = COPY
        prs.save(saved)
        print("LIVE open; wrote", saved)

    if os.path.abspath(saved) != os.path.abspath(COPY):
        try:
            shutil.copy2(saved, COPY)
        except OSError as exc:
            print("skip copy", exc)

    print("backup", bak)
    print("saved", saved, "wiped_qa", n_qa, "wiped_csat", n_csat)


if __name__ == "__main__":
    main()
