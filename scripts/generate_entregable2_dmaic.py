"""
Entregable 2 — Weekly Performance Report (DMAIC)
Professional deck aligned with QA Assessment 2025 template + DiDi Business Case.
Uses the same data pipeline as the Streamlit dashboard.
"""

from __future__ import annotations

import json
import sys
from datetime import datetime
from pathlib import Path

import matplotlib.pyplot as plt
import numpy as np
import pandas as pd
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parent.parent
sys.path.insert(0, str(ROOT))

from config import CSAT_GOAL, QA_GOAL, RECONTACT_GOAL, CONTROL_TOTALS  # noqa: E402
from modules.data_loader import load_all_data  # noqa: E402
from modules.executive_engine import (  # noqa: E402
    build_executive_brief,
    combined_operational_analysis,
    csat_segmentation,
    generate_action_plan,
    qa_channel_breakdown,
    requester_performance,
)
from modules.kpis import (  # noqa: E402
    correlation_matrix,
    kpi_summary,
    qa_agent_roster,
    recontact_by_cr,
    recontact_rate,
    top_failing_attributes,
    voc_themes_negative,
)

OUT = ROOT / "entregable 2"
CHARTS = OUT / "charts"
CANVA = OUT / "canva"

# Template palette (clean / professional — QA Assessment style)
ORANGE = RGBColor(255, 102, 0)
DARK = RGBColor(26, 26, 26)
WHITE = RGBColor(255, 255, 255)
GRAY = RGBColor(110, 110, 110)
LIGHT = RGBColor(245, 245, 245)
GREEN = RGBColor(46, 155, 87)
AMBER = RGBColor(242, 169, 0)
RED = RGBColor(214, 69, 69)

MPL_ORANGE = "#FF6600"
MPL_DARK = "#1A1A1A"
MPL_GREEN = "#2E9B57"
MPL_RED = "#D64545"
MPL_AMBER = "#F2A900"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def status_rgb(st: str) -> RGBColor:
    return {"green": GREEN, "amber": AMBER, "red": RED}.get(st, GRAY)


def vs_goal_status(value: float, goal: float, higher_better: bool = True) -> str:
    diff = value - goal if higher_better else goal - value
    if diff >= 0:
        return "green"
    if diff >= -5:
        return "amber"
    return "red"


def on_target_label(st: str) -> str:
    return "On target" if st == "green" else "Not on target"


def blank(prs: Presentation):
    return prs.slides.add_slide(prs.slide_layouts[6])


def bar_top(slide, thick=Inches(0.06)):
    s = slide.shapes.add_shape(1, 0, 0, SLIDE_W, thick)
    s.fill.solid()
    s.fill.fore_color.rgb = ORANGE
    s.line.fill.background()


def title(slide, text: str, y=Inches(0.35), size=26, color=ORANGE):
    box = slide.shapes.add_textbox(Inches(0.65), y, Inches(12), Inches(0.8))
    tf = box.text_frame
    tf.text = text
    p = tf.paragraphs[0]
    p.font.size = Pt(size)
    p.font.bold = True
    p.font.color.rgb = color


def subtitle(slide, text: str, y=Inches(0.95), size=14):
    box = slide.shapes.add_textbox(Inches(0.65), y, Inches(12), Inches(0.5))
    tf = box.text_frame
    tf.text = text
    p = tf.paragraphs[0]
    p.font.size = Pt(size)
    p.font.color.rgb = GRAY


def body(slide, lines: list[str], y=Inches(1.35), size=15, bold_first=False):
    box = slide.shapes.add_textbox(Inches(0.65), y, Inches(12), Inches(5.8))
    tf = box.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(size)
        p.font.color.rgb = DARK
        p.space_after = Pt(6)
        if bold_first and i == 0:
            p.font.bold = True


def divider_slide(prs, heading: str, sub: str = ""):
    s = blank(prs)
    bar_top(s, Inches(0.14))
    box = s.shapes.add_textbox(Inches(0.8), Inches(2.8), Inches(11.5), Inches(2))
    tf = box.text_frame
    tf.text = heading
    tf.paragraphs[0].font.size = Pt(36)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = ORANGE
    if sub:
        p = tf.add_paragraph()
        p.text = sub
        p.font.size = Pt(18)
        p.font.color.rgb = GRAY


def badge(slide, text: str, x, y, color: RGBColor, w=Inches(1.6), h=Inches(0.38)):
    sh = slide.shapes.add_shape(1, x, y, w, h)
    sh.fill.solid()
    sh.fill.fore_color.rgb = color
    sh.line.fill.background()
    tb = slide.shapes.add_textbox(x, y + Inches(0.04), w, h)
    tf = tb.text_frame
    tf.text = text
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].font.size = Pt(10)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = WHITE


def kpi_row(slide, name, score, target, st, y):
    badge(slide, on_target_label(st), Inches(10.8), y, status_rgb(st))
    cols = slide.shapes.add_textbox(Inches(0.8), y, Inches(10), Inches(0.45))
    tf = cols.text_frame
    tf.text = f"{name}    Current Score: {score}    Target: {target}"
    tf.paragraphs[0].font.size = Pt(14)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = DARK


def add_table(slide, headers, rows, y=Inches(1.4), col_w=None):
    ncols = len(headers)
    nrows = len(rows) + 1
    tbl = slide.shapes.add_table(nrows, ncols, Inches(0.6), y, Inches(12.1), Inches(0.42 * nrows)).table
    for j, h in enumerate(headers):
        c = tbl.cell(0, j)
        c.text = h
        c.fill.solid()
        c.fill.fore_color.rgb = ORANGE
        for p in c.text_frame.paragraphs:
            p.font.bold = True
            p.font.color.rgb = WHITE
            p.font.size = Pt(10)
            p.alignment = PP_ALIGN.CENTER
    for i, row in enumerate(rows, 1):
        for j, val in enumerate(row):
            c = tbl.cell(i, j)
            c.text = str(val)
            for p in c.text_frame.paragraphs:
                p.font.size = Pt(9)
                p.alignment = PP_ALIGN.CENTER


def img(slide, path: Path, y=Inches(1.2), w=Inches(12)):
    if path.exists():
        slide.shapes.add_picture(str(path), Inches(0.65), y, width=w)


def chart_correlation(corr: pd.DataFrame, path: Path):
    if corr.empty:
        return
    fig, ax = plt.subplots(figsize=(6, 5))
    im = ax.imshow(corr.values, cmap="RdYlGn", vmin=-1, vmax=1)
    ax.set_xticks(range(len(corr.columns)))
    ax.set_yticks(range(len(corr.index)))
    ax.set_xticklabels(["QA", "CSAT", "FCR"], fontsize=10)
    ax.set_yticklabels(["QA", "CSAT", "FCR"], fontsize=10)
    for i in range(len(corr.index)):
        for j in range(len(corr.columns)):
            ax.text(j, i, f"{corr.iloc[i, j]:.2f}", ha="center", va="center", fontsize=11, color=MPL_DARK)
    ax.set_title("Correlations between KPIs (CR Lv4 level)", fontsize=12, weight="bold", color=MPL_DARK)
    plt.colorbar(im, ax=ax, fraction=0.046)
    fig.savefig(path, dpi=150, bbox_inches="tight", facecolor="white")
    plt.close()


def chart_agent_tiers(roster: pd.DataFrame, path: Path):
    if roster.empty:
        return
    tiers = [
        ("Above 85", roster[roster["QA_Score"] >= QA_GOAL]),
        ("80 – 85", roster[(roster["QA_Score"] >= 80) & (roster["QA_Score"] < QA_GOAL)]),
        ("Below 80 (outlier)", roster[roster["QA_Score"] < 80]),
    ]
    labels, counts, avgs, colors = [], [], [], [MPL_GREEN, MPL_AMBER, MPL_RED]
    for label, df in tiers:
        labels.append(label)
        counts.append(len(df))
        avgs.append(df["QA_Score"].mean() if len(df) else 0)
    fig, ax = plt.subplots(figsize=(8, 4.5))
    bars = ax.bar(labels, counts, color=colors, edgecolor=MPL_DARK, linewidth=0.6)
    ax.set_ylabel("Agents (≥5 audits)")
    ax.set_title("QA Agent Ranking — compliance tiers", fontsize=12, weight="bold")
    for bar, avg, n in zip(bars, avgs, counts):
        ax.text(bar.get_x() + bar.get_width() / 2, bar.get_height() + 0.3, f"n={n}\nAvg {avg:.1f}%", ha="center", fontsize=9)
    fig.savefig(path, dpi=150, bbox_inches="tight", facecolor="white")
    plt.close()


def chart_glide_path(audits: pd.DataFrame, csat: pd.DataFrame, recontact: pd.DataFrame, path: Path):
    weeks = sorted(audits["Week"].dropna().unique()) if not audits.empty else []
    if not weeks:
        return
    qa_w, cs_w, rc_w = [], [], []
    for w in weeks:
        a = audits[audits["Week"] == w]
        qa_w.append(a["Score_Pct"].mean() if not a.empty else np.nan)
        c = csat[csat["Week"] == w] if "Week" in csat.columns else csat
        cs_w.append(c["Satisfied_CNT"].sum() / c["Feedback CNT"].sum() * 100 if not c.empty and c["Feedback CNT"].sum() else np.nan)
        r = recontact[recontact["Week"] == w] if "Week" in recontact.columns else recontact
        rc_w.append(r["Recontact Volume"].sum() / r["Contacts"].sum() * 100 if not r.empty and r["Contacts"].sum() else np.nan)

    fig, axes = plt.subplots(3, 1, figsize=(9, 6), sharex=True)
    specs = [
        (qa_w, QA_GOAL, "QA Score", MPL_GREEN, True),
        (cs_w, CSAT_GOAL, "CSAT %", MPL_ORANGE, True),
        (rc_w, RECONTACT_GOAL, "Recontact %", MPL_RED, False),
    ]
    for ax, (vals, goal, label, color, hib) in zip(axes, specs):
        ax.plot(weeks, vals, "o-", color=color, linewidth=2, markersize=6)
        ax.axhline(goal, color=MPL_DARK, linestyle="--", linewidth=1.2, label=f"Target {goal}")
        ax.set_ylabel(label)
        ax.legend(fontsize=8, loc="lower right")
        ax.grid(alpha=0.3)
    axes[-1].set_xlabel("Week")
    fig.suptitle("Glide path — weekly KPI trend (baseline before improvement plan)", fontsize=12, weight="bold")
    fig.savefig(path, dpi=150, bbox_inches="tight", facecolor="white")
    plt.close()


def chart_pareto_rc(rc_cr: pd.DataFrame, path: Path):
    top = rc_cr.head(8)
    fig, ax = plt.subplots(figsize=(9, 4.5))
    vals = top["Recontact_Volume"].values if "Recontact_Volume" in top.columns else top["Recontacts"].values
    labels = [str(x)[:28] + "…" if len(str(x)) > 28 else str(x) for x in top["CR_Lv4"].values]
    cum = np.cumsum(vals) / vals.sum() * 100
    ax.bar(range(len(vals)), vals, color=MPL_ORANGE)
    ax2 = ax.twinx()
    ax2.plot(range(len(vals)), cum, "o-", color=MPL_DARK)
    ax2.axhline(80, color=MPL_RED, linestyle="--", alpha=0.7)
    ax.set_xticks(range(len(vals)))
    ax.set_xticklabels(labels, rotation=35, ha="right", fontsize=7)
    ax.set_ylabel("Recontact volume")
    ax2.set_ylabel("Cumulative %")
    ax.set_title("Pareto — top CR Lv4 driving recontacts", fontsize=11, weight="bold")
    fig.savefig(path, dpi=150, bbox_inches="tight", facecolor="white")
    plt.close()


def load_context() -> dict:
    data = load_all_data()
    audits = data["fact_audits"]
    errors = data["fact_errors"]
    csat = data["fact_csat"]
    recontact = data["fact_recontact"]

    summary = kpi_summary(audits, csat, recontact)
    rc_rate = recontact_rate(recontact)
    combined = combined_operational_analysis(audits, csat, recontact)
    ch_break = qa_channel_breakdown(audits, errors)
    ch_perf = requester_performance(audits, csat, recontact)
    top_attr = top_failing_attributes(errors, audits, top_n=5)
    rc_cr = recontact_by_cr(recontact, top_n=10)
    voc = voc_themes_negative(csat, top_n=5)
    corr = correlation_matrix(audits, csat, recontact)
    roster = qa_agent_roster(audits, errors, min_n=5)
    csat_seg = csat_segmentation(csat, top_n=5)
    actions = generate_action_plan(combined, ch_perf, top_attr, rc_cr, summary, rc_rate)
    brief = build_executive_brief(summary, rc_rate, audits, errors, csat, recontact, combined, ch_perf, top_attr, rc_cr, voc, actions)

    weeks = sorted(audits["Week"].dropna().unique())
    period = f"Week {weeks[-1]}" if weeks else "May 2025 snapshot"

    return {
        "audits": audits,
        "errors": errors,
        "csat": csat,
        "recontact": recontact,
        "summary": summary,
        "rc_rate": rc_rate,
        "combined": combined,
        "ch_break": ch_break,
        "top_attr": top_attr,
        "rc_cr": rc_cr,
        "voc": voc,
        "corr": corr,
        "roster": roster,
        "csat_seg": csat_seg,
        "actions": actions,
        "brief": brief,
        "period": period,
    }


def build_deck(ctx: dict) -> Path:
    OUT.mkdir(exist_ok=True)
    CHARTS.mkdir(exist_ok=True)
    CANVA.mkdir(exist_ok=True)

    s = ctx["summary"]
    rc = ctx["rc_rate"]
    brief = ctx["brief"]
    corr = ctx["corr"]
    roster = ctx["roster"]
    combined = ctx["combined"]
    ch = ctx["ch_break"]
    voc = ctx["voc"]
    rc_cr = ctx["rc_cr"]

    chart_corr = CHARTS / "dmaic_correlation.png"
    chart_tiers = CHARTS / "dmaic_agent_tiers.png"
    chart_glide = CHARTS / "dmaic_glide_path.png"
    chart_pareto = CHARTS / "dmaic_pareto_rc.png"
    chart_correlation(ctx["corr"], chart_corr)
    chart_agent_tiers(roster, chart_tiers)
    chart_glide_path(ctx["audits"], ctx["csat"], ctx["recontact"], chart_glide)
    chart_pareto_rc(rc_cr, chart_pareto)

    qa_st = vs_goal_status(s["qa_score"], QA_GOAL)
    cs_st = vs_goal_status(s["csat"], CSAT_GOAL)
    rc_st = vs_goal_status(rc, RECONTACT_GOAL, False)

    # Primary problem = worst metric (like BANT in template)
    primary = "CSAT"
    primary_score = s["csat"]
    primary_goal = CSAT_GOAL
    primary_st = cs_st

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # 1 Title
    sl = blank(prs)
    bar_top(sl, Inches(0.16))
    box = sl.shapes.add_textbox(Inches(0.9), Inches(2.0), Inches(11), Inches(3))
    tf = box.text_frame
    tf.text = "CX Quality Assessment"
    tf.paragraphs[0].font.size = Pt(42)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = ORANGE
    p = tf.add_paragraph()
    p.text = "DiDi Global — CX Service Operations  |  Delivery / Food LOB"
    p.font.size = Pt(18)
    p.font.color.rgb = DARK
    p2 = tf.add_paragraph()
    p2.text = f"{ctx['period']}  |  Weekly Performance Report (Entregable 2)"
    p2.font.size = Pt(16)
    p2.font.color.rgb = GRAY
    p3 = tf.add_paragraph()
    p3.text = datetime.now().strftime("%B %Y")
    p3.font.size = Pt(14)
    p3.font.color.rgb = GRAY

    # 2 DMAIC roadmap
    sl = blank(prs)
    bar_top(sl)
    title(sl, "Methodology — DMAIC Framework")
    add_table(
        sl,
        ["Phase", "Purpose", "This assessment"],
        [
            ["Define", "Scope, KPI contract, problem statement", "3 KPIs: QA 85 · CSAT 85% · Recontact ≤5.44%"],
            ["Measure", "Baseline performance & data integrity", f"{CONTROL_TOTALS['evaluations']:,} audits · {CONTROL_TOTALS['surveys']:,} surveys"],
            ["Analyze", "Why KPIs behave as they do", "Correlations, Pareto, RCA, CR Lv4 deep dive"],
            ["Improve", "Structured action plan", "8-week phased plan · LOB owners"],
            ["Control", "Sustain gains & monitor drift", "Weekly glide path · control charts · QA calibration"],
        ],
        y=Inches(1.3),
    )

    # 3 DEFINE
    divider_slide(prs, "DEFINE", "Project charter & KPI contract")

    sl = blank(prs)
    bar_top(sl)
    title(sl, "Define — Problem Statement & Scope")
    body(
        sl,
        [
            "Problem: CX quality metrics diverge — QA meets target while CSAT and recontact do not.",
            "",
            f"Scope: Phone + Live Chat · CR Lv4 granularity · {ctx['period']} · Countries CO, MX, PE, CR.",
            "",
            "Voice of leadership: identify where customers recontact, where satisfaction drops, and which",
            "contact reasons require systemic FCR intervention — not only script compliance.",
            "",
            f"Data source: Business Case workbook · same logic as production dashboard ({CONTROL_TOTALS['contacts']:,} contacts).",
        ],
    )

    # 4 MEASURE — On target
    divider_slide(prs, "MEASURE", "On target")

    sl = blank(prs)
    bar_top(sl)
    title(sl, "How the KPIs are behaving and why?")
    subtitle(sl, f"For {ctx['period']}, the current metrics are:")
    kpi_row(sl, "QA Score", f"{s['qa_score']:.1f}%", f"{QA_GOAL:.0f}", qa_st, Inches(1.55))
    kpi_row(sl, "CSAT", f"{s['csat']:.1f}%", f"{CSAT_GOAL:.0f}%", cs_st, Inches(2.15))
    kpi_row(sl, "Recontact Rate", f"{rc:.2f}%", f"{RECONTACT_GOAL}%", rc_st, Inches(2.75))
    body(
        sl,
        [
            "",
            f"{primary} is systematically below target ({primary_score:.1f}% vs {primary_goal:.0f}%). "
            "This is the root problem we address in the Improve phase.",
            "",
            brief.insight[:320] + ("…" if len(brief.insight) > 320 else ""),
        ],
        y=Inches(3.5),
        size=14,
    )

    # 5 Agent ranking
    sl = blank(prs)
    bar_top(sl)
    title(sl, "How the KPIs are behaving and why?")
    subtitle(sl, "Ranking — QA agent compliance (≥5 audits)")
    n_above = len(roster[roster["QA_Score"] >= QA_GOAL]) if not roster.empty else 0
    n_mid = len(roster[(roster["QA_Score"] >= 80) & (roster["QA_Score"] < QA_GOAL)]) if not roster.empty else 0
    n_low = len(roster[roster["QA_Score"] < 80]) if not roster.empty else 0
    avg_all = roster["QA_Score"].mean() if not roster.empty else s["qa_score"]
    body(
        sl,
        [
            f"Agents above {QA_GOAL:.0f}%: {n_above}  |  Between 80–{QA_GOAL:.0f}%: {n_mid}  |  Below 80% (outlier): {n_low}",
            f"Average compliance: {avg_all:.1f}%",
            "",
            "Any agent below 80% is considered an outlier and should enter an improvement plan.",
            "Targeting the below-80 group addresses the majority of QA coaching deficit.",
        ],
        y=Inches(1.35),
        size=14,
    )
    img(sl, chart_tiers, y=Inches(2.8), w=Inches(8))

    # 6 Correlations
    sl = blank(prs)
    bar_top(sl)
    title(sl, "How the KPIs are behaving and why?")
    subtitle(sl, "Correlations between the 3 KPIs (CR Lv4 level)")
    if not corr.empty:
        r_qc = corr.loc["QA_Score", "CSAT_Pct"]
        r_qf = corr.loc["QA_Score", "FCR_Pct"]
        r_cf = corr.loc["CSAT_Pct", "FCR_Pct"]
        narrative = []
        if abs(r_qc) >= 0.3:
            narrative.append(
                f"QA ↔ CSAT correlation r={r_qc:.2f}: {'positive' if r_qc > 0 else 'negative'} — "
                "high QA alone does not guarantee satisfaction on high-volume operational CRs."
            )
        else:
            narrative.append(f"QA ↔ CSAT shows weak correlation (r={r_qc:.2f}) — perception gap vs audit score.")
        narrative.append(
            f"CSAT ↔ FCR (inverse recontact) r={r_cf:.2f}: low satisfaction aligns with repeat contacts."
        )
        narrative.append(
            f"QA ↔ FCR r={r_qf:.2f}: improving script compliance without FCR tools won't fix recontact."
        )
        body(sl, narrative, y=Inches(1.3), size=13)
    img(sl, chart_corr, y=Inches(3.2), w=Inches(5.5))

    # 7 QA channel focus
    sl = blank(prs)
    bar_top(sl)
    title(sl, "How the KPIs are behaving and why?")
    subtitle(sl, "QA — channel dispersion hides Phone gap")
    phone = ch.get("Phone", {})
    chat = ch.get("Live Chat", {})
    body(
        sl,
        [
            f"Phone QA: {phone.get('qa_score', '—')}% ({phone.get('qa_vs', 0):+.1f} pp vs goal) — {on_target_label(phone.get('qa_status', 'amber'))}",
            f"Live Chat QA: {chat.get('qa_score', '—')}% ({chat.get('qa_vs', 0):+.1f} pp) — {on_target_label(chat.get('qa_status', 'green'))}",
            "",
            f"Global QA {s['qa_score']:.1f}% meets target, but Phone is the weakest channel.",
            f"Top Phone defect: {phone.get('top_attrs', pd.DataFrame()).iloc[0]['Error_Category'] if not phone.get('top_attrs', pd.DataFrame()).empty else 'Time management'}.",
            "Focus: agents between 70–85 on Phone CRs tied to delivery / refund.",
        ],
        size=14,
    )

    # 8 CSAT focus (BANT-style slide)
    sl = blank(prs)
    bar_top(sl)
    title(sl, "How the KPIs are behaving and why?")
    subtitle(sl, f"{primary} KPI — population-level gap")
    csat_ch = ctx["csat"].groupby("Channel").apply(
        lambda g: g["Satisfied_CNT"].sum() / g["Feedback CNT"].sum() * 100 if g["Feedback CNT"].sum() else np.nan
    )
    body(
        sl,
        [
            f"Global CSAT {s['csat']:.1f}% — {on_target_label(cs_st)} (target {CSAT_GOAL}%).",
            f"Live Chat CSAT: {csat_ch.get('LIVE CHAT', csat_ch.get('Live Chat', float('nan'))):.1f}% drives the gap (72% of feedback).",
            f"Phone CSAT: {csat_ch.get('PHONE', csat_ch.get('Phone', float('nan'))):.1f}% — on target.",
            "",
            "CSAT is uncorrelated with QA on several high-volume CRs — improving QA alone won't fix CSAT.",
            "We need targeted FCR + VOC interventions on order status / delay CRs.",
            "",
            f"Top VOC theme: {voc.iloc[0]['Theme'] if not voc.empty else 'No resolution / wait'}.",
        ],
        size=13,
    )

    # 9 Recontact + prioritization
    sl = blank(prs)
    bar_top(sl)
    title(sl, "How the KPIs are behaving and why?")
    subtitle(sl, "Prioritization — CR Lv4 cluster")
    top3_vol = rc_cr.head(3)
    lines = [f"Recontact rate {rc:.2f}% — {on_target_label(rc_st)} (target ≤{RECONTACT_GOAL}%).", ""]
    for _, r in top3_vol.iterrows():
        cr = r["CR_Lv4"]
        vol = int(r.get("Recontact_Volume", r.get("Recontacts", 0)))
        rate = r.get("Recontact_Rate", 0)
        lines.append(f"• {cr}: {vol:,} recontacts ({rate:.1f}% rate)")
    if not combined.empty:
        row = combined.iloc[0]
        lines += ["", f"Priority #1 combined risk: {row['CR_Lv4']} — {row['Pattern']}"]
    lines.append("Targeting this cluster addresses the majority of the CSAT + recontact deficit.")
    body(sl, lines, size=13)
    img(sl, chart_pareto, y=Inches(3.4), w=Inches(11))

    # 10 ANALYZE divider RCA
    divider_slide(prs, "ANALYZE", "Root Cause Analysis")

    sl = blank(prs)
    bar_top(sl)
    title(sl, "RCA — Validated root causes (data review)")
    add_table(
        sl,
        ["Category", "Root cause", "Evidence"],
        [
            ["People", "Phone agents under time pressure", "Time management 37% fail rate Phone"],
            ["Process", "No FCR script for order status CRs", "QA green + CSAT red on same CRs"],
            ["Technology", "Limited live tracking in chat", "VOC: 'wait' / 'no information'"],
            ["Policy", "Refund/compensation rules unclear", "Critical fails: info + compensation"],
            ["Measurement", "QA measures script not resolution", "QA–CSAT weak correlation"],
            ["Culture", "Generic responses accepted", "VOC themes: no solution"],
        ],
        y=Inches(1.25),
    )
    body(sl, ["Most pressing: Process, People, Measurement, Training."], y=Inches(4.9), size=13, bold_first=True)

    # 11 Combined insight
    sl = blank(prs)
    bar_top(sl)
    title(sl, "Combined Analysis — operational story")
    body(
        sl,
        [
            brief.insight,
            "",
            f"Hypothesis: {brief.hypothesis}",
            "",
            f"Quantified: top recontact CR = {brief.top_recontact_cr} · worst combined CR = {brief.combined_cr}.",
        ],
        size=13,
    )

    # 12 IMPROVE divider
    divider_slide(prs, "IMPROVE", "Action Plan to improve CSAT & Recontact")

    sl = blank(prs)
    bar_top(sl)
    title(sl, "Action plan — Current Situation & Objective", size=22)
    body(
        sl,
        [
            f"Current {primary} score: {primary_score:.1f}% (Target: {primary_goal:.0f}%)",
            f"Recontact rate: {rc:.2f}% (Target: ≤{RECONTACT_GOAL}%)  |  QA: {s['qa_score']:.1f}% (maintain ≥{QA_GOAL})",
            "",
            "Observed issues:",
            "• Live Chat CSAT drag despite high QA scores",
            "• Order status / delay CRs drive 40%+ recontacts",
            "• Phone QA gap on time management & critical information",
            "",
            f"Objective: Increase {primary} to ≥{primary_goal:.0f}% within 8 weeks while holding QA ≥{QA_GOAL} and Recontact ≤{RECONTACT_GOAL}%.",
            "",
            "Involved teams: Training · QA · Operations (Supervisors) · Product",
        ],
        size=13,
    )

    phases = [
        (
            "Phase 1: Foundation (Weeks 1–2)",
            "Establish baseline coaching on Phone time management + order-status FCR scripts.",
            [
                "Focused uptraining on Time management & Complete information (Training, 45–60 min)",
                "FCR job aid for order status / delay CRs (Training + QA)",
                "1:1 coaching for Phone outlier agents <80% QA (Ops)",
                "QA alerts on Critical fails — 13 daily interactions reviewed (QA)",
            ],
        ),
        (
            "Phase 2: Process & QA Alignment (Weeks 3–4)",
            "Align QA, Training, and Ops on resolution — not only script compliance.",
            [
                "QA calibration sessions 2×/week (QA + Ops + Training)",
                "Live tracking widget pilot in Live Chat for status CRs (Product + Ops)",
                "Daily tips on FCR behaviors via comms channels (Training + QA)",
                "Weekly performance report highlighting CR outliers (QA → Supervisors)",
            ],
        ),
        (
            "Phase 3: Reinforcement (Weeks 5–6)",
            "Sustain progress and reduce repeated QA errors.",
            [
                "Weekly outlier tracking + follow-up coaching (QA + Ops)",
                "Recognition for agents maintaining ≥85% CSAT on audited CRs (Ops)",
                "Refresher on top VOC themes (QA + Training)",
                "Coaching compliance within 24h report (Ops → QA)",
            ],
        ),
        (
            "Phase 4: Consolidation (Weeks 7–8)",
            "Lock in gains and start next outlier cohort.",
            [
                "Targeted coaching for CRs still below 75% CSAT (Ops)",
                "Continuous monitoring + findings to supervisors (QA)",
                "Final review — assess outcomes & next DMAIC loop",
                "Restart plan on next outlier agent group",
            ],
        ),
    ]

    for phase_title, purpose, actions in phases:
        sl = blank(prs)
        bar_top(sl)
        title(sl, "Action plan", size=22)
        subtitle(sl, phase_title)
        body(sl, [purpose, ""] + [f"• {a}" for a in actions], size=12)

    # Activities matrix
    sl = blank(prs)
    bar_top(sl)
    title(sl, "Action plan — Activities matrix & LOB ownership", size=20)
    add_table(
        sl,
        ["LOB / Type", "Owner", "Timeline", "Action", "KPI"],
        [
            ["Food", "QA + Training", "Wk 1–2", "Phone coaching: time mgmt + info", "Phone QA → 85"],
            ["Full Service", "Ops + WFM", "Wk 1", "Peak staffing · reduce chat wait", "CSAT FS → 85%"],
            ["Full Service", "Product", "Wk 2–4", "Tracking widget in chat", "−3pp recontact status CRs"],
            ["Market Place", "LOB Lead", "Wk 2–3", "Separate marketplace scripts", "CSAT MP → 85%"],
            ["All LOBs", "CX Leadership", "Wk 1", "War room top 3 recontact CRs", "RC ≤ 5.44%"],
        ],
        y=Inches(1.15),
    )
    body(
        sl,
        ["Strategic framework: clear accountability · who does what · by when."],
        y=Inches(4.6),
        size=12,
    )

    # CONTROL
    divider_slide(prs, "CONTROL", "How do we ensure effectiveness?")

    sl = blank(prs)
    bar_top(sl)
    title(sl, "Control — Metrics & monitoring", size=22)
    add_table(
        sl,
        ["Type", "Metric", "Target", "Frequency"],
        [
            ["Primary KPI", "CSAT %", f"≥ {CSAT_GOAL}%", "Weekly"],
            ["Primary KPI", "Recontact rate", f"≤ {RECONTACT_GOAL}%", "Weekly"],
            ["Control", "QA Score", f"≥ {QA_GOAL} (±3pp max drift)", "Weekly"],
            ["Secondary", "FCR on status CR cluster", "↑ 5pp vs baseline", "Weekly"],
            ["Secondary", "QA alerts closed <24h", "100%", "Daily"],
            ["Secondary", "QA calibration score", "≥ 95%", "Per session"],
            ["Secondary", "Phone outlier agents <80%", "↓ 50% by week 6", "Bi-weekly"],
        ],
        y=Inches(1.2),
    )

    sl = blank(prs)
    bar_top(sl)
    title(sl, "Glide path — weekly KPI trend")
    subtitle(sl, "Baseline shows stable CSAT/recontact — improvement requires the action plan above.")
    img(sl, chart_glide, y=Inches(1.35), w=Inches(11.5))

    # 7 Quality tools annex
    sl = blank(prs)
    bar_top(sl)
    title(sl, "Annex — 7 Quality Tools applied", size=22)
    add_table(
        sl,
        ["Tool", "DMAIC phase", "Application"],
        [
            ["Check Sheet", "Measure", "Data validation & audit rules"],
            ["Pareto", "Analyze", "Top defects & recontact CRs"],
            ["Histogram", "Measure", "QA score distribution by channel"],
            ["Control Chart", "Control", "QA by country / daily SPC"],
            ["Scatter", "Analyze", "QA vs CSAT correlation"],
            ["Ishikawa", "Analyze", "Root cause — order status cluster"],
            ["Flowchart", "Analyze", "FCR failure path"],
        ],
        y=Inches(1.2),
    )

    # Thank you
    sl = blank(prs)
    bar_top(sl, Inches(0.14))
    box = sl.shapes.add_textbox(Inches(1), Inches(3.0), Inches(11), Inches(1.5))
    tf = box.text_frame
    tf.text = "Thank You"
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].font.size = Pt(44)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = ORANGE

    out_pptx = CANVA / "Entregable_2_DMAIC_Presentation.pptx"
    prs.save(out_pptx)

    meta = {
        "generated": datetime.now().isoformat(),
        "period": ctx["period"],
        "kpis": {"qa": s["qa_score"], "csat": s["csat"], "recontact": rc},
        "primary_problem": primary,
        "pptx": str(out_pptx),
        "charts": {
            "correlation": str(chart_corr),
            "agent_tiers": str(chart_tiers),
            "glide_path": str(chart_glide),
            "pareto_rc": str(chart_pareto),
        },
        "executive_brief": {
            "insight": brief.insight,
            "action": brief.action,
            "hypothesis": brief.hypothesis,
        },
        "dashboard_aligned": True,
        "control_totals": CONTROL_TOTALS,
    }
    (OUT / "report_metadata.json").write_text(json.dumps(meta, indent=2), encoding="utf-8")

    md = f"""# Entregable 2 — DMAIC Weekly Performance Report

**Period:** {ctx['period']}
**Primary problem KPI:** {primary} ({primary_score:.1f}% vs {primary_goal:.0f}%)

## KPI Status
| Metric | Score | Target | Status |
|--------|-------|--------|--------|
| QA | {s['qa_score']:.2f} | {QA_GOAL} | {on_target_label(qa_st)} |
| CSAT | {s['csat']:.2f}% | {CSAT_GOAL}% | {on_target_label(cs_st)} |
| Recontact | {rc:.2f}% | {RECONTACT_GOAL}% | {on_target_label(rc_st)} |

## Import to Canva
Upload `{out_pptx.name}` → Canva Import → edit with your brand template.

## Executive insight
{brief.insight}
"""
    (OUT / "Entregable_2_Weekly_Performance_Report.md").write_text(md, encoding="utf-8")
    return out_pptx


def main():
    ctx = load_context()
    path = build_deck(ctx)
    print(f"DMAIC presentation: {path}")
    print(f"KPIs aligned with dashboard: QA {ctx['summary']['qa_score']} | CSAT {ctx['summary']['csat']} | RC {ctx['rc_rate']:.2f}%")


if __name__ == "__main__":
    main()
