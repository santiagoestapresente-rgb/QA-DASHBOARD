"""Insert QA/CSAT volume slides into the recovered PPTX. Do not rebuild."""
from __future__ import annotations

import inspect
import os
import shutil
import sys
from datetime import datetime

from pptx import Presentation

HERE = os.path.dirname(os.path.abspath(__file__))
ROOT = os.path.dirname(HERE)
sys.path.insert(0, HERE)

import didi_deck as dd  # noqa: E402

LIVE = os.path.join(ROOT, "entregable 2", "deck", "Entregable_2_Weekly_Performance_Report.pptx")
COPIES = [
    os.path.join(ROOT, "entregable 2", "Entregable_2_Weekly_Performance_Report.pptx"),
]
FOOTER = "DiDi Global, CX Service Operations   |   Confidential. Internal use only"
QA_TITLE = "QA analysis: lowest score vs fail volume"
CSAT_TITLE = "CSAT analysis: lowest score vs detractor volume"


def pick(*names):
    for name in names:
        if hasattr(dd, name):
            return getattr(dd, name)
    raise AttributeError(names)


content_slide = pick("content_slide", "content_slide")
footnote = pick("footnote", "footnote")
data_table = pick("data_table", "data_table")
callout_fn = pick("callout")
text = pick("text")
apply_sheet_numbers = pick("apply_sheet_numbers", "apply_sheet_numbers")
INK = dd.INK
MARGIN = dd.MARGIN
CONTENT_W = pick("CONTENT_W", "CONTENT_W")
BODY_TOP = pick("BODY_TOP", "BODY_TOP")


def call_callout(slide, x, y, w, h, headline, body, head=13, body_size=11):
    kwargs = {}
    for name in inspect.signature(callout_fn).parameters:
        if "head" in name and "size" in name:
            kwargs[name] = head
        elif "body" in name and "size" in name:
            kwargs[name] = body_size
    return callout_fn(slide, x, y, w, h, headline, body, **kwargs)


def move_slide(prs, old_index, new_index):
    xml_slides = prs.slides._sldIdLst
    slides = list(xml_slides)
    el = slides[old_index]
    xml_slides.remove(el)
    xml_slides.insert(new_index, el)


def set_text(shape, new):
    tf = shape.text_frame
    tf.word_wrap = True
    first = True
    for para in tf.paragraphs:
        if first:
            if para.runs:
                para.runs[0].text = new
                for extra in para.runs[1:]:
                    extra.text = ""
            else:
                para.text = new
            first = False
        else:
            for run in para.runs:
                run.text = ""


def replace_contains(slide, needle, new_full_or_fragment, replace_fragment=True):
    n = 0
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        raw = shape.text_frame.text
        if needle not in raw:
            continue
        if replace_fragment:
            set_text(shape, raw.replace(needle, new_full_or_fragment).strip())
        else:
            set_text(shape, new_full_or_fragment)
        n += 1
    return n


def has_title(prs, title):
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame and title in shape.text_frame.text:
                return True
    return False


def add_qa_volume(prs):
    s = content_slide(
        prs, QA_TITLE,
        "A red score on 4 audits is not the same as 37 fails on one contact reason.",
        "14",
    )
    footnote(s, FOOTER)
    text(s, MARGIN, BODY_TOP, 6.05, 0.24, "Lowest QA (rate, n >= 3)",
         size=10.5, bold=True, color=INK)
    data_table(
        s, MARGIN, BODY_TOP + 0.32,
        [("Contact reason", 3.35), ("QA", 0.70), ("n", 0.55), ("Start here?", 1.45)],
        [
            ["Completed not received (market place)", "47.5", "4", "No. Too few audits."],
            ["Active order, already received", "65.8", "12", "Watch. Small n."],
            ["Completed not received (full service)", "68.2", "49", "Yes. Rate and volume."],
            ["Verbal aggression", "76.0", "5", "No. Too few audits."],
            ["Refund status and conditions", "76.4", "25", "Yes, on Phone."],
        ],
        header_h=0.36, row_h=0.46, size=8, bold_first=True,
    )
    text(s, 7.15, BODY_TOP, 5.70, 0.24, "Where the 518 attribute fails sit",
         size=10.5, bold=True, color=INK)
    data_table(
        s, 7.15, BODY_TOP + 0.32,
        [("Contact reason", 3.20), ("Fails", 0.70), ("Share", 0.70), ("Below 85", 0.90)],
        [
            ["Completed not received (full service)", "37", "7.1%", "14"],
            ["Cancel the order", "25", "4.8%", "5"],
            ["After-sales fraud review", "23", "4.4%", "4"],
            ["Cash order blocked (antifraud)", "21", "4.1%", "4"],
            ["Incomplete order", "21", "4.1%", "7"],
        ],
        header_h=0.36, row_h=0.46, size=8, bold_first=True,
    )
    y = BODY_TOP + 0.32 + 0.36 + 5 * 0.46 + 0.22
    call_callout(
        s, MARGIN, y, CONTENT_W, 1.55,
        "Phone QA volume focus is one contact reason. Chat QA fails are spread.",
        "Completed-not-received full service is 26% of Phone audits below 85 (14 of 53) "
        "and the single largest attribute-fail pile (37 of 518). Marketplace at 47.5 QA is 4 audits. "
        "Do not start coaching there. Chat below-goal audits (80) have no CR above 8% of that pile. "
        "Chat QA work stays on attributes (greeting, attitude, availability), not on one CR.",
    )
    return s


def add_csat_volume(prs):
    s = content_slide(
        prs, CSAT_TITLE,
        "15,488 unsatisfied surveys. Priority is share of that pile, not the ugliest percentage.",
        "18",
    )
    footnote(s, FOOTER)
    text(s, MARGIN, BODY_TOP, 6.05, 0.24, "Worst CSAT (rate, min 100 surveys)",
         size=10.5, bold=True, color=INK)
    data_table(
        s, MARGIN, BODY_TOP + 0.32,
        [("Contact reason", 3.20), ("CSAT", 0.72), ("Surveys", 0.85), ("Share of pile", 1.20)],
        [
            ["After-sales fraud review", "6.1%", "475", "2.9%"],
            ["Other (unmapped Business Type)", "26.5%", "558", "2.6%"],
            ["Membership program renewal", "28.6%", "248", "1.1%"],
            ["Membership program benefits", "43.1%", "109", "0.4%"],
            ["Placing an order information", "50.7%", "142", "0.5%"],
        ],
        header_h=0.36, row_h=0.46, size=8, bold_first=True,
    )
    text(s, 7.15, BODY_TOP, 5.70, 0.24, "Where the 15,488 unsatisfied sit",
         size=10.5, bold=True, color=INK)
    data_table(
        s, 7.15, BODY_TOP + 0.32,
        [("Contact reason", 3.05), ("Unsat.", 0.78), ("Share", 0.70), ("CSAT", 0.72)],
        [
            ["Order status / delay information", "3,189", "20.6%", "67.8%"],
            ["Disagrees with cancellation charge", "1,951", "12.6%", "67.4%"],
            ["Order status & delays", "1,800", "11.6%", "64.7%"],
            ["No longer wants the order", "1,229", "7.9%", "88.3%"],
            ["Refund status and conditions", "1,181", "7.6%", "67.0%"],
        ],
        header_h=0.36, row_h=0.46, size=8, bold_first=True,
    )
    y = BODY_TOP + 0.32 + 0.36 + 5 * 0.46 + 0.22
    call_callout(
        s, MARGIN, y, CONTENT_W, 1.55,
        "CSAT focus is order status, cancellation charge and refunds. Not fraud.",
        "The two order-status reasons plus cancellation charge are 44.8% of every unsatisfied survey. "
        "Add refunds and the top five are 60.3%. After-sales fraud is 6.1% CSAT and only 2.9% of the pile. "
        "No longer wants the order is 88.3% (above goal) and still 7.9% of detractors: volume without a rate problem. "
        "Do not staff that CR. ETA and refund path are the volume unlock.",
    )
    return s


def main():
    rec = os.path.join(os.path.dirname(LIVE), "recovery")
    os.makedirs(rec, exist_ok=True)
    bak = os.path.join(rec, f"Entregable_2_before_volume_{datetime.now():%Y%m%d_%H%M%S}.pptx")
    shutil.copy2(LIVE, bak)

    prs = Presentation(LIVE)
    n13 = replace_contains(
        prs.slides[12],
        'hiding a failing 68.2 Phone baseline (N=49).',
        "hiding a failing 68.2 Phone baseline (N=49). "
        "That CR is also the fail-volume leader (37 attribute fails). "
        "Lowest blended score is marketplace on 4 audits. Next slide splits rate vs volume.",
    )
    n15 = replace_contains(
        prs.slides[14],
        "Order Status & Delays: Chat drops to 60.3% CSAT vs. Phone's 79.8% (a 19.5 pp gap).",
        "Order Status & Delays: Chat drops to 60.3% CSAT vs. Phone's 79.8% (a 19.5 pp gap). "
        "Those CRs are also the detractor pile: order-status family plus cancellation charge is 44.8% of unsatisfied surveys. The volume split is two slides later (rate vs the 15,488-unsatisfied pile).",
    )
    n31 = 0
    n31 += replace_contains(
        prs.slides[30],
        "CSAT 80.74% on 46,071 surveys; order status repeats at 19.34%. Chat closes 6 of 10 assessed: this is that gap, not a staffing gap.",
        "CSAT 80.74% on 46,071 surveys; order-status family is 32.2% of unsatisfied surveys and repeats at 19.34%. Chat closes 6 of 10: ETA, not headcount.",
    )
    n31 += replace_contains(
        prs.slides[30],
        "CSAT 79.67%, the lowest of the volume lines; undelivered orders audit at 68.2 on 49 audits",
        "CSAT 79.67%, lowest volume line. Undelivered full-service is QA 68.2 (n=49) and the QA fail-volume leader (37 of 518 fails).",
    )
    n31 += replace_contains(
        prs.slides[30],
        "Refund or compensation not received is 18.8% of negative verbatims",
        "Refund or compensation not received is 18.8% of negative verbatims. Refund status is 7.6% of unsatisfied surveys (1,181).",
    )

    inserted = []
    if not has_title(prs, QA_TITLE):
        add_qa_volume(prs)
        move_slide(prs, len(list(prs.slides)) - 1, 13)
        inserted.append("qa")
    if not has_title(prs, CSAT_TITLE):
        add_csat_volume(prs)
        csat_idx = []
        for i, s in enumerate(prs.slides):
            for sh in s.shapes:
                if sh.has_text_frame and sh.text_frame.text.strip() == "CSAT analysis":
                    csat_idx.append(i)
                    break
        target = (csat_idx[-1] + 1) if csat_idx else 17
        move_slide(prs, len(list(prs.slides)) - 1, target)
        inserted.append("csat")

    apply_sheet_numbers(prs)
    try:
        prs.save(LIVE)
        saved = LIVE
    except PermissionError:
        saved = os.path.join(ROOT, "entregable 2", "Entregable_2_Weekly_Performance_Report.pptx")
        prs.save(saved)
        print("LIVE open; wrote", saved)

    print("backup", bak)
    print("saved", saved, "slides", len(list(Presentation(saved).slides)),
          "inserted", inserted, "n13", n13, "n15", n15, "n31", n31)
    for dest in COPIES:
        if os.path.abspath(dest) == os.path.abspath(saved):
            continue
        try:
            shutil.copy2(saved, dest)
            print("copy", dest)
        except OSError as exc:
            print("skip", dest, exc)


if __name__ == "__main__":
    main()
